#!/usr/bin/env python3
"""
Атомик Софт Presentation Generator
Generates branded .pptx from YAML briefs using python-pptx.
"""

import argparse
import io
import os
import sys
import yaml

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand palette ──────────────────────────────────────────────
BLUE       = RGBColor(0x00, 0x54, 0x97)
CYAN       = RGBColor(0x09, 0x97, 0xC8)
GREEN      = RGBColor(0x82, 0xC4, 0x44)
DARK_SLATE = RGBColor(0x44, 0x54, 0x69)
DARK_BLUE  = RGBColor(0x00, 0x2D, 0x52)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
RED_RISK   = RGBColor(0xCC, 0x33, 0x33)
AMBER_RISK = RGBColor(0xE8, 0x9C, 0x00)

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.625)

FONT_NAME = "Arial"


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def nl(text: str) -> str:
    """Normalize newlines in text (keep as \\n for split later)."""
    if text is None:
        return ""
    return str(text).replace("\\n", "\n")


# ── Template image extraction ──────────────────────────────────
def extract_template_images(template_path: str):
    """Extract background and logo blobs from template slide 0."""
    prs = Presentation(template_path)
    bg_blob = logo_blob = None
    bg_ct = logo_ct = None
    slide0 = prs.slides[0]
    for shape in slide0.shapes:
        if hasattr(shape, "image"):
            if "Рисунок 2" in shape.name:
                bg_blob = shape.image.blob
                bg_ct = shape.image.content_type
            elif "Рисунок 56" in shape.name:
                logo_blob = shape.image.blob
                logo_ct = shape.image.content_type
    return bg_blob, bg_ct, logo_blob, logo_ct


# ── Text helpers ───────────────────────────────────────────────
def _apply_run(run, text, size, color, bold, font_name):
    """Apply font properties to a run."""
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name


def _add_text_with_breaks(p, text, size, color, bold, font_name):
    """Add text to paragraph, splitting on \\n and inserting proper <a:br> line breaks."""
    from pptx.oxml.ns import qn
    parts = nl(text).split("\n")
    for i, part in enumerate(parts):
        if i > 0:
            # Insert XML line break element
            br = p._p.makeelement(qn("a:br"), {})
            # Copy run formatting to break element
            rPr = p._p.makeelement(qn("a:rPr"), {})
            rPr.set("lang", "ru-RU")
            rPr.set("sz", str(int(size * 100)))
            br.append(rPr)
            p._p.append(br)
        run = p.add_run()
        _apply_run(run, part, size, color, bold, font_name)


def set_text(tf, text, size=12, color=DARK_SLATE, bold=False, alignment=None, font_name=FONT_NAME):
    """Set text in a text frame, clearing existing content."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    _add_text_with_breaks(p, text, size, color, bold, font_name)
    return p


def add_para(tf, text, size=12, color=DARK_SLATE, bold=False, alignment=None, space_before=None, space_after=None):
    """Add a new paragraph to a text frame."""
    p = tf.add_paragraph()
    if alignment:
        p.alignment = alignment
    if space_before is not None:
        p.space_before = space_before
    if space_after is not None:
        p.space_after = space_after
    _add_text_with_breaks(p, text, size, color, bold, FONT_NAME)
    return p


def add_textbox(slide, left, top, width, height, text, size=12, color=DARK_SLATE,
                bold=False, alignment=None, word_wrap=True):
    """Add a textbox and set its text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    set_text(tf, text, size=size, color=color, bold=bold, alignment=alignment)
    return txBox


# ── Common slide elements ──────────────────────────────────────
def std_header(slide):
    """Two thin blue accent lines at y=0.47in."""
    y = Inches(0.47)
    h = Pt(2)
    # Left line
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y, Inches(0.5), h)
    line1.fill.solid()
    line1.fill.fore_color.rgb = BLUE
    line1.line.fill.background()
    # Right line
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), y, Inches(1.5), h)
    line2.fill.solid()
    line2.fill.fore_color.rgb = BLUE
    line2.line.fill.background()


def std_title(slide, text):
    """Standard header + title text."""
    std_header(slide)
    add_textbox(slide, Inches(0.55), Inches(0.22), Inches(8.0), Inches(0.4),
                text, size=22, color=BLUE, bold=True)


# ── Rounded rect helper ───────────────────────────────────────
def add_rounded_rect(slide, left, top, width, height, fill_color=LIGHT_GRAY):
    """Add a rounded rectangle with no border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Reduce corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_accent_stripe(slide, left, top, width, color):
    """Thin colored stripe at top of a card."""
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()
    return stripe


# ── Icon helper ────────────────────────────────────────────────
def add_icon(slide, icon_name, left, top, size, icons_dir):
    """Add an icon PNG to the slide."""
    icon_path = os.path.join(icons_dir, icon_name)
    if not os.path.exists(icon_path):
        print(f"  Warning: icon not found: {icon_path}", file=sys.stderr)
        return None
    pic = slide.shapes.add_picture(icon_path, left, top, size, size)
    return pic


# ── Slide builders ─────────────────────────────────────────────

def build_title_slide(prs, slide_data, brief, bg_blob, logo_blob, icons_dir):
    """Title slide with full-bleed background image."""
    layout = prs.slide_layouts[0]  # Helium_Break_2
    slide = prs.slides.add_slide(layout)
    # Remove placeholder shapes
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Background image
    if bg_blob:
        img_stream = io.BytesIO(bg_blob)
        slide.shapes.add_picture(img_stream, 0, 0, SLIDE_W, Inches(5.76))

    # Logo
    if logo_blob:
        logo_stream = io.BytesIO(logo_blob)
        slide.shapes.add_picture(logo_stream, Inches(0.3), Inches(0.27), Inches(2.5))

    # Title
    title_text = brief.get("title", "")
    add_textbox(slide, Inches(0.3), Inches(1.1), Inches(6), Inches(1.5),
                title_text, size=32, color=WHITE, bold=True)

    # Subtitle
    subtitle_text = brief.get("subtitle", "")
    if subtitle_text:
        add_textbox(slide, Inches(0.3), Inches(3.2), Inches(7), Inches(0.5),
                    subtitle_text, size=16, color=hex_to_rgb("80D0F0"))

    # Footer bar
    footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.95), SLIDE_W, Inches(0.675))
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = DARK_BLUE
    footer_bar.line.fill.background()

    # Footer text
    parts = []
    if brief.get("company"):
        parts.append(brief["company"])
    if brief.get("confidential"):
        parts.append("Конфиденциально")
    if brief.get("presenter"):
        parts.append(brief["presenter"])
    if brief.get("presenter_contact"):
        parts.append(brief["presenter_contact"])
    footer_text = "  |  ".join(parts) if parts else ""
    if footer_text:
        add_textbox(slide, Inches(0.55), Inches(5.05), Inches(8.9), Inches(0.4),
                    footer_text, size=9, color=hex_to_rgb("80B0D0"))

    return slide


def build_cards_slide(prs, slide_data, brief, bg_blob, logo_blob, icons_dir):
    """Cards slide: 2-3 cards side by side."""
    layout = prs.slide_layouts[3]  # General Slide
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    std_title(slide, slide_data.get("title", ""))

    cards = slide_data.get("cards", [])
    n = len(cards)
    if n <= 2:
        card_w = Inches(4.4)
        gap = Inches(0.2)
        start_x = Inches(0.5)
    else:
        card_w = Inches(2.95)
        gap = Inches(0.2)
        start_x = Inches(0.5)

    card_h = Inches(3.0)
    card_y = Inches(0.85)

    for i, card in enumerate(cards):
        x = start_x + i * (card_w + gap)
        accent = hex_to_rgb(card.get("accent", "#005497"))

        # Card background
        add_rounded_rect(slide, x, card_y, card_w, card_h)
        # Accent stripe
        add_accent_stripe(slide, x, card_y, card_w, accent)

        # Icon
        icon_name = card.get("icon", "")
        icon_size = Cm(1.2)
        if icon_name:
            add_icon(slide, icon_name, x + Inches(0.2), card_y + Inches(0.25), icon_size, icons_dir)

        # Heading
        heading = card.get("heading", "")
        add_textbox(slide, x + Inches(0.2) + icon_size + Inches(0.1), card_y + Inches(0.22),
                    card_w - Inches(0.5) - icon_size, Inches(0.4),
                    heading, size=13, color=accent, bold=True)

        # Body
        body = card.get("body", "")
        add_textbox(slide, x + Inches(0.2), card_y + Inches(0.85),
                    card_w - Inches(0.4), card_h - Inches(1.1),
                    body, size=10, color=DARK_SLATE)

    # Bottom line
    bottom_line = slide_data.get("bottom_line", "")
    if bottom_line:
        bc = hex_to_rgb(slide_data.get("bottom_color", "#CC3333"))
        add_textbox(slide, Inches(0.5), Inches(4.1), Inches(9.0), Inches(0.4),
                    bottom_line, size=11, color=bc, bold=True)

    return slide


def build_metrics_slide(prs, slide_data, brief, bg_blob, logo_blob, icons_dir):
    """Metrics slide with big numbers in card blocks."""
    layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    std_title(slide, slide_data.get("title", ""))

    # Subtitle
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.55), Inches(8.0), Inches(0.3),
                    subtitle, size=11, color=DARK_SLATE)

    metrics = slide_data.get("metrics", [])
    n = len(metrics)
    card_w = Inches(2.1)
    card_h = Inches(1.55)
    total_w = n * card_w + (n - 1) * Inches(0.2)
    start_x = (SLIDE_W - total_w) // 2
    card_y = Inches(1.0)

    for i, m in enumerate(metrics):
        x = start_x + i * (card_w + Inches(0.2))
        add_rounded_rect(slide, x, card_y, card_w, card_h)

        icon_name = m.get("icon", "")
        icon_size = Cm(1.2)
        if icon_name:
            icon_x = x + (card_w - icon_size) // 2
            add_icon(slide, icon_name, icon_x, card_y + Inches(0.1), icon_size, icons_dir)

        # Value
        add_textbox(slide, x, card_y + Inches(0.6), card_w, Inches(0.4),
                    m.get("value", ""), size=24, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

        # Label
        add_textbox(slide, x + Inches(0.1), card_y + Inches(1.05), card_w - Inches(0.2), Inches(0.4),
                    m.get("label", ""), size=9, color=DARK_SLATE, alignment=PP_ALIGN.CENTER)

    # Bullets
    bullets = slide_data.get("bullets", [])
    if bullets:
        bullet_y = Inches(2.8)
        for j, b in enumerate(bullets):
            by = bullet_y + j * Inches(0.35)
            icon_name = b.get("icon", "")
            if icon_name:
                add_icon(slide, icon_name, Inches(0.8), by, Cm(0.5), icons_dir)
            add_textbox(slide, Inches(1.2), by, Inches(8.0), Inches(0.3),
                        b.get("text", ""), size=10, color=DARK_SLATE)

    return slide


def build_contrast_slide(prs, slide_data, brief, bg_blob, logo_blob, icons_dir):
    """Two contrasting cards side by side."""
    layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    std_title(slide, slide_data.get("title", ""))

    card_w = Inches(4.4)
    card_h = Inches(2.2)
    card_y = Inches(0.85)
    gap = Inches(0.2)

    for idx, side_key in enumerate(["left", "right"]):
        side = slide_data.get(side_key, {})
        if not side:
            continue
        x = Inches(0.5) + idx * (card_w + gap)
        accent = hex_to_rgb(side.get("accent", "#005497"))

        add_rounded_rect(slide, x, card_y, card_w, card_h)
        add_accent_stripe(slide, x, card_y, card_w, accent)

        icon_name = side.get("icon", "")
        icon_size = Cm(1.2)
        if icon_name:
            add_icon(slide, icon_name, x + Inches(0.2), card_y + Inches(0.25), icon_size, icons_dir)

        # Heading
        add_textbox(slide, x + Inches(0.2) + icon_size + Inches(0.1), card_y + Inches(0.22),
                    card_w - Inches(0.5) - icon_size, Inches(0.35),
                    side.get("heading", ""), size=13, color=accent, bold=True)

        # Subtitle
        sub = side.get("subtitle", "")
        if sub:
            add_textbox(slide, x + Inches(0.2) + icon_size + Inches(0.1), card_y + Inches(0.52),
                        card_w - Inches(0.5) - icon_size, Inches(0.25),
                        sub, size=9, color=DARK_SLATE, bold=True)

        # Body
        add_textbox(slide, x + Inches(0.2), card_y + Inches(0.85),
                    card_w - Inches(0.4), card_h - Inches(1.0),
                    side.get("body", ""), size=10, color=DARK_SLATE)

    # Bottom section
    bottom_title = slide_data.get("bottom_title", "")
    bottom_points = slide_data.get("bottom_points", [])
    if bottom_title or bottom_points:
        # Divider line
        div_y = Inches(3.3)
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), div_y, Inches(9.0), Pt(1))
        div.fill.solid()
        div.fill.fore_color.rgb = BLUE
        div.line.fill.background()

        if bottom_title:
            add_textbox(slide, Inches(0.55), Inches(3.4), Inches(8.5), Inches(0.3),
                        bottom_title, size=13, color=BLUE, bold=True)

        for k, pt_text in enumerate(bottom_points):
            add_textbox(slide, Inches(0.8), Inches(3.75) + k * Inches(0.28), Inches(8.5), Inches(0.25),
                        f"• {pt_text}", size=10, color=DARK_SLATE)

    return slide


def build_columns_slide(prs, slide_data, brief, bg_blob, logo_blob, icons_dir):
    """Columns with colored headers."""
    layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    std_title(slide, slide_data.get("title", ""))

    columns = slide_data.get("columns", [])
    n = len(columns)
    if n == 0:
        return slide

    col_gap = Inches(0.15)
    total_gap = col_gap * (n - 1)
    usable = Inches(9.0)
    col_w = (usable - total_gap) // n
    start_x = Inches(0.5)
    header_h = Inches(0.4)
    col_y = Inches(0.85)

    for i, col in enumerate(columns):
        x = start_x + i * (col_w + col_gap)
        accent = hex_to_rgb(col.get("accent", "#005497"))

        # Header bar
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, col_y, col_w, header_h)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = accent
        hdr.line.fill.background()
        hdr.adjustments[0] = 0.1

        # Header text
        add_textbox(slide, x + Inches(0.1), col_y + Inches(0.05), col_w - Inches(0.2), header_h - Inches(0.1),
                    col.get("heading", ""), size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Items
        items = col.get("items", [])
        for j, item in enumerate(items):
            iy = col_y + header_h + Inches(0.15) + j * Inches(0.35)
            add_textbox(slide, x + Inches(0.1), iy, col_w - Inches(0.2), Inches(0.35),
                        f"• {nl(item)}", size=9, color=DARK_SLATE)

    return slide


def build_timeline_slide(prs, slide_data, brief, bg_blob, logo_blob, icons_dir):
    """Timeline with colored circles and arrows."""
    layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    std_title(slide, slide_data.get("title", ""))

    phases = slide_data.get("phases", [])
    n = len(phases)
    if n == 0:
        return slide

    circle_d = Inches(1.1)
    circle_y = Inches(1.1)
    total_w = Inches(8.5)
    spacing = total_w // n if n > 0 else total_w
    start_x = Inches(0.75)

    for i, phase in enumerate(phases):
        cx = start_x + i * spacing
        color = hex_to_rgb(phase.get("color", "#005497"))

        # Arrow connector (between circles)
        if i > 0:
            arrow_y = circle_y + circle_d // 2 - Pt(2)
            prev_cx = start_x + (i - 1) * spacing
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           prev_cx + circle_d + Inches(0.05), arrow_y,
                                           spacing - circle_d - Inches(0.1), Pt(8))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            arrow.line.fill.background()

        # Circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, circle_y, circle_d, circle_d)
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # White icon inside circle
        icon_name = phase.get("icon", "")
        # Try white variant first
        icon_w_name = icon_name.replace(".png", "-w.png") if icon_name else ""
        icon_size = Cm(1.0)
        icon_x = cx + (circle_d - icon_size) // 2
        icon_y_pos = circle_y + (circle_d - icon_size) // 2
        if icon_w_name and os.path.exists(os.path.join(icons_dir, icon_w_name)):
            add_icon(slide, icon_w_name, icon_x, icon_y_pos, icon_size, icons_dir)
        elif icon_name:
            add_icon(slide, icon_name, icon_x, icon_y_pos, icon_size, icons_dir)

        # Phase name
        text_x = cx - Inches(0.1)
        text_w = circle_d + Inches(0.2)
        add_textbox(slide, text_x, circle_y + circle_d + Inches(0.1), text_w, Inches(0.25),
                    phase.get("name", ""), size=9, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Duration
        add_textbox(slide, text_x, circle_y + circle_d + Inches(0.32), text_w, Inches(0.2),
                    phase.get("duration", ""), size=8, color=DARK_SLATE, alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, text_x - Inches(0.1), circle_y + circle_d + Inches(0.5),
                    text_w + Inches(0.2), Inches(0.7),
                    phase.get("description", ""), size=8, color=DARK_SLATE, alignment=PP_ALIGN.CENTER)

    # Summary
    summary = slide_data.get("summary", "")
    if summary:
        sy = Inches(4.4)
        add_rounded_rect(slide, Inches(0.5), sy, Inches(9.0), Inches(0.5), LIGHT_GRAY)
        add_textbox(slide, Inches(0.7), sy + Inches(0.05), Inches(8.6), Inches(0.4),
                    summary, size=10, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    return slide


def build_table_slide(prs, slide_data, brief, bg_blob, logo_blob, icons_dir):
    """Table with zebra striping."""
    layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    std_title(slide, slide_data.get("title", ""))

    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])
    if not headers:
        return slide

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header
    table_w = Inches(9.0)
    table_h = Inches(0.4) * n_rows
    table_x = Inches(0.5)
    table_y = Inches(0.9)

    table_shape = slide.shapes.add_table(n_rows, n_cols, table_x, table_y, table_w, table_h)
    table = table_shape.table

    # Set column widths
    col_w = table_w // n_cols
    for c in range(n_cols):
        table.columns[c].width = col_w

    # Header row
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        set_text(cell.text_frame, h, size=10, color=WHITE, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = ""
            set_text(cell.text_frame, str(val), size=9, color=DARK_SLATE)
            # Zebra striping
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return slide


def build_cta_slide(prs, slide_data, brief, bg_blob, logo_blob, icons_dir):
    """Call to action slide."""
    layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    std_title(slide, slide_data.get("title", ""))

    # Intro text
    intro = slide_data.get("intro", "")
    if intro:
        add_textbox(slide, Inches(0.55), Inches(0.6), Inches(8.5), Inches(0.3),
                    intro, size=11, color=DARK_SLATE)

    actions = slide_data.get("actions", [])
    action_y_start = Inches(1.1)
    action_spacing = Inches(0.95)

    for i, action in enumerate(actions):
        ay = action_y_start + i * action_spacing

        # Numbered circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), ay, Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE
        circle.line.fill.background()

        # Number text in circle
        tf = circle.text_frame
        tf.word_wrap = False
        set_text(tf, str(i + 1), size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Icon
        icon_name = action.get("icon", "")
        if icon_name:
            add_icon(slide, icon_name, Inches(1.15), ay + Inches(0.02), Cm(0.9), icons_dir)

        # Heading
        heading_x = Inches(1.65)
        add_textbox(slide, heading_x, ay, Inches(7.5), Inches(0.3),
                    action.get("heading", ""), size=12, color=BLUE, bold=True)

        # Body
        add_textbox(slide, heading_x, ay + Inches(0.3), Inches(7.5), Inches(0.55),
                    action.get("body", ""), size=10, color=DARK_SLATE)

    return slide


def build_closing_slide(prs, slide_data, brief, bg_blob, logo_blob, icons_dir):
    """Closing slide with dark background."""
    layout = prs.slide_layouts[4]  # 1_Contact Us
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Closing text from slide data or brief
    closing_text = slide_data.get("text", brief.get("title", ""))
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(2.0),
                closing_text, size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    closing_sub = slide_data.get("subtitle", "")
    if closing_sub:
        add_textbox(slide, Inches(0.5), Inches(2.8), Inches(9.0), Inches(0.5),
                    closing_sub, size=16, color=hex_to_rgb("80D0F0"), alignment=PP_ALIGN.CENTER)

    # Contact info
    contact_y = Inches(3.8)
    company = brief.get("company", "")
    presenter = brief.get("presenter", "")
    contact = brief.get("presenter_contact", "")

    if company:
        add_textbox(slide, Inches(1.0), contact_y, Inches(3.0), Inches(0.3),
                    company, size=12, color=WHITE, bold=True)
    if presenter:
        add_textbox(slide, Inches(4.0), contact_y, Inches(3.0), Inches(0.3),
                    presenter, size=12, color=WHITE)
    if contact:
        add_textbox(slide, Inches(7.0), contact_y, Inches(2.5), Inches(0.3),
                    contact, size=12, color=hex_to_rgb("80D0F0"))

    # Logo
    if logo_blob:
        logo_stream = io.BytesIO(logo_blob)
        slide.shapes.add_picture(logo_stream, Inches(3.75), Inches(4.5), Inches(2.5))

    return slide


# ── Builder dispatcher ─────────────────────────────────────────
BUILDERS = {
    "title":    build_title_slide,
    "cards":    build_cards_slide,
    "metrics":  build_metrics_slide,
    "contrast": build_contrast_slide,
    "columns":  build_columns_slide,
    "timeline": build_timeline_slide,
    "table":    build_table_slide,
    "cta":      build_cta_slide,
    "closing":  build_closing_slide,
}


def generate(brief: dict, template_path: str, icons_dir: str, output_path: str):
    """Main generation pipeline."""
    # Extract template images
    bg_blob, bg_ct, logo_blob, logo_ct = extract_template_images(template_path)

    # Create presentation from template (for theme/masters)
    prs = Presentation(template_path)

    # Set slide dimensions
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Remove existing slides properly
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Build slides
    slides_data = brief.get("slides", [])
    for i, sd in enumerate(slides_data):
        stype = sd.get("type", "")
        builder = BUILDERS.get(stype)
        if builder is None:
            print(f"  Warning: unknown slide type '{stype}' at index {i}, skipping", file=sys.stderr)
            continue
        print(f"  Building slide {i + 1}: {stype}")
        builder(prs, sd, brief, bg_blob, logo_blob, icons_dir)

    prs.save(output_path)
    print(f"  ✓ Saved: {output_path} ({len(brief.get('slides', []))} slides)")


def validate_brief(brief: dict):
    """Basic validation of the YAML brief."""
    if not isinstance(brief, dict):
        raise ValueError("YAML must be a mapping (dict) at top level")
    if "slides" not in brief:
        raise ValueError("Missing 'slides' key in YAML")
    if not isinstance(brief["slides"], list):
        raise ValueError("'slides' must be a list")
    for i, s in enumerate(brief["slides"]):
        if not isinstance(s, dict):
            raise ValueError(f"Slide {i} must be a mapping")
        if "type" not in s:
            raise ValueError(f"Slide {i} missing 'type' key")
        if s["type"] not in BUILDERS:
            print(f"  Warning: unknown slide type '{s['type']}' at index {i}", file=sys.stderr)


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    default_template = os.path.join(script_dir, "..", "templates", "atomicsoft-about.pptx")
    default_icons = os.path.join(script_dir, "..", "assets", "icons")

    parser = argparse.ArgumentParser(description="Generate branded PPTX from YAML brief")
    parser.add_argument("input", help="Path to YAML brief file")
    parser.add_argument("--output", "-o", default=None, help="Output PPTX path (default: <input_stem>.pptx)")
    parser.add_argument("--template", "-t", default=default_template, help="Template PPTX path")
    parser.add_argument("--icons", default=default_icons, help="Icons directory path")
    args = parser.parse_args()

    if args.output is None:
        stem = os.path.splitext(args.input)[0]
        args.output = stem + ".pptx"

    # Validate paths
    if not os.path.exists(args.input):
        print(f"Error: input file not found: {args.input}", file=sys.stderr)
        sys.exit(1)
    if not os.path.exists(args.template):
        print(f"Error: template not found: {args.template}", file=sys.stderr)
        sys.exit(1)
    if not os.path.isdir(args.icons):
        print(f"Error: icons directory not found: {args.icons}", file=sys.stderr)
        sys.exit(1)

    # Load YAML
    with open(args.input, "r", encoding="utf-8") as f:
        brief = yaml.safe_load(f)

    validate_brief(brief)

    print(f"Generating presentation from: {args.input}")
    print(f"  Template: {args.template}")
    print(f"  Icons: {args.icons}")
    print(f"  Output: {args.output}")
    generate(brief, args.template, args.icons, args.output)


if __name__ == "__main__":
    main()
