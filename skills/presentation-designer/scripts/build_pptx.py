#!/usr/bin/env python3
"""
Presentation builder for industrial automation presentations.
Uses python-pptx to generate .pptx files with professional styling.

Usage:
    python3 build_pptx.py --output presentation.pptx --style neutral
    python3 build_pptx.py --output presentation.pptx --style siemens
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# --- Color Palettes ---
PALETTES = {
    "neutral": {
        "primary": "#1B2A4A",
        "secondary": "#5A6B7F",
        "accent": "#00A5B5",
        "bg": "#FFFFFF",
        "text_light": "#FFFFFF",
        "text_dark": "#1B2A4A",
    },
    "siemens": {
        "primary": "#009999",
        "secondary": "#333333",
        "accent": "#009999",
        "bg": "#FFFFFF",
        "text_light": "#FFFFFF",
        "text_dark": "#333333",
    },
    "aveva": {
        "primary": "#1B365D",
        "secondary": "#3DCD58",
        "accent": "#3DCD58",
        "bg": "#FFFFFF",
        "text_light": "#FFFFFF",
        "text_dark": "#1B365D",
    },
    "honeywell": {
        "primary": "#E10600",
        "secondary": "#333333",
        "accent": "#E10600",
        "bg": "#FFFFFF",
        "text_light": "#FFFFFF",
        "text_dark": "#333333",
    },
    "emerson": {
        "primary": "#7B2D8E",
        "secondary": "#4A4A4A",
        "accent": "#7B2D8E",
        "bg": "#FFFFFF",
        "text_light": "#FFFFFF",
        "text_dark": "#4A4A4A",
    },
    "abb": {
        "primary": "#FF000F",
        "secondary": "#666666",
        "accent": "#FF000F",
        "bg": "#FFFFFF",
        "text_light": "#FFFFFF",
        "text_dark": "#333333",
    },
    "yokogawa": {
        "primary": "#003B71",
        "secondary": "#5A6B7F",
        "accent": "#0077C8",
        "bg": "#FFFFFF",
        "text_light": "#FFFFFF",
        "text_dark": "#003B71",
    },
    "rockwell": {
        "primary": "#C8102E",
        "secondary": "#666666",
        "accent": "#C8102E",
        "bg": "#FFFFFF",
        "text_light": "#FFFFFF",
        "text_dark": "#333333",
    },
    "owen": {
        "primary": "#F37021",
        "secondary": "#4A4A4A",
        "accent": "#F37021",
        "bg": "#FFFFFF",
        "text_light": "#FFFFFF",
        "text_dark": "#333333",
    },
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_title_slide(prs, title: str, subtitle: str = "", palette: dict = None):
    """Add a title slide with colored header bar."""
    p = palette or PALETTES["neutral"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Header bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(p["primary"])
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.text = title
    para.font.size = Pt(40)
    para.font.bold = True
    para.font.color.rgb = hex_to_rgb(p["text_light"])

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.5), Inches(1.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        para2 = tf2.paragraphs[0]
        para2.text = subtitle
        para2.font.size = Pt(20)
        para2.font.color.rgb = hex_to_rgb(p["secondary"])

    return slide


def add_content_slide(prs, title: str, bullets: list, notes: str = "", palette: dict = None):
    """Add a content slide with title bar and bullet points."""
    p = palette or PALETTES["neutral"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(p["primary"])
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    para = tf.paragraphs[0]
    para.text = title
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = hex_to_rgb(p["text_light"])

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(1.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(p["accent"])
    line.line.fill.background()

    # Bullets
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            para = tf2.paragraphs[0]
        else:
            para = tf2.add_paragraph()
        para.text = f"  •  {bullet}"
        para.font.size = Pt(18)
        para.font.color.rgb = hex_to_rgb(p["text_dark"])
        para.space_after = Pt(8)

    # Speaker notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_two_column_slide(prs, title: str, left_title: str, left_items: list,
                          right_title: str, right_items: list, notes: str = "", palette: dict = None):
    """Add a two-column slide (Problem→Solution, comparison, etc.)."""
    p = palette or PALETTES["neutral"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(p["primary"])
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.9))
    para = txBox.text_frame.paragraphs[0]
    para.text = title
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = hex_to_rgb(p["text_light"])

    # Left column
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    left_box.line.fill.background()

    ltx = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.3), Inches(5))
    ltf = ltx.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.text = left_title
    lp.font.size = Pt(22)
    lp.font.bold = True
    lp.font.color.rgb = hex_to_rgb(p["secondary"])
    for item in left_items:
        lp2 = ltf.add_paragraph()
        lp2.text = f"  •  {item}"
        lp2.font.size = Pt(16)
        lp2.space_after = Pt(6)

    # Right column
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = hex_to_rgb(p["accent"])
    right_box.line.fill.background()

    rtx = slide.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.3), Inches(5))
    rtf = rtx.text_frame
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rp.text = right_title
    rp.font.size = Pt(22)
    rp.font.bold = True
    rp.font.color.rgb = hex_to_rgb(p["text_light"])
    for item in right_items:
        rp2 = rtf.add_paragraph()
        rp2.text = f"  •  {item}"
        rp2.font.size = Pt(16)
        rp2.font.color.rgb = hex_to_rgb(p["text_light"])
        rp2.space_after = Pt(6)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_metrics_slide(prs, title: str, metrics: list, notes: str = "", palette: dict = None):
    """Add a metrics/KPI slide with large numbers. metrics = [{'value': '99.9%', 'label': 'Uptime'}, ...]"""
    p = palette or PALETTES["neutral"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(p["primary"])
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.9))
    para = txBox.text_frame.paragraphs[0]
    para.text = title
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = hex_to_rgb(p["text_light"])

    # Metrics grid
    count = len(metrics)
    col_width = min(3.0, 12.0 / max(count, 1))
    start_x = (13.333 - col_width * count) / 2

    for i, m in enumerate(metrics):
        x = start_x + i * col_width
        # Value
        vbox = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(col_width - 0.2), Inches(2))
        vp = vbox.text_frame.paragraphs[0]
        vp.text = str(m.get("value", ""))
        vp.font.size = Pt(48)
        vp.font.bold = True
        vp.font.color.rgb = hex_to_rgb(p["accent"])
        vp.alignment = PP_ALIGN.CENTER

        # Label
        lbox = slide.shapes.add_textbox(Inches(x), Inches(4.5), Inches(col_width - 0.2), Inches(1))
        lp = lbox.text_frame.paragraphs[0]
        lp.text = str(m.get("label", ""))
        lp.font.size = Pt(16)
        lp.font.color.rgb = hex_to_rgb(p["secondary"])
        lp.alignment = PP_ALIGN.CENTER

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_cta_slide(prs, title: str, contacts: list, palette: dict = None):
    """Add a final CTA slide with contact info."""
    p = palette or PALETTES["neutral"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(p["primary"])
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    para = tf.paragraphs[0]
    para.text = title
    para.font.size = Pt(36)
    para.font.bold = True
    para.font.color.rgb = hex_to_rgb(p["text_light"])
    para.alignment = PP_ALIGN.CENTER

    # Contacts
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, contact in enumerate(contacts):
        if i == 0:
            cp = tf2.paragraphs[0]
        else:
            cp = tf2.add_paragraph()
        cp.text = contact
        cp.font.size = Pt(20)
        cp.font.color.rgb = hex_to_rgb(p["text_light"])
        cp.alignment = PP_ALIGN.CENTER
        cp.space_after = Pt(12)

    return slide


def build_from_json(json_path: str, output: str, style: str = "neutral"):
    """Build presentation from a JSON definition file.
    
    JSON format:
    {
        "slides": [
            {"type": "title", "title": "...", "subtitle": "..."},
            {"type": "content", "title": "...", "bullets": [...], "notes": "..."},
            {"type": "two_column", "title": "...", "left_title": "...", "left_items": [...], 
             "right_title": "...", "right_items": [...], "notes": "..."},
            {"type": "metrics", "title": "...", "metrics": [{"value": "...", "label": "..."}], "notes": "..."},
            {"type": "cta", "title": "...", "contacts": [...]}
        ]
    }
    """
    with open(json_path) as f:
        data = json.load(f)

    palette = PALETTES.get(style, PALETTES["neutral"])
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for s in data["slides"]:
        t = s["type"]
        if t == "title":
            add_title_slide(prs, s["title"], s.get("subtitle", ""), palette)
        elif t == "content":
            add_content_slide(prs, s["title"], s.get("bullets", []), s.get("notes", ""), palette)
        elif t == "two_column":
            add_two_column_slide(prs, s["title"], s.get("left_title", ""),
                                  s.get("left_items", []), s.get("right_title", ""),
                                  s.get("right_items", []), s.get("notes", ""), palette)
        elif t == "metrics":
            add_metrics_slide(prs, s["title"], s.get("metrics", []), s.get("notes", ""), palette)
        elif t == "cta":
            add_cta_slide(prs, s["title"], s.get("contacts", []), palette)

    prs.save(output)
    print(f"Presentation saved: {output}")


def main():
    parser = argparse.ArgumentParser(description="Build industrial automation presentations")
    parser.add_argument("--output", "-o", default="presentation.pptx", help="Output .pptx path")
    parser.add_argument("--style", "-s", default="neutral", choices=list(PALETTES.keys()),
                        help="Visual style / vendor palette")
    parser.add_argument("--json", "-j", help="Build from JSON definition file")
    args = parser.parse_args()

    if args.json:
        build_from_json(args.json, args.output, args.style)
    else:
        # Demo presentation
        palette = PALETTES.get(args.style, PALETTES["neutral"])
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        add_title_slide(prs, "Industrial Automation", "Digital Transformation & Industry 4.0", palette)
        add_content_slide(prs, "Overview", [
            "SCADA/DCS integration",
            "MES implementation",
            "OPC UA connectivity",
            "Cybersecurity (IEC 62443)"
        ], "Speaker notes go here", palette)
        add_cta_slide(prs, "Let's Connect", ["email@company.com", "+7 (XXX) XXX-XX-XX"], palette)

        prs.save(args.output)
        print(f"Demo presentation saved: {args.output}")


if __name__ == "__main__":
    main()
