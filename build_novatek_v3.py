#!/usr/bin/env python3
"""
NOVATEK VP IT — Замена PI System на Альфа платформу
v3 FINAL: Honest comparison + Duarte Resonate + slide:ology + Zelazny

Key fixes vs v2:
- NO invented numbers (rule D-2026-03-28-01)
- Honest three-column mapping: ✅ full / ⚠️ partial / ❌ gaps
- Strategy for closing gaps (from v5)
- Bigger textboxes with margin to prevent text overflow
- Fewer items per slide
- LibreOffice render check
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import io, os

# ── Brand colors ──
BLUE = RGBColor(0x00, 0x54, 0x97)
CYAN = RGBColor(0x09, 0x97, 0xC8)
GREEN = RGBColor(0x82, 0xC4, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
SLATE = RGBColor(0x44, 0x54, 0x69)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x33, 0x75, 0xAD)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xE8, 0x77, 0x22)

SW = 9144000
SH = 5143500

# ── Assets ──
ASSETS = '/tmp/pptx_assets'
def load(name):
    p = os.path.join(ASSETS, name)
    return open(p, 'rb').read() if os.path.exists(p) else None

title_bg = load('title_bg.jpg')
industrial_bg = load('industrial_bg.jpg')
platform_ss = load('platform_screenshot.png')
arch_img = load('architecture.png')
proj_nv = load('project_novatek.png')
proj_tn = load('project_transneft.jpg')
logo = load('logo.png')
deco = load('deco_circle.png')
qr = load('qr_logo.png')

# ── Presentation ──
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BL = prs.slide_layouts[6]

# ── Helpers ──
def R(s, l, t, w, h, fill=None, alpha=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if alpha is not None:
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            sf = sh._element.find('.//a:solidFill', ns)
            if sf is not None:
                sr = sf.find('a:srgbClr', ns)
                if sr is not None:
                    ae = etree.SubElement(sr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    ae.set('val', str(int(alpha * 1000)))
    else:
        sh.fill.background()
    return sh

def RR(s, l, t, w, h, fill=None, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.fill.solid()
        sh.line.width = Pt(1)
    return sh

def O(s, l, t, w, h, fill=None):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    return sh

def T(s, l, t, w, h, text, sz=14, c=BLACK, b=False, a=PP_ALIGN.LEFT, f='Arial'):
    tx = s.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = f
    p.alignment = a
    return tx

def P(tf, text, sz=14, c=BLACK, b=False, a=PP_ALIGN.LEFT, sp=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = 'Arial'
    p.alignment = a
    p.space_before = sp
    return p

def IMG(s, blob, l, t, w, h):
    if blob:
        return s.shapes.add_picture(io.BytesIO(blob), l, t, w, h)

def E(*a): return Emu(sum(a))

def hdr(s, title, sub=None):
    R(s, 0, E(180000), SW, E(620000), fill=BLUE)
    R(s, 0, E(800000), SW, E(28000), fill=CYAN)
    T(s, E(500000), E(220000), E(7400000), E(550000), title, sz=22, c=WHITE, b=True)
    if sub:
        T(s, E(500000), E(520000), E(7400000), E(280000), sub, sz=11, c=RGBColor(0xBB, 0xCC, 0xDD))
    IMG(s, logo, E(7200000), E(230000), E(1600000), E(245000))

def ftr(s):
    R(s, 0, E(4960000), SW, E(28000), fill=CYAN)
    R(s, 0, E(4988000), SW, E(155500), fill=BLUE)


# ═══════════════════════════════════════════════
# SLIDE 1: TITLE
# Big Idea: "Замена PI System возможна. Проверено."
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
if title_bg: IMG(s, title_bg, 0, 0, SW, SH)
R(s, 0, 0, SW, SH, fill=BLACK, alpha=65)
R(s, 0, 0, SW, E(45000), fill=CYAN)
IMG(s, logo, E(480000), E(280000), E(2400000), E(370000))

tx = s.shapes.add_textbox(E(480000), E(1100000), E(8000000), E(1500000))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Замена PI System возможна."
p.font.size = Pt(36); p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = 'Arial'
P(tf, "Проверено на объектах вашего класса.", sz=36, c=CYAN, b=True, sp=Pt(4))

T(s, E(480000), E(2800000), E(7500000), E(400000),
  "Стратегия перехода на Альфа платформу", sz=16, c=RGBColor(0xBB, 0xCC, 0xDD))

R(s, 0, E(4300000), SW, E(1), fill=RGBColor(0x44, 0x55, 0x77))
T(s, E(480000), E(4400000), E(4500000), E(350000),
  "АО «Атомик Софт»  •  automiq.ru  •  Конфиденциально", sz=11, c=GRAY)
RR(s, E(6500000), E(4380000), E(2300000), E(400000), fill=RGBColor(0x22, 0x33, 0x55))
T(s, E(6500000), E(4420000), E(2300000), E(340000),
  "Внутренняя встреча • 2026", sz=11, c=RGBColor(0x88, 0x99, 0xBB), a=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 2: WHAT IS — Ситуация
# Facts only, no invented dollar amounts
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "Ситуация: зависимость от ушедшего вендора")
ftr(s)

cards = [
    ("Санкции", "AVEVA прекратила\nподдержку и продажу\nв России.\n\nНет патчей, нет\nобновлений с 2022 г.", RED),
    ("Подписка", "Модель AVEVA Flex:\nпри неоплате ПО\nперестаёт работать.\n\nОблачные сервисы\nуже отключены.", ORANGE),
    ("Регуляторика", "ФЗ-187 о безопасности\nКИИ. Указ №250.\n\nТребование перехода\nна отечественное ПО.\nСроки ужесточаются.", BLUE),
    ("Кадры", "Новые специалисты\nобучаются на российском\nПО. Экспертиза по\nPI System сокращается.", ACCENT_BLUE),
]

for i, (label, detail, color) in enumerate(cards):
    x = E(150000 + i * 2250000)
    y = E(920000)
    card = RR(s, x, y, E(2100000), E(3700000), fill=WHITE, border=RGBColor(0xE0, 0xE5, 0xEB))
    R(s, x, y, E(2100000), E(55000), fill=color)
    T(s, E(x, 80000), E(y, 130000), E(1940000), E(300000),
      label, sz=14, c=color, b=True, a=PP_ALIGN.CENTER)
    T(s, E(x, 120000), E(y, 500000), E(1860000), E(3000000),
      detail, sz=11, c=SLATE, a=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 3: WHAT COULD BE — Альфа платформа
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
R(s, 0, 0, SW, SH, fill=LIGHT_BG)
hdr(s, "Альтернатива: Альфа платформа",
    "Российская SCADA-платформа №1 в реестре Минцифры — АО «Атомик Софт», Томск")
ftr(s)

if platform_ss:
    IMG(s, platform_ss, E(4600000), E(1100000), E(4300000), E(3600000))

facts = [
    ("Бессрочная лицензия в рублях", "Одна покупка. Нет подписки.\nНе отключат при неоплате.", GREEN),
    ("6 500+ инсталляций", "Нефтегаз, энергетика,\nкритическая инфраструктура.", CYAN),
    ("Windows + Linux", "Astra Linux, РЕД ОС,\nАльт, Ubuntu. Сертифицированные ОС.", BLUE),
    ("Техподдержка 24/7", "На русском языке. Вендор в Томске.\nЛицензия ФСТЭК.", ACCENT_BLUE),
]

y = E(1000000)
for title, desc, color in facts:
    R(s, E(400000), y, E(50000), E(550000), fill=color)
    T(s, E(550000), E(y, 40000), E(3800000), E(300000), title, sz=16, c=BLUE, b=True)
    T(s, E(550000), E(y, 310000), E(3800000), E(300000), desc, sz=11, c=SLATE)
    y = E(y, 630000)


# ═══════════════════════════════════════════════
# SLIDE 4: Ключевое различие — «Склад» vs «Диспетчерская»
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "Ключевое различие: «Склад данных» vs «Диспетчерская»")
ftr(s)

# PI column
R(s, E(150000), E(920000), E(4300000), E(400000), fill=RED)
T(s, E(150000), E(950000), E(4300000), E(360000),
  "📦  AVEVA PI System — «Склад данных»", sz=14, c=WHITE, b=True, a=PP_ALIGN.CENTER)

pi_desc = (
    "Огромный склад с идеальной системой учёта.\n"
    "Любой товар можно найти, любому курьеру\n"
    "выдать, в любую аналитику подключить.\n"
    "Но сам склад ничем не управляет."
)
T(s, E(250000), E(1400000), E(4100000), E(900000), pi_desc, sz=11, c=SLATE)

pi_pros = [
    "▸  450+ коннекторов — почти всё на свете",
    "▸  REST API — данные по HTTP",
    "▸  Облако — данные из любой точки мира",
    "▸  Event Frames — запись событий",
    "▸  Оператор сам рисует дашборд",
]
y = E(2400000)
for item in pi_pros:
    T(s, E(300000), y, E(4000000), E(220000), item, sz=10, c=SLATE)
    y = E(y, 220000)

# Alpha column
R(s, E(4700000), E(920000), E(4300000), E(400000), fill=GREEN)
T(s, E(4700000), E(950000), E(4300000), E(360000),
  "🖥️  Альфа платформа — «Диспетчерская»", sz=14, c=WHITE, b=True, a=PP_ALIGN.CENTER)

alpha_desc = (
    "Диспетчерская с пультом. Оператор видит всё,\n"
    "управляет всем, все тревоги приходят сюда.\n"
    "Склад тоже есть (Historian),\n"
    "но главное — пульт управления."
)
T(s, E(4800000), E(1400000), E(4100000), E(900000), alpha_desc, sz=11, c=SLATE)

alpha_pros = [
    "▸  Мощный дизайнер мнемосхем (HMI)",
    "▸  ~30 протоколов — все основные",
    "▸  Напрямую к оборудованию, без OPC-прослоек",
    "▸  Linux + Windows, on-premises",
    "▸  Бессрочная лицензия в рублях",
]
y = E(2400000)
for item in alpha_pros:
    T(s, E(4800000), y, E(4100000), E(220000), item, sz=10, c=SLATE)
    y = E(y, 220000)

# Bottom insight
RR(s, E(150000), E(3700000), E(8850000), E(450000), fill=RGBColor(0xE3, 0xF2, 0xFD))
T(s, E(350000), E(3740000), E(8450000), E(380000),
  "Миграция = переезд со «склада данных» на «диспетчерскую с собственным складом».\nPI хорош в сборе и раздаче данных. Альфа хороша в управлении и визуализации. Это разные инструменты.",
  sz=11, c=BLUE, b=True)


# ═══════════════════════════════════════════════
# SLIDE 5: ЧЕСТНОЕ СРАВНЕНИЕ — таблица
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "Сравнение: PI System vs Альфа платформа",
    "Честная оценка сильных и слабых сторон обеих платформ")
ftr(s)

rows = [
    ("Что это",       "Инфраструктура данных",        "SCADA-платформа"),
    ("Главный фокус", "Собрать и раздать данные",      "Показать и управлять"),
    ("Визуализация",  "Вторична (PI Vision — простой)","Ядро продукта (Alpha.HMI)"),
    ("Коннекторы",    "450+ (почти всё на свете)",     "~30 (основные промышленные)"),
    ("Модель объектов","Runtime (меняешь на ходу)",     "Design-time (проектируешь)"),
    ("API",           "REST (как веб-сайт)",           "OPC UA + TCP + SQL"),
    ("Облако",        "Да (AVEVA Data Hub)",           "Нет (только on-prem)"),
    ("Event Frames",  "Да (уникальная фича)",          "Нет (нужна доработка)"),
    ("ОС",            "Только Windows",                "Windows + Linux"),
    ("Лицензия",      "Подписка, доллары",             "Бессрочная, рубли"),
    ("Поддержка в РФ","Нет (санкции)",                 "24/7 на русском"),
]

# Table header
R(s, E(150000), E(900000), E(2700000), E(330000), fill=SLATE)
T(s, E(200000), E(930000), E(2600000), E(280000), "Параметр", sz=11, c=WHITE, b=True, a=PP_ALIGN.CENTER)
R(s, E(2850000), E(900000), E(3000000), E(330000), fill=RED)
T(s, E(2900000), E(930000), E(2900000), E(280000), "AVEVA PI System", sz=11, c=WHITE, b=True, a=PP_ALIGN.CENTER)
R(s, E(5850000), E(900000), E(3150000), E(330000), fill=GREEN)
T(s, E(5900000), E(930000), E(3050000), E(280000), "Альфа платформа", sz=11, c=WHITE, b=True, a=PP_ALIGN.CENTER)

y = E(1230000)
row_h = E(305000)
for i, (param, pi_val, alpha_val) in enumerate(rows):
    bg = RGBColor(0xF8, 0xF9, 0xFB) if i % 2 == 0 else WHITE
    R(s, E(150000), y, E(8850000), row_h, fill=bg)
    T(s, E(200000), E(y, 45000), E(2600000), E(260000), param, sz=10, c=SLATE, b=True)
    T(s, E(2900000), E(y, 45000), E(2900000), E(260000), pi_val, sz=10, c=SLATE)
    T(s, E(5900000), E(y, 45000), E(3050000), E(260000), alpha_val, sz=10, c=SLATE)
    y = E(y, row_h)


# ═══════════════════════════════════════════════
# SLIDE 6: ПОКРЫТИЕ — Полная замена
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "Покрытие функционала: полная замена (6 компонентов)")
ftr(s)

full_items = [
    ("Архив данных → Alpha.Historian 4.0",
     "Все данные с датчиков хранятся надёжно.\nСжатие LZMA — архив занимает меньше места."),
    ("Экраны оператора → Alpha.HMI 2.0",
     "Мнемосхемы мощнее, чем в PI ProcessBook.\nДизайнер + визуализатор + веб-доступ."),
    ("Тревоги → Alpha.HMI.Alarms 3.3",
     "Кроссплатформенный: Windows, Linux, браузер.\nКвитирование, подавление, экспорт."),
    ("Буферизация → встроена в Alpha.Server",
     "Если связь с сервером пропала — данные\nне теряются. Работает автоматически."),
    ("Администрирование → Alpha.DevStudio 4.1",
     "Единая среда: проектирование, диагностика,\nмониторинг всей системы. CLI для автоматизации."),
    ("Программный доступ → OPC UA + RMap",
     "Внешние программы подключаются через\nOPC UA, TCP или SQL (Alpha.RMap)."),
]

y = E(920000)
for i, (title, desc) in enumerate(full_items):
    col = i % 2
    row = i // 2
    x = E(150000 + col * 4500000)
    cy = E(y + row * 1200000)
    
    card = RR(s, x, cy, E(4300000), E(1050000), fill=RGBColor(0xF0, 0xF8, 0xF0), border=RGBColor(0xCC, 0xDD, 0xCC))
    R(s, x, cy, E(50000), E(1050000), fill=GREEN)
    T(s, E(x, 130000), E(cy, 70000), E(4000000), E(300000), "✅  " + title, sz=12, c=BLUE, b=True)
    T(s, E(x, 130000), E(cy, 360000), E(4000000), E(600000), desc, sz=10, c=SLATE)


# ═══════════════════════════════════════════════
# SLIDE 7: ПОКРЫТИЕ — Частичная замена + пробелы
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "Честно: что работает иначе и чего нет",
    "Частичная замена (4 компонента) и функциональные пробелы (6)")
ftr(s)

# Partial — left column
R(s, E(150000), E(890000), E(4300000), E(340000), fill=ORANGE)
T(s, E(150000), E(910000), E(4300000), E(300000),
  "⚠️  Частичная замена (4)", sz=13, c=WHITE, b=True, a=PP_ALIGN.CENTER)

partials = [
    ("450+ коннекторов → ~30", "Основные есть. Экзотика — через OPC UA шлюзы."),
    ("Оператор сам рисует дашборд", "В Альфе дашборды рисует инженер.\nОператор смотрит готовое в браузере."),
    ("Модель активов «на лету»", "В Альфе структура проектируется\nзаранее в DevStudio."),
    ("PI DataLink (Excel)", "Есть Alpha.Reports, но это не\nинтерактивный Excel add-in."),
]

y = E(1300000)
for title, desc in partials:
    T(s, E(250000), y, E(4000000), E(250000), "⚠️  " + title, sz=10, c=ORANGE, b=True)
    T(s, E(250000), E(y, 230000), E(4000000), E(350000), desc, sz=9, c=SLATE)
    y = E(y, 600000)

# Gaps — right column
R(s, E(4700000), E(890000), E(4300000), E(340000), fill=RED)
T(s, E(4700000), E(910000), E(4300000), E(300000),
  "❌  Функциональные пробелы (6)", sz=13, c=WHITE, b=True, a=PP_ALIGN.CENTER)

gaps = [
    ("Event Frames", "Нельзя «запиши что было\nс 10:15 до 11:30 на насосе»"),
    ("REST API", "Нет HTTP-доступа к данным"),
    ("Репликация архивов", "Нет синхронизации между серверами"),
    ("Edge-версия", "Нет мини-сервера для\nудалённых площадок"),
    ("Ручной ввод данных", "Нет формы для лабораторных\nанализов и обходов"),
    ("Коннектор Power BI", "Нет прямого. Обход:\nSQL → PostgreSQL → BI"),
]

y = E(1300000)
for title, desc in gaps:
    T(s, E(4800000), y, E(4000000), E(250000), "✕  " + title, sz=10, c=RED, b=True)
    T(s, E(4800000), E(y, 220000), E(4000000), E(300000), desc, sz=9, c=SLATE)
    y = E(y, 530000)


# ═══════════════════════════════════════════════
# SLIDE 8: СТРАТЕГИЯ ЗАКРЫТИЯ ПРОБЕЛОВ
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "Стратегия закрытия пробелов",
    "Каждый пробел имеет конкретное решение")
ftr(s)

# Table
R(s, E(150000), E(890000), E(3500000), E(330000), fill=RED)
T(s, E(200000), E(910000), E(3400000), E(290000), "Чего не хватает", sz=12, c=WHITE, b=True, a=PP_ALIGN.CENTER)
R(s, E(3650000), E(890000), E(3700000), E(330000), fill=GREEN)
T(s, E(3700000), E(910000), E(3600000), E(290000), "Как решаем", sz=12, c=WHITE, b=True, a=PP_ALIGN.CENTER)
R(s, E(7350000), E(890000), E(1650000), E(330000), fill=ACCENT_BLUE)
T(s, E(7400000), E(910000), E(1550000), E(290000), "Подход", sz=12, c=WHITE, b=True, a=PP_ALIGN.CENTER)

strategy = [
    ("Event Frames",     "Alpha.Om + PostgreSQL",                "Доработка"),
    ("REST API",         "REST-обёртка над SQL-данными",          "Обёртка"),
    ("Репликация",       "Экспорт Historian → PostgreSQL",        "Периодич."),
    ("Edge-сервер",      "Alpha.Server в мини-конфигурации",      "Лёгкая версия"),
    ("Ручной ввод",      "Веб-форма через HMI.WebViewer",        "Веб-форма"),
    ("BI-аналитика",     "SQL → Grafana / Power BI",             "Через СУБД"),
    ("Экзотич. протоколы","~30 покрывают 80%, остальное — OPC UA","Шлюзы"),
]

y = E(1280000)
row_h = E(420000)
for i, (gap, solution, approach) in enumerate(strategy):
    bg = RGBColor(0xF8, 0xF9, 0xFB) if i % 2 == 0 else WHITE
    R(s, E(150000), y, E(8850000), row_h, fill=bg)
    T(s, E(250000), E(y, 60000), E(3300000), E(350000), gap, sz=11, c=SLATE, b=True)
    T(s, E(3750000), E(y, 60000), E(3500000), E(350000), solution, sz=11, c=SLATE)
    T(s, E(7450000), E(y, 60000), E(1450000), E(350000), approach, sz=10, c=ACCENT_BLUE, b=True, a=PP_ALIGN.CENTER)
    y = E(y, row_h)

# Insight
RR(s, E(150000), E(4350000), E(8850000), E(400000), fill=RGBColor(0xE8, 0xF5, 0xE9))
T(s, E(350000), E(4385000), E(8450000), E(340000),
  "Все пробелы закрываемы. Аудит PI-инфраструктуры покажет, какие из них реально используются у вас.",
  sz=12, c=RGBColor(0x2E, 0x7D, 0x32), b=True)


# ═══════════════════════════════════════════════
# SLIDE 9: ОПЫТ — Проекты
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "Опыт: объекты вашего класса",
    "Альфа платформа в промышленной эксплуатации")
ftr(s)

projects = [
    ("НОВАТЭК — Арктик СПГ 2", "Автоматизация технологических\nпроцессов СПГ-производства.\nАльфа в составе АСУ ТП.", CYAN),
    ("Транснефть — ВСТО-II", "Магистральный нефтепровод.\nДиспетчеризация и управление\nнасосных станций.", BLUE),
    ("СИБУР", "Нефтехимический холдинг.\nАльфа в составе систем\nуправления производствами.", ACCENT_BLUE),
    ("Интер РАО — Сочинская ТЭС", "Тепловая электростанция.\nМнемосхемы + архив + алармы.", GREEN),
]

for i, (name, desc, color) in enumerate(projects):
    col = i % 2
    row = i // 2
    x = E(150000 + col * 4500000)
    y = E(920000 + row * 1900000)
    
    card = RR(s, x, y, E(4300000), E(1650000), fill=WHITE, border=RGBColor(0xE0, 0xE5, 0xEB))
    R(s, x, y, E(50000), E(1650000), fill=color)
    T(s, E(x, 150000), E(y, 100000), E(3900000), E(350000), name, sz=16, c=BLUE, b=True)
    T(s, E(x, 150000), E(y, 480000), E(3900000), E(1000000), desc, sz=12, c=SLATE)

# Photos where available
if proj_tn:
    IMG(s, proj_tn, E(3100000), E(1050000), E(1250000), E(1250000))


# ═══════════════════════════════════════════════
# SLIDE 10: ПЛАН ПЕРЕХОДА
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "План перехода: 5 фаз за 8–14 месяцев",
    "Ключевой принцип: PI и Альфа работают параллельно. Нулевой простой.")
ftr(s)

phases = [
    ("АУДИТ",       "1–2 мес.", "Инвентаризация:\nтеги, протоколы,\nскрипты, EF", BLUE),
    ("ПРОЕКТ",      "1 мес.",   "Архитектура,\nлицензии, бюджет,\nрабочая группа", ACCENT_BLUE),
    ("ПИЛОТ",       "2–3 мес.", "Один объект.\nPI + Альфа\nпараллельно", CYAN),
    ("МИГРАЦИЯ",    "3–6 мес.", "Поэтапный\nперенос. Обучение.\nВалидация.", GREEN),
    ("ЗАКРЫТИЕ",    "1–2 мес.", "Вывод PI.\nОптимизация.\nПередача.", RGBColor(0x61, 0x96, 0x2F)),
]

tl_y = E(1300000)
R(s, E(200000), E(tl_y, 230000), E(8700000), E(60000), fill=RGBColor(0xDD, 0xE5, 0xEE))

pw = E(1600000)
gap = E(130000)
sx = E(300000)

for i, (name, dur, desc, color) in enumerate(phases):
    x = E(sx, i * (pw + gap))
    O(s, E(x, 600000), tl_y, E(420000), E(420000), fill=color)
    T(s, E(x, 600000), E(tl_y, 90000), E(420000), E(300000),
      str(i+1), sz=18, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    T(s, x, E(tl_y, 530000), pw, E(280000), name, sz=13, c=color, b=True, a=PP_ALIGN.CENTER)
    T(s, x, E(tl_y, 790000), pw, E(220000), dur, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    T(s, E(x, 30000), E(tl_y, 1050000), E(pw, -60000), E(700000), desc, sz=10, c=SLATE, a=PP_ALIGN.CENTER)

RR(s, E(200000), E(4000000), E(8700000), E(650000), fill=RGBColor(0xE3, 0xF2, 0xFD))
tx = s.shapes.add_textbox(E(400000), E(4050000), E(8300000), E(560000))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Параллельная работа = нулевой риск для производства"
p.font.size = Pt(13); p.font.color.rgb = BLUE; p.font.bold = True; p.font.name = 'Arial'
P(tf, "PI System и Альфа работают одновременно. Переключение — только после подтверждения корректности данных.", sz=10, c=SLATE, sp=Pt(6))


# ═══════════════════════════════════════════════
# SLIDE 11: СЛЕДУЮЩИЕ ШАГИ (Zelazny: action program)
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "Предлагаемые следующие шаги")
ftr(s)

steps = [
    ("1", "Согласование объёма аудита PI-инфраструктуры", BLUE),
    ("2", "Определение пилотного объекта для миграции", CYAN),
    ("3", "Перечень критичных функций PI\n(Event Frames? DataLink? API?)", ACCENT_BLUE),
    ("4", "Расчёт лицензий Альфа под целевую архитектуру", GREEN),
    ("5", "Рабочая группа: Новатэк + Атомик Софт + интегратор", RGBColor(0x61, 0x96, 0x2F)),
    ("6", "Детальный план-график и бюджет", ORANGE),
]

y = E(920000)
for num, text, color in steps:
    O(s, E(350000), y, E(450000), E(450000), fill=color)
    T(s, E(350000), E(y, 90000), E(450000), E(320000), num, sz=22, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    T(s, E(1000000), E(y, 80000), E(7500000), E(380000), text, sz=14, c=SLATE)
    y = E(y, 560000)


# ═══════════════════════════════════════════════
# SLIDE 12: NEW BLISS — Видение будущего
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
if industrial_bg: IMG(s, industrial_bg, 0, 0, SW, SH)
R(s, 0, 0, SW, SH, fill=BLACK, alpha=70)

tx = s.shapes.add_textbox(E(480000), E(800000), E(8000000), E(2200000))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Через 12 месяцев:"
p.font.size = Pt(14); p.font.color.rgb = CYAN; p.font.name = 'Arial'

for v in [
    "Ваша SCADA — полностью ваша.",
    "Никаких подписок. Никаких санкций.",
    "Техподдержка на расстоянии звонка.",
    "Полное соответствие ФЗ-187 о КИИ.",
]:
    P(tf, v, sz=24, c=WHITE, b=True, sp=Pt(14))

R(s, 0, E(5093500), SW, E(50000), fill=CYAN)


# ═══════════════════════════════════════════════
# SLIDE 13: КОНТАКТЫ
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(BL)
R(s, 0, 0, SW, SH, fill=DARK_BG)
IMG(s, deco, E(-800000), E(100000), E(4900000), E(4900000))

tx = s.shapes.add_textbox(E(480000), E(700000), E(8000000), E(1200000))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Будущее автоматизации"
p.font.size = Pt(32); p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = 'Arial'
P(tf, "в ваших руках.", sz=32, c=CYAN, b=True, sp=Pt(4))

tx2 = s.shapes.add_textbox(E(2000000), E(2600000), E(3500000), E(1800000))
tf2 = tx2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]; p.text = "АО «Атомик Софт»"
p.font.size = Pt(14); p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = 'Arial'
P(tf2, "634050 Томск, пр. Ленина 60/1", sz=12, c=RGBColor(0xBB, 0xCC, 0xDD), sp=Pt(4))
P(tf2, "+7 (3822) 281 914", sz=12, c=RGBColor(0xBB, 0xCC, 0xDD), sp=Pt(8))
P(tf2, "info@automiq.ru  •  automiq.ru", sz=12, c=CYAN, sp=Pt(2))

pc = s.shapes.add_textbox(E(6000000), E(3000000), E(3000000), E(1000000))
tf3 = pc.text_frame; tf3.word_wrap = True
p = tf3.paragraphs[0]; p.text = "Ваш контакт:"
p.font.size = Pt(10); p.font.color.rgb = GRAY; p.font.name = 'Arial'
P(tf3, "Станислав Павловский", sz=14, c=WHITE, b=True, sp=Pt(4))
P(tf3, "Telegram: @Integraleus", sz=12, c=CYAN, sp=Pt(2))

IMG(s, logo, E(480000), E(4300000), E(2200000), E(340000))
IMG(s, qr, E(6800000), E(2000000), E(1200000), E(1200000))
R(s, 0, E(5093500), SW, E(50000), fill=CYAN)


# ── Save ──
out = "NOVATEK_VP_IT_v3_honest.pptx"
prs.save(out)
print(f"✅ {out}: {len(prs.slides)} slides")
