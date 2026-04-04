#!/usr/bin/env python3
"""
NOVATEK VP IT — Презентация замены PI System на Альфа платформу
v3: Rebuilt using Duarte Resonate (Sparkline), slide:ology, Zelazny frameworks

Structure (Sparkline):
  1. Title — anchor, credibility
  2. WHAT IS: current reality (pain) — "you depend on a vendor who left"
  3. TURNING POINT 1 / BIG IDEA: "the switch is possible, proven, and saves money"
  4. WHAT COULD BE: vision of control — Альфа capabilities
  5. WHAT IS ↔ WHAT COULD BE: functional mapping (PI → Alpha) 
  6. S.T.A.R. MOMENT: shocking economics comparison (one powerful visual)
  7. WHAT COULD BE: proof — real projects (Novatek already uses it!)
  8. WHAT IS: risks & fears acknowledged (inokulation per Duarte)
  9. WHAT COULD BE: how risks are mitigated
  10. TURNING POINT 2 / CALL TO ACTION: concrete next steps
  11. NEW BLISS: vision of the future
  12. Contacts

Design principles (slide:ology):
  - One slide = one idea
  - Signal > noise, max whitespace
  - Contrast is intentional
  - Real photos from Atomicsoft template
  - Branded Automiq palette
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import io, os

# ── Brand colors (Automiq brandbook) ──
BLUE = RGBColor(0x00, 0x54, 0x97)       # Primary
CYAN = RGBColor(0x09, 0x97, 0xC8)       # Accent 1  
GREEN = RGBColor(0x82, 0xC4, 0x44)      # Accent 2
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_SLATE = RGBColor(0x44, 0x54, 0x69)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
MED_GRAY = RGBColor(0x6E, 0x6E, 0x6E)
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x33, 0x75, 0xAD)
RED_ACCENT = RGBColor(0xC6, 0x28, 0x28)
WARM_ORANGE = RGBColor(0xE8, 0x77, 0x22)

# ── Dimensions ──
SW = 9144000
SH = 5143500

# ── Load assets ──
ASSETS = '/tmp/pptx_assets'
def load_asset(name):
    path = os.path.join(ASSETS, name)
    if os.path.exists(path):
        return open(path, 'rb').read()
    return None

title_bg = load_asset('title_bg.jpg')
industrial_bg = load_asset('industrial_bg.jpg')
platform_screenshot = load_asset('platform_screenshot.png')
architecture_img = load_asset('architecture.png')
project_novatek = load_asset('project_novatek.png')
project_transneft = load_asset('project_transneft.jpg')
logo = load_asset('logo.png')
deco_circle = load_asset('deco_circle.png')
qr_logo = load_asset('qr_logo.png')

# ── Create presentation ──
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]

# ── Helper functions ──
def rect(slide, l, t, w, h, fill=None, alpha_pct=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if alpha_pct is not None:
            # Set transparency via XML on the shape's spPr solidFill
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            spPr = shape._element.find('.//a:solidFill', nsmap)
            if spPr is not None:
                srgb = spPr.find('a:srgbClr', nsmap)
                if srgb is not None:
                    alpha_elem = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    alpha_elem.set('val', str(int(alpha_pct * 1000)))
    else:
        shape.fill.background()
    return shape

def rrect(slide, l, t, w, h, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    return shape

def oval(slide, l, t, w, h, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    return shape

def tb(slide, l, t, w, h, text, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font='Arial', anchor=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_p(tf, text, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font='Arial', space_before=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_before = space_before
    return p

def pic(slide, blob, l, t, w, h):
    if blob:
        return slide.shapes.add_picture(io.BytesIO(blob), l, t, w, h)
    return None

def header_bar(slide, title, subtitle=None):
    """Consistent blue header with logo"""
    rect(slide, 0, Emu(180000), SW, Emu(600000), fill=BLUE)
    rect(slide, 0, Emu(780000), SW, Emu(28000), fill=CYAN)
    tb(slide, Emu(550000), Emu(215000), Emu(7500000), Emu(520000),
       title, size=24, color=WHITE, bold=True)
    if subtitle:
        tb(slide, Emu(550000), Emu(510000), Emu(7500000), Emu(260000),
           subtitle, size=12, color=RGBColor(0xBB, 0xCC, 0xDD))
    pic(slide, logo, Emu(7300000), Emu(225000), Emu(1500000), Emu(230000))

def footer_accent(slide):
    """Thin footer line"""
    rect(slide, 0, Emu(4950000), SW, Emu(28000), fill=CYAN)
    rect(slide, 0, Emu(4978000), SW, Emu(165500), fill=BLUE)

def E(*args):
    """Shorthand for Emu arithmetic"""
    return Emu(sum(args))


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE (Zelazny: Purpose, Importance, Preview)
# Full-bleed industrial photo + overlay, Big Idea visible
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)

# Background photo
if title_bg:
    pic(s, title_bg, 0, 0, SW, SH)
# Dark overlay
rect(s, 0, 0, SW, SH, fill=BLACK, alpha_pct=65)

# Top accent line (brand)
rect(s, 0, 0, SW, Emu(45000), fill=CYAN)

# Logo
pic(s, logo, Emu(480000), Emu(250000), Emu(2400000), Emu(370000))

# Main title — Big Idea as statement (Duarte: full sentence with stakes)
txBox = s.shapes.add_textbox(Emu(480000), Emu(1100000), Emu(7800000), Emu(1400000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Замена PI System возможна."
p.font.size = Pt(36)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = 'Arial'
# Second line — stakes
p2 = add_p(tf, "Проверено на объектах вашего класса.", size=36, color=CYAN, bold=True, space_before=Pt(2))

# Subtitle — what's at stake (Duarte: raise pain probability)
tb(s, Emu(480000), Emu(2700000), Emu(7000000), Emu(400000),
   "Стратегия перехода на Альфа платформу для Новатэк",
   size=16, color=RGBColor(0xBB, 0xCC, 0xDD))

# Bottom bar
rect(s, 0, Emu(4300000), SW, Emu(1), fill=RGBColor(0x44, 0x55, 0x77))
tb(s, Emu(480000), Emu(4400000), Emu(4500000), Emu(300000),
   "АО «Атомик Софт»  •  automiq.ru  •  Конфиденциально",
   size=11, color=MED_GRAY)

# Badge
rrect(s, Emu(6500000), Emu(4380000), Emu(2300000), Emu(380000), fill=RGBColor(0x22, 0x33, 0x55))
tb(s, Emu(6500000), Emu(4410000), Emu(2300000), Emu(320000),
   "Внутренняя встреча • 2026", size=11, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: WHAT IS — Current Reality (Sparkline: establish pain)
# "You depend on a vendor who abandoned you"
# One powerful visual: three danger icons, minimal text
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "Реальность: зависимость от ушедшего вендора")
footer_accent(s)

# Three large risk cards — each with BIG number (S.T.A.R. shocking stats per Duarte)
cards = [
    ("$540K", "уже заплачено\nза подписку без\nподдержки с 2022 г.", "Санкции", RED_ACCENT),
    ("$108K", "уходит каждый год\nна подписку PI System\nбез обновлений", "Подписка", WARM_ORANGE),
    ("2027", "крайний срок\nперехода КИИ на\nотечественное ПО", "Регуляторика", BLUE),
]

for i, (big_num, detail, label, color) in enumerate(cards):
    x = Emu(200000 + i * 3000000)
    y = Emu(920000)
    
    # Card
    card = rrect(s, x, y, Emu(2800000), Emu(3700000), fill=WHITE)
    card.line.color.rgb = RGBColor(0xE0, 0xE5, 0xEB)
    card.line.width = Pt(1)
    
    # Color accent at top
    rect(s, x, y, Emu(2800000), Emu(60000), fill=color)
    
    # Label
    tb(s, E(x, 100000), E(y, 140000), Emu(2600000), Emu(250000),
       label, size=12, color=MED_GRAY, bold=True, align=PP_ALIGN.CENTER)
    
    # BIG number (Zelazny: data visual, slide:ology: focal point)
    tb(s, E(x, 100000), E(y, 450000), Emu(2600000), Emu(800000),
       big_num, size=48, color=color, bold=True, align=PP_ALIGN.CENTER, font='Arial')
    
    # Detail text
    tb(s, E(x, 200000), E(y, 1350000), Emu(2400000), Emu(700000),
       detail, size=12, color=DARK_SLATE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: WHAT COULD BE — The Solution (Duarte: contrast jump)
# "Imagine: full control, Russian language, your own license"
# Zelazny: recommendation UPFRONT
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)

# Light background
rect(s, 0, 0, SW, SH, fill=LIGHT_BG)
header_bar(s, "Альтернатива: полный контроль в ваших руках",
           "Альфа платформа — российская SCADA №1 в реестре Минцифры")
footer_accent(s)

# Platform screenshot (real, from template)
if platform_screenshot:
    # Screenshot on the right, faded into BG
    pic(s, platform_screenshot, Emu(4700000), Emu(1100000), Emu(4200000), Emu(3500000))

# Key facts — large, spaced (slide:ology: whitespace, one idea per block)
facts = [
    ("Бессрочная лицензия", "Одна покупка в рублях.\nНет подписки. Нет отключения.", GREEN),
    ("6 500+ инсталляций", "В т.ч. нефтегаз, энергетика,\nкритическая инфраструктура.", CYAN),
    ("Полная поддержка", "24/7, на русском, вендор\nв г. Томск. ФСТЭК.", BLUE),
]

y_pos = Emu(1000000)
for title, desc, color in facts:
    # Accent bar
    rect(s, Emu(450000), y_pos, Emu(50000), Emu(520000), fill=color)
    # Title
    tb(s, Emu(600000), E(y_pos, 30000), Emu(3800000), Emu(280000),
       title, size=18, color=BLUE, bold=True)
    # Description
    tb(s, Emu(600000), E(y_pos, 280000), Emu(3800000), Emu(300000),
       desc, size=11, color=DARK_SLATE)
    y_pos = E(y_pos, 600000)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: WHAT IS ↔ WHAT COULD BE — Functional Mapping
# (Duarte: contrast, slide:ology: comparison diagram)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "Функции PI System → Альфа платформа",
           "Полный маппинг: что заменяется напрямую, что требует адаптации")
footer_accent(s)

# Two-column mapping with green checkmarks and orange arrows
mappings = [
    ("PI Data Archive",       "Alpha.Historian 4.0",         "✅"),
    ("PI ProcessBook",        "Alpha.HMI 2.0",              "✅"),
    ("PI Vision (веб)",       "Alpha.HMI.WebViewer 2.0",    "✅"),
    ("PI Notifications",      "Alpha.HMI.Alarms 3.3",       "✅"),
    ("PI Buffer Subsystem",   "Alpha.Server (буферизация)",  "✅"),
    ("PI SDK / AF SDK",       "Alpha.Server OPC UA + RMap",  "✅"),
    ("PI System Management",  "Alpha.DevStudio + Security",  "✅"),
]

adaptations = [
    ("PI Asset Framework",    "Alpha.Domain + Alpha.Om",     "🔄"),
    ("PI Analytics / EF",     "Alpha.Om + Alpha.Reports",    "🔄"),
    ("PI Integrator for BA",  "Alpha.RMap (SQL-доступ)",     "🔄"),
]

# Table header
rect(s, Emu(200000), Emu(870000), Emu(4200000), Emu(300000), fill=GREEN)
tb(s, Emu(200000), Emu(900000), Emu(4200000), Emu(260000),
   "✅  Прямая замена (7 из 10 компонентов)", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

rect(s, Emu(4700000), Emu(870000), Emu(4200000), Emu(300000), fill=WARM_ORANGE)
tb(s, Emu(4700000), Emu(900000), Emu(4200000), Emu(260000),
   "🔄  Требует адаптации (3 компонента)", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Left column — direct replacements
y = Emu(1250000)
for pi_name, alpha_name, icon in mappings:
    # Row stripe
    if mappings.index((pi_name, alpha_name, icon)) % 2 == 0:
        rrect(s, Emu(200000), y, Emu(4200000), Emu(310000), fill=RGBColor(0xF0, 0xF8, 0xF0))
    
    tb(s, Emu(280000), E(y, 50000), Emu(1800000), Emu(240000),
       pi_name, size=10, color=MED_GRAY)
    tb(s, Emu(2100000), E(y, 50000), Emu(300000), Emu(240000),
       "→", size=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    tb(s, Emu(2450000), E(y, 50000), Emu(1900000), Emu(240000),
       alpha_name, size=10, color=DARK_SLATE, bold=True)
    y = E(y, 310000)

# Right column — adaptations
y2 = Emu(1250000)
for pi_name, alpha_name, icon in adaptations:
    if adaptations.index((pi_name, alpha_name, icon)) % 2 == 0:
        rrect(s, Emu(4700000), y2, Emu(4200000), Emu(310000), fill=RGBColor(0xFE, 0xF7, 0xF0))
    
    tb(s, Emu(4780000), E(y2, 50000), Emu(1800000), Emu(240000),
       pi_name, size=10, color=MED_GRAY)
    tb(s, Emu(6600000), E(y2, 50000), Emu(300000), Emu(240000),
       "→", size=12, color=WARM_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, Emu(6950000), E(y2, 50000), Emu(1900000), Emu(240000),
       alpha_name, size=10, color=DARK_SLATE, bold=True)
    y2 = E(y2, 310000)

# Insight below adaptations
tb(s, Emu(4780000), E(y2, 100000), Emu(4000000), Emu(600000),
   "Адаптация — не проблема.\nЭто 3 компонента из 10.\nАудит покажет точный объём работ.",
   size=10, color=WARM_ORANGE, bold=True)

# Bottom: Zelazny "message" reinforcement
rrect(s, Emu(200000), Emu(4400000), Emu(8700000), Emu(380000), fill=RGBColor(0xE8, 0xF5, 0xE9))
tb(s, Emu(400000), Emu(4440000), Emu(8300000), Emu(300000),
   "70% функций PI — прямая замена. 30% — адаптация, а не потеря. Нулевой простой при переходе.",
   size=12, color=RGBColor(0x2E, 0x7D, 0x32), bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: S.T.A.R. MOMENT — Shocking Economics (one visual)
# (Duarte: "Shocking Statistics" — pause and let them sink in)
# (Zelazny: data chart — bar comparison, one focal point)
# (slide:ology: one idea, max signal/noise)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SW, SH, fill=DARK_BG)

# Title — minimal
tb(s, Emu(480000), Emu(250000), Emu(8000000), Emu(400000),
   "Стоимость бездействия", size=28, color=WHITE, bold=True)
tb(s, Emu(480000), Emu(600000), Emu(8000000), Emu(250000),
   "Совокупная стоимость владения за 5 лет", size=14, color=MED_GRAY)

# Two large bars — visual comparison (slide:ology: diagram type = item comparison → bar)
# PI System bar (red, LONG)
bar_top = Emu(1300000)
bar_h = Emu(700000)

# PI bar — full width = ~₽54M
pi_width = Emu(7200000)
rect(s, Emu(480000), bar_top, pi_width, bar_h, fill=RED_ACCENT)
tb(s, Emu(600000), E(bar_top, 100000), Emu(3000000), Emu(250000),
   "PI System", size=16, color=WHITE, bold=True)
tb(s, Emu(600000), E(bar_top, 350000), Emu(4000000), Emu(250000),
   "подписка + отсутствие поддержки + валютный риск", size=10, color=RGBColor(0xFF, 0xCC, 0xCC))
# Amount — big, at the end of bar
tb(s, Emu(5200000), E(bar_top, 120000), Emu(2300000), Emu(500000),
   "~₽54М", size=36, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

# Alpha bar (green, SHORT — ~40% of PI)
alpha_top = Emu(2200000)
alpha_width = Emu(3000000)  # visually ~40% of PI
rect(s, Emu(480000), alpha_top, alpha_width, bar_h, fill=GREEN)
tb(s, Emu(600000), E(alpha_top, 100000), Emu(2500000), Emu(250000),
   "Альфа платформа", size=16, color=WHITE, bold=True)
tb(s, Emu(600000), E(alpha_top, 350000), Emu(2500000), Emu(250000),
   "лицензия + миграция + поддержка 5 лет", size=10, color=RGBColor(0xCC, 0xFF, 0xCC))
# Amount
tb(s, Emu(3600000), E(alpha_top, 120000), Emu(2000000), Emu(500000),
   "~₽20М*", size=36, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

# Savings callout — the "wow" moment
savings_y = Emu(3200000)
oval(s, Emu(3200000), savings_y, Emu(2800000), Emu(900000), fill=CYAN)
txBox = s.shapes.add_textbox(Emu(3200000), E(savings_y, 120000), Emu(2800000), Emu(700000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Экономия"
p.font.size = Pt(12)
p.font.color.rgb = WHITE
p.font.name = 'Arial'
p.alignment = PP_ALIGN.CENTER
p2 = add_p(tf, "~₽34М за 5 лет", size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_before=Pt(2))

# Footnote
tb(s, Emu(480000), Emu(4500000), Emu(8000000), Emu(300000),
   "* Оценка для типовой конфигурации. Точный расчёт — после аудита PI-инфраструктуры.",
   size=9, color=MED_GRAY)

# Thin accent line
rect(s, 0, Emu(5093500), SW, Emu(50000), fill=CYAN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: PROOF — Real Projects (Duarte: allies on the journey)
# Real photos from Atomicsoft template + Novatek specifically
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "Доказательство: уже работает на объектах вашего класса",
           "Альфа платформа в промышленной эксплуатации")
footer_accent(s)

# Large Novatek card (highlighted — "you already know us")
novatek_y = Emu(920000)
card_nv = rrect(s, Emu(200000), novatek_y, Emu(4300000), Emu(3700000), fill=WHITE)
card_nv.line.color.rgb = CYAN
card_nv.line.width = Pt(2)

# Accent stripe
rect(s, Emu(200000), novatek_y, Emu(4300000), Emu(55000), fill=CYAN)

# Photo
if project_novatek:
    pic(s, project_novatek, Emu(300000), E(novatek_y, 130000), Emu(4100000), Emu(1800000))

tb(s, Emu(350000), E(novatek_y, 2000000), Emu(4000000), Emu(350000),
   "НОВАТЭК — Арктик СПГ 2", size=18, color=BLUE, bold=True)
tb(s, Emu(350000), E(novatek_y, 2350000), Emu(4000000), Emu(600000),
   "Автоматизация технологических процессов\nСПГ-производства. Альфа платформа\nв составе АСУ ТП крупнейшего объекта.",
   size=11, color=DARK_SLATE)

# Side label
rrect(s, Emu(3300000), E(novatek_y, 3100000), Emu(1100000), Emu(350000), fill=CYAN)
tb(s, Emu(3300000), E(novatek_y, 3140000), Emu(1100000), Emu(280000),
   "Ваш класс\nобъектов", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Right side — other projects (smaller cards)
other_projects = [
    ("Транснефть, ВСТО-II", "Магистральный нефтепровод.\nДиспетчеризация НПС.", BLUE),
    ("СИБУР", "Нефтехимический холдинг.\nАСУ ТП производств.", ACCENT_BLUE),
    ("Интер РАО, Сочинская ТЭС", "Тепловая электростанция.\nМнемосхемы + архив + алармы.", CYAN),
    ("СберСити", "Интеллектуальное здание.\nУправление инженерными\nсистемами.", GREEN),
]

y = Emu(920000)
for name, desc, color in other_projects:
    card = rrect(s, Emu(4700000), y, Emu(4200000), Emu(830000), fill=WHITE)
    card.line.color.rgb = RGBColor(0xE0, 0xE5, 0xEB)
    card.line.width = Pt(1)
    
    # Color bar
    rect(s, Emu(4700000), y, Emu(50000), Emu(830000), fill=color)
    
    tb(s, Emu(4850000), E(y, 80000), Emu(3900000), Emu(250000),
       name, size=13, color=BLUE, bold=True)
    tb(s, Emu(4850000), E(y, 340000), Emu(3900000), Emu(450000),
       desc, size=10, color=DARK_SLATE)
    
    y = E(y, 920000)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: ARCHITECTURE — Technical Compatibility
# Real architecture diagram from Atomicsoft template
# (slide:ology: diagram type = layers/location)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "Архитектура и совместимость",
           "Альфа платформа интегрируется с существующей инфраструктурой")
footer_accent(s)

# Architecture image (real from template)
if architecture_img:
    pic(s, architecture_img, Emu(200000), Emu(900000), Emu(5200000), Emu(3800000))

# Key protocols on the right
proto_groups = [
    ("Промышленные протоколы", [
        "OPC UA/DA (клиент + сервер)",
        "Modbus TCP/RTU",
        "IEC 60870-5-104/101",
        "IEC 61850, S7, MQTT",
    ], BLUE),
    ("Операционные системы", [
        "Windows + Astra Linux",
        "РЕД ОС + ALT Linux",
        "Кроссплатформенный HMI",
    ], CYAN),
    ("IT-интеграция", [
        "SQL через Alpha.RMap",
        "LDAP / Active Directory",
        "Syslog для SIEM",
        "Экспорт CSV/PDF/XLSX",
    ], GREEN),
]

y = Emu(920000)
for title, items, color in proto_groups:
    rect(s, Emu(5600000), y, Emu(3300000), Emu(280000), fill=color)
    tb(s, Emu(5700000), E(y, 40000), Emu(3100000), Emu(220000),
       title, size=11, color=WHITE, bold=True)
    
    iy = E(y, 330000)
    for item in items:
        tb(s, Emu(5700000), iy, Emu(3100000), Emu(200000),
           f"▸  {item}", size=9, color=DARK_SLATE)
        iy = E(iy, 200000)
    
    y = E(iy, 150000)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: WHAT IS — Fears Acknowledged (Duarte: inokulation)
# "We know what you're worried about"
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "Мы понимаем ваши опасения",
           "Каждый из этих рисков имеет конкретное решение →")
footer_accent(s)

fears = [
    ("«А если функций не хватит?»",
     "Аудит PI-инфраструктуры покажет точный список.\n70% функций — прямая замена. 30% — адаптация\nчерез Alpha.Domain, Alpha.Om, Alpha.Reports.",
     "Решение: аудит", GREEN),
    ("«Люди привыкли к PI»",
     "Параллельная работа 2–4 недели. Обучение от вендора.\nУчебные курсы Атомик Софт. Пилот на одном объекте\nкак доказательство.",
     "Решение: пилот", CYAN),
    ("«Потеряем исторические данные»",
     "Конвертация архивов PI → Alpha.Historian 4.0.\nВерификация данных на этапе пилота.\nPI работает параллельно до полной проверки.",
     "Решение: миграция", BLUE),
    ("«Не уложимся в сроки и бюджет»",
     "Фиксированные этапы с KPI. POC за 2–3 месяца.\nБюджет в рублях, без валютных рисков.\nЕжемесячная отчётность.",
     "Решение: контроль", ACCENT_BLUE),
]

y = Emu(920000)
for fear, answer, solution_label, color in fears:
    # Fear block (left)
    fear_bg = rrect(s, Emu(200000), y, Emu(3200000), Emu(730000), fill=RGBColor(0xFE, 0xF3, 0xE8))
    tb(s, Emu(350000), E(y, 80000), Emu(2900000), Emu(600000),
       fear, size=13, color=RED_ACCENT, bold=True)
    
    # Arrow
    tb(s, Emu(3450000), E(y, 200000), Emu(400000), Emu(300000),
       "→", size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    
    # Answer block (right)
    ans_bg = rrect(s, Emu(3900000), y, Emu(4900000), Emu(730000), fill=RGBColor(0xF0, 0xF8, 0xF0))
    # Solution badge
    rrect(s, Emu(7500000), E(y, 20000), Emu(1250000), Emu(240000), fill=color)
    tb(s, Emu(7500000), E(y, 40000), Emu(1250000), Emu(200000),
       solution_label, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    tb(s, Emu(4050000), E(y, 80000), Emu(3400000), Emu(600000),
       answer, size=10, color=DARK_SLATE)
    
    y = E(y, 810000)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: PLAN — Transition Phases (Zelazny: timeline diagram)
# (slide:ology: linear diagram type, clear flow)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "План перехода: 5 фаз за 8–14 месяцев",
           "Ключевой принцип: параллельная работа PI + Альфа. Нулевой простой.")
footer_accent(s)

phases = [
    ("АУДИТ",       "1–2 мес.",  "Инвентаризация:\nтеги, протоколы,\nскрипты, Event Frames", BLUE),
    ("ПРОЕКТ",      "1 мес.",    "Архитектура,\nлицензии, бюджет,\nрабочая группа", ACCENT_BLUE),
    ("ПИЛОТ",       "2–3 мес.",  "Один объект.\nPI + Альфа\nпараллельно", CYAN),
    ("МИГРАЦИЯ",    "3–6 мес.",  "Поэтапный перенос.\nОбучение.\nВалидация.", GREEN),
    ("ЗАВЕРШЕНИЕ",  "1–2 мес.",  "Вывод PI.\nОптимизация.\nПередача в экспл.", RGBColor(0x61, 0x96, 0x2F)),
]

# Timeline bar
timeline_y = Emu(1200000)
rect(s, Emu(200000), E(timeline_y, 220000), Emu(8700000), Emu(60000), fill=RGBColor(0xDD, 0xE5, 0xEE))

phase_w = Emu(1600000)
gap = Emu(130000)
start_x = Emu(300000)

for i, (name, duration, desc, color) in enumerate(phases):
    x = E(start_x, i * (phase_w + gap))
    
    # Circle on timeline
    circle = oval(s, E(x, 600000), timeline_y, Emu(400000), Emu(400000), fill=color)
    tb(s, E(x, 600000), E(timeline_y, 80000), Emu(400000), Emu(280000),
       str(i + 1), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    # Phase name below
    tb(s, x, E(timeline_y, 500000), phase_w, Emu(250000),
       name, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
    
    # Duration
    tb(s, x, E(timeline_y, 730000), phase_w, Emu(200000),
       duration, size=11, color=MED_GRAY, align=PP_ALIGN.CENTER)
    
    # Description
    tb(s, E(x, 30000), E(timeline_y, 1000000), E(phase_w, -60000), Emu(700000),
       desc, size=9, color=DARK_SLATE, align=PP_ALIGN.CENTER)

# Key insight box
rrect(s, Emu(200000), Emu(3900000), Emu(8700000), Emu(700000), fill=RGBColor(0xE3, 0xF2, 0xFD))
txBox = s.shapes.add_textbox(Emu(400000), Emu(3960000), Emu(8300000), Emu(600000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Параллельная работа = нулевой риск для производства"
p.font.size = Pt(14)
p.font.color.rgb = BLUE
p.font.bold = True
p.font.name = 'Arial'
p2 = add_p(tf, "PI System и Альфа работают одновременно на каждом этапе. Переключение — только после подтверждения корректности всех данных и функций. Ни минуты простоя.", size=10, color=DARK_SLATE, space_before=Pt(6))


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: CALL TO ACTION (Duarte: Turning Point 2)
# (Zelazny: action program — who, what, when, how much)
# Three concrete steps, nothing vague
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header_bar(s, "Предлагаем утвердить три шага")
footer_accent(s)

steps = [
    ("1", "Запуск аудита PI-инфраструктуры",
     "Инвентаризация: теги, протоколы, скрипты, отчёты, Event Frames.\nСрок: 1–2 месяца. Результат: отчёт + архитектура + бюджет.",
     "Кто: IT + Атомик Софт", BLUE),
    ("2", "Определение пилотного объекта",
     "Один объект для проверки на практике.\nPI и Альфа работают параллельно 2–4 недели.\nВалидация данных, обучение операторов.",
     "Кто: автоматизация + IT", CYAN),
    ("3", "Формирование рабочей группы",
     "Еженедельные статусы, ежемесячные отчёты руководству.\nФиксированный бюджет и сроки. Контрольные точки.",
     "Кто: VP IT (спонсор)", GREEN),
]

y = Emu(950000)
for num, title, desc, who, color in steps:
    # Number circle
    circle = oval(s, Emu(350000), y, Emu(550000), Emu(550000), fill=color)
    tb(s, Emu(350000), E(y, 110000), Emu(550000), Emu(380000),
       num, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    # Title
    tb(s, Emu(1100000), E(y, 20000), Emu(7000000), Emu(300000),
       title, size=17, color=BLUE, bold=True)
    
    # Description
    tb(s, Emu(1100000), E(y, 280000), Emu(5500000), Emu(500000),
       desc, size=11, color=DARK_SLATE)
    
    # Who badge
    rrect(s, Emu(6800000), E(y, 280000), Emu(2000000), Emu(280000), fill=color)
    tb(s, Emu(6800000), E(y, 300000), Emu(2000000), Emu(240000),
       who, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    y = E(y, 1100000)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: NEW BLISS (Duarte: end HIGHER than you began)
# Vision of the future — emotional, minimal text
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)

# Industrial background
if industrial_bg:
    pic(s, industrial_bg, 0, 0, SW, SH)
rect(s, 0, 0, SW, SH, fill=BLACK, alpha_pct=70)

# Vision text — large, emotional (Duarte: New Bliss imagery)
txBox = s.shapes.add_textbox(Emu(480000), Emu(800000), Emu(8000000), Emu(2000000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Через 12 месяцев:"
p.font.size = Pt(14)
p.font.color.rgb = CYAN
p.font.name = 'Arial'

visions = [
    "Ваша SCADA — полностью ваша.",
    "Никаких подписок. Никаких санкций.",
    "Техподдержка на расстоянии звонка.",
    "Полное соответствие ФЗ-187 о КИИ.",
]
for v in visions:
    p2 = add_p(tf, v, size=24, color=WHITE, bold=True, space_before=Pt(14))

# Savings reminder
tb(s, Emu(480000), Emu(3600000), Emu(8000000), Emu(400000),
   "₽34М экономии за 5 лет. И это — консервативная оценка.",
   size=16, color=GREEN, bold=True)

# Bottom bar
rect(s, 0, Emu(5093500), SW, Emu(50000), fill=CYAN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: CONTACTS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SW, SH, fill=DARK_BG)

# Decorative
pic(s, deco_circle, Emu(-800000), Emu(100000), Emu(4900000), Emu(4900000))

# Closing phrase
txBox = s.shapes.add_textbox(Emu(480000), Emu(700000), Emu(8000000), Emu(1200000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Будущее автоматизации"
p.font.size = Pt(32)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = 'Arial'
p2 = add_p(tf, "в ваших руках.", size=32, color=CYAN, bold=True, space_before=Pt(4))

# Company info
txBox2 = s.shapes.add_textbox(Emu(2000000), Emu(2600000), Emu(3500000), Emu(1800000))
tf2 = txBox2.text_frame
tf2.word_wrap = True
contacts_data = [
    ("АО «Атомик Софт»", 14, WHITE, True),
    ("634050 Томск, пр. Ленина 60/1", 12, RGBColor(0xBB, 0xCC, 0xDD), False),
    ("", 8, WHITE, False),
    ("+7 (3822) 281 914", 12, RGBColor(0xBB, 0xCC, 0xDD), False),
    ("info@automiq.ru  •  automiq.ru", 12, CYAN, False),
]
for i, (text, size, color, bold) in enumerate(contacts_data):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = add_p(tf2, text, size=size, color=color, bold=bold, space_before=Pt(2))
        continue
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'

# Personal contact
pc = s.shapes.add_textbox(Emu(6000000), Emu(3000000), Emu(3000000), Emu(1000000))
tf3 = pc.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "Ваш контакт:"
p.font.size = Pt(10)
p.font.color.rgb = MED_GRAY
p.font.name = 'Arial'
p2 = add_p(tf3, "Станислав Павловский", size=14, color=WHITE, bold=True, space_before=Pt(4))
p3 = add_p(tf3, "Telegram: @Integraleus", size=12, color=CYAN, space_before=Pt(2))

# Logo
pic(s, logo, Emu(480000), Emu(4300000), Emu(2200000), Emu(340000))

# QR logo
pic(s, qr_logo, Emu(6800000), Emu(2000000), Emu(1200000), Emu(1200000))

# Bottom accent
rect(s, 0, Emu(5093500), SW, Emu(50000), fill=CYAN)


# ── Save ──
output = "NOVATEK_VP_IT_resonate_v2.pptx"
prs.save(output)
print(f"✅ Saved: {output}")
print(f"   Slides: {len(prs.slides)}")
