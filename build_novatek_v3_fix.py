#!/usr/bin/env python3
"""
NOVATEK VP IT — v3 FIXED
Fixes: emoji rendering, text overflow, photo overlap, truncated text,
unverified percentages, QR/contact overlap
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import io, os

# ── Colors ──
BLUE = RGBColor(0x00, 0x54, 0x97)
CYAN = RGBColor(0x09, 0x97, 0xC8)
GREEN = RGBColor(0x82, 0xC4, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
SLATE = RGBColor(0x44, 0x54, 0x69)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)
ABLU = RGBColor(0x33, 0x75, 0xAD)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xE8, 0x77, 0x22)

SW = 9144000; SH = 5143500

# Assets
A = '/tmp/pptx_assets'
def ld(n):
    p = os.path.join(A, n)
    return open(p,'rb').read() if os.path.exists(p) else None

title_bg = ld('title_bg.jpg')
ind_bg = ld('industrial_bg.jpg')
plat_ss = ld('platform_screenshot.png')
arch_img = ld('architecture.png')
logo = ld('logo.png')
deco = ld('deco_circle.png')
qr = ld('qr_logo.png')

prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH
BL = prs.slide_layouts[6]

def R(s,l,t,w,h,fill=None,alpha=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb=fill
        if alpha:
            ns={'a':'http://schemas.openxmlformats.org/drawingml/2006/main'}
            sf=sh._element.find('.//a:solidFill',ns)
            if sf is not None:
                sr=sf.find('a:srgbClr',ns)
                if sr is not None:
                    ae=etree.SubElement(sr,'{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    ae.set('val',str(int(alpha*1000)))
    else: sh.fill.background()
    return sh

def RR(s,l,t,w,h,fill=None,brd=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if brd: sh.line.color.rgb=brd; sh.line.fill.solid(); sh.line.width=Pt(1)
    return sh

def O(s,l,t,w,h,fill=None):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,w,h)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    return sh

def T(s,l,t,w,h,text,sz=14,c=BLACK,b=False,a=PP_ALIGN.LEFT):
    tx=s.shapes.add_textbox(l,t,w,h)
    tf=tx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text
    p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name='Arial'; p.alignment=a
    return tx

def P(tf,text,sz=14,c=BLACK,b=False,a=PP_ALIGN.LEFT,sp=Pt(4)):
    p=tf.add_paragraph(); p.text=text
    p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name='Arial'; p.alignment=a; p.space_before=sp
    return p

def IMG(s,blob,l,t,w,h):
    if blob: return s.shapes.add_picture(io.BytesIO(blob),l,t,w,h)

def E(*a): return Emu(sum(a))

def hdr(s,title,sub=None):
    R(s,0,E(180000),SW,E(620000),fill=BLUE)
    R(s,0,E(800000),SW,E(28000),fill=CYAN)
    T(s,E(500000),E(220000),E(7400000),E(550000),title,sz=20,c=WHITE,b=True)
    if sub: T(s,E(500000),E(530000),E(7400000),E(260000),sub,sz=11,c=RGBColor(0xBB,0xCC,0xDD))
    IMG(s,logo,E(7200000),E(230000),E(1600000),E(245000))

def ftr(s):
    R(s,0,E(4960000),SW,E(28000),fill=CYAN)
    R(s,0,E(4988000),SW,E(155500),fill=BLUE)


# ═══════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
if title_bg: IMG(s,title_bg,0,0,SW,SH)
R(s,0,0,SW,SH,fill=BLACK,alpha=65)
R(s,0,0,SW,E(45000),fill=CYAN)
IMG(s,logo,E(480000),E(280000),E(2400000),E(370000))

tx=s.shapes.add_textbox(E(480000),E(1100000),E(8000000),E(1500000))
tf=tx.text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; p.text="Замена PI System возможна."
p.font.size=Pt(36); p.font.color.rgb=WHITE; p.font.bold=True; p.font.name='Arial'
P(tf,"Проверено на объектах вашего класса.",sz=36,c=CYAN,b=True,sp=Pt(4))

T(s,E(480000),E(2800000),E(7500000),E(400000),
  "Стратегия перехода на Альфа платформу",sz=16,c=RGBColor(0xBB,0xCC,0xDD))

R(s,0,E(4300000),SW,E(1),fill=RGBColor(0x44,0x55,0x77))
T(s,E(480000),E(4400000),E(4500000),E(350000),
  "АО «Атомик Софт»   automiq.ru   Конфиденциально",sz=11,c=GRAY)
RR(s,E(6500000),E(4380000),E(2300000),E(400000),fill=RGBColor(0x22,0x33,0x55))
T(s,E(6500000),E(4420000),E(2300000),E(340000),
  "Внутренняя встреча   2026",sz=11,c=RGBColor(0x88,0x99,0xBB),a=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 2: СИТУАЦИЯ (no invented $$$)
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
hdr(s,"Ситуация: зависимость от ушедшего вендора")
ftr(s)

cards=[
    ("Санкции","AVEVA прекратила\nподдержку и продажу\nв России.\n\nНет патчей, нет\nобновлений с 2022 г.",RED),
    ("Подписка","Модель AVEVA Flex:\nпри неоплате ПО\nперестаёт работать.\n\nОблачные сервисы\nуже отключены.",ORANGE),
    ("Регуляторика","ФЗ-187 о безопасности\nКИИ. Указ №250.\n\nТребование перехода\nна отечественное ПО.\nСроки ужесточаются.",BLUE),
    ("Кадры","Новые специалисты\nобучаются на российском\nПО. Экспертиза по\nPI System сокращается.",ABLU),
]

for i,(label,detail,color) in enumerate(cards):
    x=E(150000+i*2250000); y=E(920000)
    RR(s,x,y,E(2100000),E(3700000),fill=WHITE,brd=RGBColor(0xE0,0xE5,0xEB))
    R(s,x,y,E(2100000),E(55000),fill=color)
    T(s,E(x,80000),E(y,130000),E(1940000),E(300000),label,sz=14,c=color,b=True,a=PP_ALIGN.CENTER)
    T(s,E(x,120000),E(y,500000),E(1860000),E(3000000),detail,sz=11,c=SLATE,a=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 3: АЛЬФА ПЛАТФОРМА (screenshot smaller, facts bigger)
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
R(s,0,0,SW,SH,fill=LIGHT_BG)
hdr(s,"Альтернатива: Альфа платформа",
    "Российская SCADA №1 в реестре Минцифры — Атомик Софт, Томск")
ftr(s)

if plat_ss: IMG(s,plat_ss,E(5000000),E(1300000),E(3800000),E(3200000))

facts=[
    ("Бессрочная лицензия в рублях","Одна покупка. Нет подписки.\nНе отключат при неоплате.",GREEN),
    ("6 500+ инсталляций","Нефтегаз, энергетика,\nкритическая инфраструктура.",CYAN),
    ("Windows + Linux","Astra Linux, РЕД ОС, Альт.\nСертифицированные ОС.",BLUE),
    ("Техподдержка 24/7","На русском. Вендор в Томске.\nЛицензия ФСТЭК.",ABLU),
]

y=E(1000000)
for title,desc,color in facts:
    R(s,E(400000),y,E(50000),E(550000),fill=color)
    T(s,E(550000),E(y,40000),E(4200000),E(300000),title,sz=16,c=BLUE,b=True)
    T(s,E(550000),E(y,320000),E(4200000),E(300000),desc,sz=11,c=SLATE)
    y=E(y,650000)


# ═══════════════════════════════════════════════
# SLIDE 4: СКЛАД VS ДИСПЕТЧЕРСКАЯ (no emoji, shorter title)
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
hdr(s,"Ключевое различие: Склад vs Диспетчерская")
ftr(s)

# PI
R(s,E(150000),E(920000),E(4300000),E(380000),fill=RED)
T(s,E(150000),E(945000),E(4300000),E(340000),
  "AVEVA PI System — «Склад данных»",sz=14,c=WHITE,b=True,a=PP_ALIGN.CENTER)

T(s,E(250000),E(1380000),E(4100000),E(800000),
  "Огромный склад с идеальной системой учёта.\nЛюбой товар можно найти, любому курьеру\nвыдать, в любую аналитику подключить.\nНо сам склад ничем не управляет.",sz=11,c=SLATE)

pi_pros=["+ 450+ коннекторов","+ REST API — данные по HTTP",
         "+ Облако — данные отовсюду","+ Event Frames — запись событий",
         "+ Оператор сам рисует дашборд"]
y=E(2300000)
for item in pi_pros:
    T(s,E(300000),y,E(4000000),E(230000),item,sz=10,c=SLATE); y=E(y,230000)

# Alpha
R(s,E(4700000),E(920000),E(4300000),E(380000),fill=GREEN)
T(s,E(4700000),E(945000),E(4300000),E(340000),
  "Альфа платформа — «Диспетчерская»",sz=14,c=WHITE,b=True,a=PP_ALIGN.CENTER)

T(s,E(4800000),E(1380000),E(4100000),E(800000),
  "Диспетчерская с пультом. Оператор видит всё,\nуправляет всем, все тревоги приходят сюда.\nСклад тоже есть (Historian),\nно главное — пульт управления.",sz=11,c=SLATE)

al_pros=["+ Мощный дизайнер мнемосхем (HMI)","+ ~30 протоколов — все основные",
         "+ Напрямую к оборудованию","+ Linux + Windows, on-premises",
         "+ Бессрочная лицензия в рублях"]
y=E(2300000)
for item in al_pros:
    T(s,E(4800000),y,E(4100000),E(230000),item,sz=10,c=SLATE); y=E(y,230000)

RR(s,E(150000),E(3700000),E(8850000),E(500000),fill=RGBColor(0xE3,0xF2,0xFD))
T(s,E(350000),E(3740000),E(8450000),E(420000),
  "Миграция = переезд со «склада данных» на «диспетчерскую с собственным складом».\nPI хорош в сборе/раздаче данных. Альфа — в управлении и визуализации. Это разные инструменты.",
  sz=11,c=BLUE,b=True)


# ═══════════════════════════════════════════════
# SLIDE 5: ТАБЛИЦА СРАВНЕНИЯ (honest, no invented stats)
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
hdr(s,"Сравнение: PI System vs Альфа платформа",
    "Честная оценка обеих платформ")
ftr(s)

rows=[
    ("Что это","Инфраструктура данных","SCADA-платформа"),
    ("Главный фокус","Собрать и раздать данные","Показать и управлять"),
    ("Визуализация","Вторична (PI Vision)","Ядро продукта (Alpha.HMI)"),
    ("Коннекторы","450+ (по данным AVEVA)","~30 встроенных модулей"),
    ("Модель объектов","Runtime (на ходу)","Design-time (заранее)"),
    ("API","REST (Web API)","OPC UA + TCP + SQL"),
    ("Облако","Да (AVEVA Data Hub)","Нет (только on-prem)"),
    ("Event Frames","Да (уникальная фича)","Нет (нужна доработка)"),
    ("ОС","Только Windows","Windows + Linux"),
    ("Лицензия","Подписка, доллары","Бессрочная, рубли"),
    ("Поддержка в РФ","Нет (санкции)","24/7 на русском"),
]

R(s,E(150000),E(900000),E(2700000),E(330000),fill=SLATE)
T(s,E(200000),E(930000),E(2600000),E(280000),"Параметр",sz=11,c=WHITE,b=True,a=PP_ALIGN.CENTER)
R(s,E(2850000),E(900000),E(3000000),E(330000),fill=RED)
T(s,E(2900000),E(930000),E(2900000),E(280000),"AVEVA PI System",sz=11,c=WHITE,b=True,a=PP_ALIGN.CENTER)
R(s,E(5850000),E(900000),E(3150000),E(330000),fill=GREEN)
T(s,E(5900000),E(930000),E(3050000),E(280000),"Альфа платформа",sz=11,c=WHITE,b=True,a=PP_ALIGN.CENTER)

y=E(1230000); rh=E(305000)
for i,(par,pi,al) in enumerate(rows):
    bg=RGBColor(0xF8,0xF9,0xFB) if i%2==0 else WHITE
    R(s,E(150000),y,E(8850000),rh,fill=bg)
    T(s,E(200000),E(y,45000),E(2600000),E(260000),par,sz=10,c=SLATE,b=True)
    T(s,E(2900000),E(y,45000),E(2900000),E(260000),pi,sz=10,c=SLATE)
    T(s,E(5900000),E(y,45000),E(3050000),E(260000),al,sz=10,c=SLATE)
    y=E(y,rh)


# ═══════════════════════════════════════════════
# SLIDE 6: ПОЛНАЯ ЗАМЕНА (6 компонентов) — shorter title
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
hdr(s,"Полная замена: 6 компонентов PI System")
ftr(s)

fulls=[
    ("Архив данных  ->  Alpha.Historian 4.0",
     "Все данные с датчиков хранятся надёжно.\nСжатие LZMA — архив занимает меньше места."),
    ("Экраны оператора  ->  Alpha.HMI 2.0",
     "Мнемосхемы мощнее, чем в PI ProcessBook.\nДизайнер + визуализатор + веб-доступ."),
    ("Тревоги  ->  Alpha.HMI.Alarms 3.3",
     "Кроссплатформенный: Windows, Linux, браузер.\nКвитирование, подавление, экспорт."),
    ("Буферизация  ->  Alpha.Server",
     "Если связь пропала — данные не теряются.\nРаботает автоматически."),
    ("Администрирование  ->  Alpha.DevStudio 4.1",
     "Единая среда: проектирование, диагностика,\nмониторинг. CLI для автоматизации."),
    ("Программный доступ  ->  OPC UA + RMap",
     "Внешние программы через OPC UA,\nTCP или SQL (Alpha.RMap)."),
]

y=E(920000)
for i,(title,desc) in enumerate(fulls):
    col=i%2; row=i//2
    x=E(150000+col*4500000); cy=E(y+row*1250000)
    RR(s,x,cy,E(4300000),E(1100000),fill=RGBColor(0xF0,0xF8,0xF0),brd=RGBColor(0xCC,0xDD,0xCC))
    R(s,x,cy,E(50000),E(1100000),fill=GREEN)
    T(s,E(x,130000),E(cy,80000),E(4050000),E(300000),"[OK]  "+title,sz=11,c=BLUE,b=True)
    T(s,E(x,130000),E(cy,380000),E(4050000),E(650000),desc,sz=10,c=SLATE)


# ═══════════════════════════════════════════════
# SLIDE 7: ЧАСТИЧНАЯ + ПРОБЕЛЫ (honest)
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
hdr(s,"Честно: что работает иначе и чего нет",
    "Частичная замена (4) и функциональные пробелы (6)")
ftr(s)

R(s,E(150000),E(890000),E(4300000),E(340000),fill=ORANGE)
T(s,E(150000),E(910000),E(4300000),E(300000),
  "[!]  Частичная замена (4)",sz=13,c=WHITE,b=True,a=PP_ALIGN.CENTER)

partials=[
    ("450+ коннекторов  ->  ~30","Основные есть. Экзотика —\nчерез OPC UA шлюзы."),
    ("Оператор рисует дашборд","В Альфе дашборды рисует инженер.\nОператор смотрит готовое."),
    ("Модель активов «на лету»","В Альфе структура проектируется\nзаранее в DevStudio."),
    ("PI DataLink (Excel)","Alpha.Reports есть, но это не\nинтерактивный Excel add-in."),
]

y=E(1300000)
for title,desc in partials:
    T(s,E(250000),y,E(4000000),E(250000),"(!)  "+title,sz=10,c=ORANGE,b=True)
    T(s,E(250000),E(y,240000),E(4000000),E(350000),desc,sz=9,c=SLATE)
    y=E(y,610000)

R(s,E(4700000),E(890000),E(4300000),E(340000),fill=RED)
T(s,E(4700000),E(910000),E(4300000),E(300000),
  "[X]  Функциональные пробелы (6)",sz=13,c=WHITE,b=True,a=PP_ALIGN.CENTER)

gaps=[
    ("Event Frames","Нельзя «запиши что было\nс 10:15 до 11:30 на насосе»"),
    ("REST API","Нет HTTP-доступа к данным"),
    ("Репликация архивов","Нет синхронизации\nмежду серверами"),
    ("Edge-версия","Нет мини-сервера для\nудалённых площадок"),
    ("Ручной ввод данных","Нет формы для лабораторных\nанализов и обходов"),
    ("Коннектор Power BI","Нет прямого. Обход:\nSQL -> PostgreSQL -> BI"),
]

y=E(1300000)
for title,desc in gaps:
    T(s,E(4800000),y,E(4000000),E(250000),"(X)  "+title,sz=10,c=RED,b=True)
    T(s,E(4800000),E(y,230000),E(4000000),E(300000),desc,sz=9,c=SLATE)
    y=E(y,540000)


# ═══════════════════════════════════════════════
# SLIDE 8: СТРАТЕГИЯ ЗАКРЫТИЯ (no unverified %)
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
hdr(s,"Стратегия закрытия пробелов",
    "Каждый пробел имеет конкретное решение")
ftr(s)

R(s,E(150000),E(890000),E(3500000),E(330000),fill=RED)
T(s,E(200000),E(910000),E(3400000),E(290000),"Чего не хватает",sz=12,c=WHITE,b=True,a=PP_ALIGN.CENTER)
R(s,E(3650000),E(890000),E(3700000),E(330000),fill=GREEN)
T(s,E(3700000),E(910000),E(3600000),E(290000),"Как решаем",sz=12,c=WHITE,b=True,a=PP_ALIGN.CENTER)
R(s,E(7350000),E(890000),E(1650000),E(330000),fill=ABLU)
T(s,E(7400000),E(910000),E(1550000),E(290000),"Подход",sz=12,c=WHITE,b=True,a=PP_ALIGN.CENTER)

strat=[
    ("Event Frames",       "Alpha.Om + PostgreSQL",                  "Доработка"),
    ("REST API",           "REST-обёртка над SQL-данными",            "Обёртка"),
    ("Репликация",         "Экспорт Historian -> PostgreSQL",         "Периодический"),
    ("Edge-сервер",        "Alpha.Server в мини-конфигурации",        "Лёгкая версия"),
    ("Ручной ввод",        "Веб-форма через HMI.WebViewer",          "Веб-форма"),
    ("BI-аналитика",       "SQL -> Grafana / Power BI",              "Через СУБД"),
    ("Экзотич. протоколы", "Основные встроены, остальное OPC UA",     "Шлюзы"),
]

y=E(1280000); rh=E(420000)
for i,(gap,sol,app) in enumerate(strat):
    bg=RGBColor(0xF8,0xF9,0xFB) if i%2==0 else WHITE
    R(s,E(150000),y,E(8850000),rh,fill=bg)
    T(s,E(250000),E(y,65000),E(3300000),E(350000),gap,sz=11,c=SLATE,b=True)
    T(s,E(3750000),E(y,65000),E(3500000),E(350000),sol,sz=11,c=SLATE)
    T(s,E(7450000),E(y,65000),E(1450000),E(350000),app,sz=10,c=ABLU,b=True,a=PP_ALIGN.CENTER)
    y=E(y,rh)

RR(s,E(150000),E(4300000),E(8850000),E(430000),fill=RGBColor(0xE8,0xF5,0xE9))
T(s,E(350000),E(4340000),E(8450000),E(360000),
  "Все пробелы закрываемы. Аудит покажет, какие из них реально используются у вас.",
  sz=12,c=RGBColor(0x2E,0x7D,0x32),b=True)


# ═══════════════════════════════════════════════
# SLIDE 9: ПРОЕКТЫ (no overlapping photo)
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
hdr(s,"Опыт: объекты вашего класса",
    "Альфа платформа в промышленной эксплуатации")
ftr(s)

projects=[
    ("НОВАТЭК — Арктик СПГ 2","Автоматизация технологических\nпроцессов СПГ-производства.\nАльфа в составе АСУ ТП.",CYAN),
    ("Транснефть — ВСТО-II","Магистральный нефтепровод.\nДиспетчеризация и управление\nнасосных станций.",BLUE),
    ("СИБУР","Нефтехимический холдинг.\nАльфа в составе систем\nуправления производствами.",ABLU),
    ("Интер РАО — Сочинская ТЭС","Тепловая электростанция.\nМнемосхемы + архив +\nтревоги.",GREEN),
]

for i,(name,desc,color) in enumerate(projects):
    col=i%2; row=i//2
    x=E(150000+col*4500000); y=E(920000+row*1900000)
    RR(s,x,y,E(4300000),E(1650000),fill=WHITE,brd=RGBColor(0xE0,0xE5,0xEB))
    R(s,x,y,E(50000),E(1650000),fill=color)
    T(s,E(x,150000),E(y,100000),E(4000000),E(350000),name,sz=16,c=BLUE,b=True)
    T(s,E(x,150000),E(y,500000),E(4000000),E(1000000),desc,sz=12,c=SLATE)


# ═══════════════════════════════════════════════
# SLIDE 10: ПЛАН ПЕРЕХОДА
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
hdr(s,"План перехода: 5 фаз, 8-14 месяцев",
    "PI и Альфа работают параллельно. Нулевой простой.")
ftr(s)

phases=[
    ("АУДИТ","1-2 мес.","Инвентаризация:\nтеги, протоколы,\nскрипты, EF",BLUE),
    ("ПРОЕКТ","1 мес.","Архитектура,\nлицензии, бюджет,\nрабочая группа",ABLU),
    ("ПИЛОТ","2-3 мес.","Один объект.\nPI + Альфа\nпараллельно",CYAN),
    ("МИГРАЦИЯ","3-6 мес.","Поэтапный\nперенос. Обучение.\nВалидация.",GREEN),
    ("ЗАКРЫТИЕ","1-2 мес.","Вывод PI.\nОптимизация.\nПередача.",RGBColor(0x61,0x96,0x2F)),
]

tl_y=E(1300000)
R(s,E(200000),E(tl_y,230000),E(8700000),E(60000),fill=RGBColor(0xDD,0xE5,0xEE))

pw=E(1600000); gp=E(130000); sx=E(300000)
for i,(nm,dur,desc,col) in enumerate(phases):
    x=E(sx,i*(pw+gp))
    O(s,E(x,600000),tl_y,E(420000),E(420000),fill=col)
    T(s,E(x,600000),E(tl_y,90000),E(420000),E(300000),str(i+1),sz=18,c=WHITE,b=True,a=PP_ALIGN.CENTER)
    T(s,x,E(tl_y,530000),pw,E(280000),nm,sz=13,c=col,b=True,a=PP_ALIGN.CENTER)
    T(s,x,E(tl_y,790000),pw,E(220000),dur,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    T(s,E(x,30000),E(tl_y,1050000),E(pw,-60000),E(700000),desc,sz=10,c=SLATE,a=PP_ALIGN.CENTER)

RR(s,E(200000),E(4000000),E(8700000),E(650000),fill=RGBColor(0xE3,0xF2,0xFD))
tx=s.shapes.add_textbox(E(400000),E(4050000),E(8300000),E(560000))
tf=tx.text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; p.text="Параллельная работа = нулевой риск для производства"
p.font.size=Pt(13); p.font.color.rgb=BLUE; p.font.bold=True; p.font.name='Arial'
P(tf,"PI и Альфа работают одновременно. Переключение — только после подтверждения корректности.",sz=10,c=SLATE,sp=Pt(6))


# ═══════════════════════════════════════════════
# SLIDE 11: СЛЕДУЮЩИЕ ШАГИ
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
hdr(s,"Предлагаемые следующие шаги")
ftr(s)

steps=[
    ("1","Согласование объёма аудита PI-инфраструктуры",BLUE),
    ("2","Определение пилотного объекта для миграции",CYAN),
    ("3","Перечень критичных функций PI\n(Event Frames? DataLink? API?)",ABLU),
    ("4","Расчёт лицензий Альфа под целевую архитектуру",GREEN),
    ("5","Рабочая группа: Новатэк + Атомик Софт + интегратор",RGBColor(0x61,0x96,0x2F)),
    ("6","Детальный план-график и бюджет",ORANGE),
]

y=E(920000)
for num,text,col in steps:
    O(s,E(350000),y,E(450000),E(450000),fill=col)
    T(s,E(350000),E(y,90000),E(450000),E(320000),num,sz=22,c=WHITE,b=True,a=PP_ALIGN.CENTER)
    T(s,E(1000000),E(y,80000),E(7500000),E(400000),text,sz=14,c=SLATE)
    y=E(y,570000)


# ═══════════════════════════════════════════════
# SLIDE 12: ВИДЕНИЕ (bigger text)
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
if ind_bg: IMG(s,ind_bg,0,0,SW,SH)
R(s,0,0,SW,SH,fill=BLACK,alpha=75)

T(s,E(480000),E(600000),E(8000000),E(400000),
  "Через 12 месяцев:",sz=18,c=CYAN,b=True)

visions=["Ваша SCADA — полностью ваша.",
         "Никаких подписок. Никаких санкций.",
         "Техподдержка на расстоянии звонка.",
         "Полное соответствие ФЗ-187 о КИИ."]
y=E(1100000)
for v in visions:
    T(s,E(480000),y,E(8000000),E(450000),v,sz=24,c=WHITE,b=True); y=E(y,500000)

R(s,0,E(5093500),SW,E(50000),fill=CYAN)


# ═══════════════════════════════════════════════
# SLIDE 13: КОНТАКТЫ (QR moved, no overlap)
# ═══════════════════════════════════════════════
s=prs.slides.add_slide(BL)
R(s,0,0,SW,SH,fill=DARK_BG)
IMG(s,deco,E(-800000),E(100000),E(4900000),E(4900000))

tx=s.shapes.add_textbox(E(480000),E(700000),E(8000000),E(1200000))
tf=tx.text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; p.text="Будущее автоматизации"
p.font.size=Pt(32); p.font.color.rgb=WHITE; p.font.bold=True; p.font.name='Arial'
P(tf,"в ваших руках.",sz=32,c=CYAN,b=True,sp=Pt(4))

tx2=s.shapes.add_textbox(E(2000000),E(2600000),E(3500000),E(1800000))
tf2=tx2.text_frame; tf2.word_wrap=True
p=tf2.paragraphs[0]; p.text="АО «Атомик Софт»"
p.font.size=Pt(14); p.font.color.rgb=WHITE; p.font.bold=True; p.font.name='Arial'
P(tf2,"634050 Томск, пр. Ленина 60/1",sz=12,c=RGBColor(0xBB,0xCC,0xDD),sp=Pt(4))
P(tf2,"+7 (3822) 281 914",sz=12,c=RGBColor(0xBB,0xCC,0xDD),sp=Pt(8))
P(tf2,"info@automiq.ru   automiq.ru",sz=12,c=CYAN,sp=Pt(2))

# Contact info — moved DOWN, QR moved to NOT overlap
pc=s.shapes.add_textbox(E(6200000),E(3500000),E(2800000),E(900000))
tf3=pc.text_frame; tf3.word_wrap=True
p=tf3.paragraphs[0]; p.text="Ваш контакт:"
p.font.size=Pt(10); p.font.color.rgb=GRAY; p.font.name='Arial'
P(tf3,"Станислав Павловский",sz=14,c=WHITE,b=True,sp=Pt(4))
P(tf3,"Telegram: @Integraleus",sz=12,c=CYAN,sp=Pt(2))

IMG(s,logo,E(480000),E(4300000),E(2200000),E(340000))
# QR moved higher, not overlapping contact text
IMG(s,qr,E(7000000),E(2000000),E(1100000),E(1100000))
R(s,0,E(5093500),SW,E(50000),fill=CYAN)


# ── Save ──
out="NOVATEK_VP_IT_v3_final.pptx"
prs.save(out)
print(f"OK: {out}, {len(prs.slides)} slides")
