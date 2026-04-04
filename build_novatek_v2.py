#!/usr/bin/env python3
"""
NOVATEK VP IT — Презентация замены PI System на Альфа платформу
Бренд: Атомик Софт / Automiq
Аудитория: Вице-президент по IT и автоматизации
Цель: Заинтересовать, показать что замена возможна
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
import io

# ── Brand colors ──
BLUE = RGBColor(0x00, 0x54, 0x97)       # #005497 — primary
CYAN = RGBColor(0x09, 0x97, 0xC8)       # #0997C8 — accent
GREEN = RGBColor(0x82, 0xC4, 0x44)      # #82C444 — secondary accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_SLATE = RGBColor(0x44, 0x54, 0x69) # body text
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5) # light bg
MED_GRAY = RGBColor(0x6E, 0x6E, 0x6E)   # secondary text
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)    # dark navy for contrast slides

# ── Dimensions (same as atomicsoft template: 10×5.625 in = 16:9) ──
SLIDE_W = 9144000
SLIDE_H = 5143500

# ── Extract resources from template ──
template_prs = Presentation('skills/presentation-designer/templates/atomicsoft-about.pptx')
template_slides = list(template_prs.slides)

# Extract logo from title slide (Рисунок 56 = Atomik Soft logo)
logo_blob = None
logo_ct = None
for shape in template_slides[0].shapes:
    if shape.name == "Рисунок 56":
        logo_blob = shape.image.blob
        logo_ct = shape.image.content_type
        break

# Extract QR/logo from contact slide
contact_logo_blob = None
for shape in template_slides[-1].shapes:
    if shape.name == "Рисунок 12":
        contact_logo_blob = shape.image.blob
        break

# Decorative element (large translucent logo from contacts)
deco_blob = None
for shape in template_slides[-1].shapes:
    if shape.name == "Рисунок 5":
        deco_blob = shape.image.blob
        break

# ── Create new presentation ──
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank


def add_shape_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, font_color=BLACK, bold=False, space_before=Pt(4), space_after=Pt(2), alignment=PP_ALIGN.LEFT, font_name='Arial'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p


def add_header_bar(slide, title_text, subtitle_text=None):
    """Add blue header bar with title (like atomicsoft template)"""
    # Header background
    add_shape_rect(slide, 0, Emu(198120), SLIDE_W, Emu(586740), fill_color=BLUE)
    # Title
    add_text_box(slide, Emu(623128), Emu(230000), Emu(8000000), Emu(520000),
                 title_text, font_size=24, font_color=WHITE, bold=True)
    if subtitle_text:
        add_text_box(slide, Emu(623128), Emu(530000), Emu(8000000), Emu(250000),
                     subtitle_text, font_size=13, font_color=RGBColor(0xCC, 0xDD, 0xEE), bold=False)
    # Thin accent line at bottom of header
    add_shape_rect(slide, 0, Emu(784860), SLIDE_W, Emu(25000), fill_color=CYAN)
    # Logo (small) in top-right if available
    if logo_blob:
        img_stream = io.BytesIO(logo_blob)
        slide.shapes.add_picture(img_stream, Emu(7400000), Emu(240000), Emu(1500000), Emu(230000))


def add_bottom_bar(slide):
    """Thin bottom accent bar"""
    add_shape_rect(slide, 0, Emu(4950000), SLIDE_W, Emu(40000), fill_color=CYAN)
    # Page indicator line
    add_shape_rect(slide, 0, Emu(4990000), SLIDE_W, Emu(153500), fill_color=BLUE)


def add_metric_block(slide, left, top, number, label, color=CYAN):
    """Add a metric tile (number + label)"""
    # Number
    add_text_box(slide, left, top, Emu(1800000), Emu(450000),
                 number, font_size=30, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left, Emu(top + 380000), Emu(1800000), Emu(300000),
                 label, font_size=10, font_color=DARK_SLATE, bold=False, alignment=PP_ALIGN.CENTER)


def add_bullet_list(slide, left, top, width, items, font_size=12, font_color=DARK_SLATE, bullet_char="▪"):
    """Add bullet list"""
    txBox = slide.shapes.add_textbox(left, top, width, Emu(len(items) * 280000))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = 'Arial'
        p.space_before = Pt(4)
        p.space_after = Pt(2)
    return txBox


def add_two_column_block(slide, left_title, left_items, right_title, right_items, 
                          left_color=BLUE, right_color=CYAN, y_start=Emu(950000)):
    """Two-column comparison"""
    col_w = Emu(3900000)
    gap = Emu(200000)
    left_x = Emu(400000)
    right_x = Emu(4700000)
    
    # Left column header
    add_shape_rect(slide, left_x, y_start, col_w, Emu(350000), fill_color=left_color)
    add_text_box(slide, Emu(left_x + 100000), Emu(y_start + 60000), Emu(col_w - 200000), Emu(250000),
                 left_title, font_size=15, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Right column header
    add_shape_rect(slide, right_x, y_start, col_w, Emu(350000), fill_color=right_color)
    add_text_box(slide, Emu(right_x + 100000), Emu(y_start + 60000), Emu(col_w - 200000), Emu(250000),
                 right_title, font_size=15, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Left items
    add_bullet_list(slide, Emu(left_x + 80000), Emu(y_start + 420000), Emu(col_w - 160000), 
                    left_items, font_size=11)
    # Right items
    add_bullet_list(slide, Emu(right_x + 80000), Emu(y_start + 420000), Emu(col_w - 160000),
                    right_items, font_size=11)


# ═══════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)

# Dark navy background
add_shape_rect(slide1, 0, 0, SLIDE_W, SLIDE_H, fill_color=DARK_BG)

# Cyan accent line top
add_shape_rect(slide1, 0, 0, SLIDE_W, Emu(50000), fill_color=CYAN)

# Logo
if logo_blob:
    img_stream = io.BytesIO(logo_blob)
    slide1.shapes.add_picture(img_stream, Emu(500000), Emu(300000), Emu(2600000), Emu(400000))

# Main title
txBox = slide1.shapes.add_textbox(Emu(500000), Emu(1200000), Emu(8000000), Emu(1200000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Альфа платформа"
p.font.size = Pt(38)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = 'Arial'
p2 = tf.add_paragraph()
p2.text = "Стратегия замены PI System"
p2.font.size = Pt(28)
p2.font.color.rgb = CYAN
p2.font.bold = True
p2.font.name = 'Arial'
p2.space_before = Pt(6)

# Subtitle
add_text_box(slide1, Emu(500000), Emu(2700000), Emu(7000000), Emu(400000),
             "От зависимости к контролю. Подход к импортозамещению для IT и автоматизации.",
             font_size=15, font_color=RGBColor(0xBB, 0xCC, 0xDD), bold=False)

# Bottom info line
add_shape_rect(slide1, 0, Emu(4400000), SLIDE_W, Emu(1), fill_color=RGBColor(0x33, 0x44, 0x66))
add_text_box(slide1, Emu(500000), Emu(4500000), Emu(5000000), Emu(300000),
             "АО «Атомик Софт»  •  automiq.ru  •  Конфиденциально",
             font_size=11, font_color=MED_GRAY)

# Meeting badge
add_shape_rect(slide1, Emu(6500000), Emu(4450000), Emu(2300000), Emu(350000), fill_color=RGBColor(0x22, 0x33, 0x55))
add_text_box(slide1, Emu(6500000), Emu(4480000), Emu(2300000), Emu(300000),
             "Внутренняя встреча • 2026",
             font_size=11, font_color=RGBColor(0x88, 0x99, 0xBB), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 2: СИТУАЦИЯ — ПОЧЕМУ НУЖНО ДЕЙСТВОВАТЬ
# ═══════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
add_header_bar(slide2, "Ситуация: почему нужно действовать сейчас")
add_bottom_bar(slide2)

# Three risk columns
risks = [
    ("⚠️  Санкционный риск", BLUE, [
        "AVEVA прекратила продажи и поддержку в РФ",
        "Нет обновлений безопасности с 2022 г.",
        "Нет техподдержки на русском языке",
        "Контрактные обязательства не исполняются",
    ]),
    ("💰  Финансовый риск", CYAN, [
        "AVEVA Flex: ~$108K/год за подписку",
        "При неоплате — ПО перестаёт работать",
        "Валютный риск (расчёты в долларах)",
        "Непредсказуемый бюджет, нет контроля",
    ]),
    ("📋  Регуляторный риск", GREEN, [
        "ФЗ-187 о безопасности КИИ",
        "Указ Президента №250 (2022)",
        "Сроки перехода ужесточаются",
        "Регулярные проверки ФСТЭК",
    ]),
]

col_w = Emu(2700000)
gap = Emu(150000)
start_x = Emu(250000)

for i, (title, color, items) in enumerate(risks):
    x = Emu(start_x + i * (col_w + gap))
    y = Emu(950000)
    
    # Column header
    add_shape_rect(slide2, x, y, col_w, Emu(300000), fill_color=color)
    add_text_box(slide2, Emu(x + 80000), Emu(y + 55000), Emu(col_w - 160000), Emu(200000),
                 title, font_size=12, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Items
    add_bullet_list(slide2, Emu(x + 80000), Emu(y + 370000), Emu(col_w - 160000),
                    items, font_size=10, bullet_char="•")

# Bottom callout
callout = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(250000), Emu(4350000), Emu(8600000), Emu(380000))
callout.fill.solid()
callout.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xE0)
callout.line.fill.background()
add_text_box(slide2, Emu(450000), Emu(4400000), Emu(8200000), Emu(300000),
             "⏱  Каждый месяц без решения — $9K подписки + растущий риск отключения + нарушение требований КИИ",
             font_size=11, font_color=RGBColor(0x8B, 0x4A, 0x00), bold=True)


# ═══════════════════════════════════════════════
# SLIDE 3: РЕШЕНИЕ — АЛЬФА ПЛАТФОРМА
# ═══════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
add_header_bar(slide3, "Решение: Альфа платформа",
               "Российская SCADA-платформа №1 в реестре Минцифры — АО «Атомик Софт», Томск")
add_bottom_bar(slide3)

# Metrics row
metrics = [
    ("6 500+", "инсталляций\nв России"),
    ("120+", "сотрудников"),
    ("25+", "коммуникационных\nпротоколов"),
    ("2 000 000", "тегов\nмакс. ёмкость"),
]
for i, (num, label) in enumerate(metrics):
    add_metric_block(slide3, Emu(300000 + i * 2200000), Emu(950000), num, label)

# Key advantages
advantages = [
    "✅  Реестр Минцифры + лицензия ФСТЭК на разработку средств защиты",
    "✅  Бессрочная лицензия в рублях — предсказуемый бюджет, нет подписок",
    "✅  Windows + Linux (Astra Linux, РЕД ОС) — сертифицированные ОС",
    "✅  Техподдержка 24/7 на русском языке, вендор в г. Томск",
    "✅  Собственный Historian для архивирования данных",
    "✅  Веб-доступ к мнемосхемам через Alpha.HMI.WebViewer",
]
y = Emu(1900000)
for adv in advantages:
    add_text_box(slide3, Emu(400000), y, Emu(8300000), Emu(280000),
                 adv, font_size=12, font_color=DARK_SLATE)
    y += Emu(310000)


# ═══════════════════════════════════════════════
# SLIDE 4: ЧТО МЕНЯЕТСЯ — СУТЬ ПЕРЕХОДА
# ═══════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
add_header_bar(slide4, "Что меняется: суть перехода",
               "PI System и Альфа платформа — разные по архитектуре, но совместимые по задачам")
add_bottom_bar(slide4)

add_two_column_block(
    slide4,
    "AVEVA PI System", [
        "Платформа управления данными (Data Infrastructure)",
        "Фокус: сбор → хранение → раздача",
        "Аналитикам, BI, облаку, Excel",
        "Подписка в долларах, облако за рубежом",
        "Управление процессом — не его задача",
    ],
    "Альфа платформа", [
        "SCADA-платформа (управление + данные)",
        "Фокус: оператор видит и управляет",
        "HMI + Historian + алармы + отчёты",
        "Бессрочная лицензия, рубли, on-premise",
        "Хранение данных тоже есть (Historian)",
    ],
    y_start=Emu(950000)
)

# Bottom insight
add_text_box(slide4, Emu(400000), Emu(3750000), Emu(8300000), Emu(800000),
             "Ключевой момент для руководства:\n"
             "PI System — это «умный склад данных». Альфа — это «диспетчерская с пультом управления».\n"
             "Данные PI можно мигрировать в Alpha.Historian. Экраны — пересоздать в Alpha.HMI.",
             font_size=11, font_color=DARK_SLATE)


# ═══════════════════════════════════════════════
# SLIDE 5: ПОКРЫТИЕ ФУНКЦИОНАЛА PI SYSTEM
# ═══════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)
add_header_bar(slide5, "Покрытие функционала PI System")
add_bottom_bar(slide5)

# Full replacement column
left_x = Emu(300000)
add_shape_rect(slide5, left_x, Emu(900000), Emu(4000000), Emu(300000), fill_color=GREEN)
add_text_box(slide5, Emu(left_x + 100000), Emu(935000), Emu(3800000), Emu(250000),
             "✅  Полная замена", font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

full_replace = [
    ("PI Data Archive", "Alpha.Historian 4.0"),
    ("PI ProcessBook", "Alpha.HMI 2.0 (мнемосхемы)"),
    ("PI Notifications", "Alpha.HMI.Alarms 3.3"),
    ("PI Buffer Subsystem", "Alpha.Server (буферизация)"),
    ("PI System Management", "Alpha.DevStudio + Security"),
    ("PI SDK / AF SDK", "Alpha.Server OPC UA + RMap"),
    ("PI Vision (веб)", "Alpha.HMI.WebViewer 2.0"),
]

y = Emu(1300000)
for pi_comp, alpha_comp in full_replace:
    add_text_box(slide5, Emu(left_x + 100000), y, Emu(1700000), Emu(230000),
                 pi_comp, font_size=10, font_color=MED_GRAY)
    add_text_box(slide5, Emu(left_x + 1900000), y, Emu(50000), Emu(230000),
                 "→", font_size=10, font_color=CYAN, bold=True)
    add_text_box(slide5, Emu(left_x + 2050000), y, Emu(1900000), Emu(230000),
                 alpha_comp, font_size=10, font_color=DARK_SLATE, bold=True)
    y += Emu(250000)

# Needs adaptation column
right_x = Emu(4800000)
add_shape_rect(slide5, right_x, Emu(900000), Emu(4000000), Emu(300000), fill_color=RGBColor(0xE8, 0x77, 0x22))
add_text_box(slide5, Emu(right_x + 100000), Emu(935000), Emu(3800000), Emu(250000),
             "🔄  Требует проработки", font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

adaptation = [
    "PI AF (Asset Framework) — иерархия через Alpha.Domain + Alpha.Om",
    "PI Analytics / Event Frames — Alpha.Om формулы + Alpha.Reports",
    "PI Integrator for BA — интеграция через Alpha.RMap (SQL-доступ)",
    "PI DataLink (Excel) — Excel-надстройка через OPC UA или RMap",
    "Кастомные скрипты — переработка на Alpha.Om / JavaScript",
]
add_bullet_list(slide5, Emu(right_x + 100000), Emu(1300000), Emu(3800000),
                adaptation, font_size=10, bullet_char="•")


# ═══════════════════════════════════════════════
# SLIDE 6: ОПЫТ В НЕФТЕГАЗЕ
# ═══════════════════════════════════════════════
slide6 = prs.slides.add_slide(blank_layout)
add_header_bar(slide6, "Опыт в нефтегазовой отрасли",
               "Альфа платформа уже работает на объектах уровня Новатэк")
add_bottom_bar(slide6)

cases = [
    ("НОВАТЭК Арктик СПГ 2", "Автоматизация технологических процессов СПГ-производства.\nАльфа платформа в составе АСУ ТП объекта."),
    ("ПАО «Транснефть» ВСТО-II", "Магистральный нефтепровод Восточная Сибирь — Тихий океан.\nДиспетчеризация и управление НПС."),
    ("Сочинская ТЭС", "Автоматизация тепловой электростанции.\nМнемосхемы + архивирование + тревоги."),
    ("СИБУР", "Крупнейший нефтехимический холдинг.\nАльфа в составе систем управления."),
]

y = Emu(950000)
for i, (name, desc) in enumerate(cases):
    col = i % 2
    row = i // 2
    x = Emu(300000 + col * 4400000)
    cy = Emu(y + row * 1650000)
    
    # Card background
    card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, Emu(4100000), Emu(1400000))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
    card.line.color.rgb = RGBColor(0xDD, 0xE5, 0xEE)
    card.line.width = Pt(1)
    
    # Accent stripe
    add_shape_rect(slide6, x, cy, Emu(50000), Emu(1400000), fill_color=CYAN)
    
    # Name
    add_text_box(slide6, Emu(x + 130000), Emu(cy + 120000), Emu(3800000), Emu(280000),
                 name, font_size=14, font_color=BLUE, bold=True)
    # Description
    add_text_box(slide6, Emu(x + 130000), Emu(cy + 420000), Emu(3800000), Emu(900000),
                 desc, font_size=10, font_color=DARK_SLATE)


# ═══════════════════════════════════════════════
# SLIDE 7: ТЕХНОЛОГИЧЕСКАЯ СОВМЕСТИМОСТЬ
# ═══════════════════════════════════════════════
slide7 = prs.slides.add_slide(blank_layout)
add_header_bar(slide7, "Технологическая совместимость",
               "Альфа платформа интегрируется с существующей инфраструктурой Новатэк")
add_bottom_bar(slide7)

# Protocol groups
groups = [
    ("Промышленные\nпротоколы", BLUE, [
        "OPC UA/DA (клиент + сервер)",
        "Modbus TCP/RTU",
        "IEC 60870-5-104/101",
        "IEC 61850",
        "S7 (Siemens)",
        "MQTT, SNMP, BACnet",
    ]),
    ("Операционные\nсистемы", CYAN, [
        "Windows 10/11, Server",
        "Astra Linux (сертиф.)",
        "РЕД ОС (сертиф.)",
        "Ubuntu, ALT Linux",
        "Кросс-платформенный",
        "HMI — включая веб",
    ]),
    ("Интеграция\nс IT", GREEN, [
        "SQL-доступ через Alpha.RMap",
        "OPC UA — стандарт Industry 4.0",
        "Syslog для SIEM",
        "LDAP/Active Directory",
        "REST API (планируется)",
        "Экспорт CSV/PDF/XLSX",
    ]),
]

for i, (title, color, items) in enumerate(groups):
    x = Emu(200000 + i * 3000000)
    y = Emu(950000)
    
    # Icon block
    add_shape_rect(slide7, x, y, Emu(2800000), Emu(380000), fill_color=color)
    add_text_box(slide7, Emu(x + 80000), Emu(y + 40000), Emu(2640000), Emu(320000),
                 title, font_size=12, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_bullet_list(slide7, Emu(x + 80000), Emu(y + 450000), Emu(2640000),
                    items, font_size=10, bullet_char="▸")


# ═══════════════════════════════════════════════
# SLIDE 8: ЭКОНОМИКА ПЕРЕХОДА
# ═══════════════════════════════════════════════
slide8 = prs.slides.add_slide(blank_layout)
add_header_bar(slide8, "Экономика перехода: TCO за 5 лет")
add_bottom_bar(slide8)

# PI System costs
add_text_box(slide8, Emu(400000), Emu(900000), Emu(3800000), Emu(300000),
             "PI System (AVEVA Flex)", font_size=16, font_color=RGBColor(0xCC, 0x33, 0x33), bold=True)
pi_costs = [
    "Подписка: ~$108K/год (~₽10.8М при ₽100/$)",
    "За 5 лет: ~₽54М только подписка",
    "Нет собственности — ПО арендуется",
    "Риск деактивации при неоплате",
    "Валютный риск при ослаблении рубля",
    "Поддержка недоступна (санкции)",
]
add_bullet_list(slide8, Emu(400000), Emu(1250000), Emu(3800000), pi_costs, font_size=11, bullet_char="✗")

# Alpha costs
add_text_box(slide8, Emu(4800000), Emu(900000), Emu(3800000), Emu(300000),
             "Альфа платформа", font_size=16, font_color=GREEN, bold=True)
alpha_costs = [
    "Разовая покупка лицензий в рублях",
    "Проект миграции: 8–14 месяцев",
    "Годовая техподдержка: ~15% от лицензии",
    "За 5 лет: значительно ниже TCO",
    "Бессрочная собственность на ПО",
    "Полная техподдержка в России",
]
add_bullet_list(slide8, Emu(4800000), Emu(1250000), Emu(3800000), alpha_costs, font_size=11, bullet_char="✓")

# ROI callout
roi_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(400000), Emu(3800000), Emu(8300000), Emu(750000))
roi_box.fill.solid()
roi_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
roi_box.line.fill.background()

txBox = slide8.shapes.add_textbox(Emu(600000), Emu(3870000), Emu(7900000), Emu(650000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Точка окупаемости: 2–3 года"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
p.font.bold = True
p.font.name = 'Arial'
p2 = tf.add_paragraph()
p2.text = "С учётом стоимости миграции, обучения персонала и годовой техподдержки — совокупная стоимость владения Альфа платформой за 5 лет существенно ниже, чем продление подписки PI System. Точные цифры — после аудита инфраструктуры."
p2.font.size = Pt(10)
p2.font.color.rgb = DARK_SLATE
p2.font.name = 'Arial'
p2.space_before = Pt(4)


# ═══════════════════════════════════════════════
# SLIDE 9: ПЛАН ПЕРЕХОДА
# ═══════════════════════════════════════════════
slide9 = prs.slides.add_slide(blank_layout)
add_header_bar(slide9, "План перехода: фазы и сроки",
               "Общая длительность 8–14 месяцев. Нулевой простой — PI и Альфа работают параллельно.")
add_bottom_bar(slide9)

phases = [
    ("1", "АУДИТ", "1–2 мес.", "Инвентаризация тегов, протоколов,\nскриптов, отчётов, Event Frames", BLUE),
    ("2", "ПРОЕКТ", "1 мес.", "Архитектура решения, лицензии,\nбюджет, рабочая группа", RGBColor(0x33, 0x75, 0xAD)),
    ("3", "ПИЛОТ", "2–3 мес.", "Один объект, параллельная работа\nPI + Альфа, валидация данных", CYAN),
    ("4", "МИГРАЦИЯ", "3–6 мес.", "Поэтапный перенос объектов,\nобучение персонала", GREEN),
    ("5", "ЗАВЕРШЕНИЕ", "1–2 мес.", "Вывод PI, оптимизация,\nпередача в эксплуатацию", RGBColor(0x61, 0x96, 0x2F)),
]

phase_w = Emu(1650000)
gap = Emu(60000)
start_x = Emu(200000)
y = Emu(950000)

for i, (num, name, duration, desc, color) in enumerate(phases):
    x = Emu(start_x + i * (phase_w + gap))
    
    # Phase number circle
    circle = slide9.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + 650000), y, Emu(350000), Emu(350000))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide9, Emu(x + 650000), Emu(y + 55000), Emu(350000), Emu(250000),
                 num, font_size=16, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Phase name
    add_text_box(slide9, x, Emu(y + 430000), phase_w, Emu(250000),
                 name, font_size=13, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Duration
    add_text_box(slide9, x, Emu(y + 680000), phase_w, Emu(200000),
                 duration, font_size=11, font_color=MED_GRAY, bold=False, alignment=PP_ALIGN.CENTER)
    
    # Description
    add_text_box(slide9, Emu(x + 30000), Emu(y + 950000), Emu(phase_w - 60000), Emu(600000),
                 desc, font_size=9, font_color=DARK_SLATE, alignment=PP_ALIGN.CENTER)

# Arrow connectors (simplified as thin rectangles)
for i in range(4):
    x = Emu(start_x + (i + 1) * (phase_w + gap) - gap)
    add_shape_rect(slide9, x, Emu(y + 170000), gap, Emu(12000), fill_color=MED_GRAY)

# Key principle callout
add_text_box(slide9, Emu(400000), Emu(4050000), Emu(8300000), Emu(600000),
             "Ключевой принцип: параллельная работа. PI System и Альфа работают одновременно на каждом этапе.\n"
             "Переключение только после подтверждения корректности данных. Нулевой простой.",
             font_size=11, font_color=BLUE, bold=True)


# ═══════════════════════════════════════════════
# SLIDE 10: РИСКИ И МИТИГАЦИЯ
# ═══════════════════════════════════════════════
slide10 = prs.slides.add_slide(blank_layout)
add_header_bar(slide10, "Риски и как мы их снимаем")
add_bottom_bar(slide10)

risks_table = [
    ("Функциональные разрывы\n(PI AF, Analytics)", "Маппинг функций на этапе аудита.\nPI AF → Alpha.Domain + Alpha.Om.\nНет 1:1 копии — есть решение задач."),
    ("Сопротивление персонала\n(привычка к PI)", "Обучение от вендора. Параллельная работа.\nПилот на одном объекте — доказательство.\nУчебные курсы Атомик Софт."),
    ("Потеря исторических данных", "Конвертация архивов PI → Alpha.Historian.\nФормат 4.0 поддерживает импорт.\nВерификация на этапе пилота."),
    ("Сроки и бюджет", "Фиксированные этапы с контрольными точками.\nPOC за 2–3 месяца. Бюджет в рублях.\nЕжемесячная отчётность."),
]

y = Emu(920000)
for i, (risk, mitigation) in enumerate(risks_table):
    row_y = Emu(y + i * 820000)
    
    # Risk
    risk_bg = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(250000), row_y, Emu(3400000), Emu(720000))
    risk_bg.fill.solid()
    risk_bg.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    risk_bg.line.fill.background()
    add_text_box(slide10, Emu(370000), Emu(row_y + 80000), Emu(3200000), Emu(600000),
                 risk, font_size=10, font_color=RGBColor(0xC6, 0x28, 0x28), bold=True)
    
    # Arrow
    add_text_box(slide10, Emu(3700000), Emu(row_y + 200000), Emu(400000), Emu(300000),
                 "→", font_size=18, font_color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Mitigation
    mit_bg = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(4200000), row_y, Emu(4700000), Emu(720000))
    mit_bg.fill.solid()
    mit_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    mit_bg.line.fill.background()
    add_text_box(slide10, Emu(4350000), Emu(row_y + 80000), Emu(4400000), Emu(600000),
                 mitigation, font_size=10, font_color=RGBColor(0x2E, 0x7D, 0x32))


# ═══════════════════════════════════════════════
# SLIDE 11: ПРЕДЛАГАЕМЫЕ ШАГИ
# ═══════════════════════════════════════════════
slide11 = prs.slides.add_slide(blank_layout)
add_header_bar(slide11, "Предлагаемые следующие шаги")
add_bottom_bar(slide11)

steps = [
    ("1", "Запуск аудита PI-инфраструктуры",
     "Инвентаризация: теги, протоколы, скрипты, отчёты, Event Frames.\nСрок: 1–2 месяца. Результат: отчёт с рекомендациями и бюджетом.", BLUE),
    ("2", "Определение пилотного объекта",
     "Один объект для практической проверки.\nPI и Альфа работают параллельно 2–4 недели.\nВалидация данных, обучение операторов.", CYAN),
    ("3", "Формирование рабочей группы",
     "IT + автоматизация + Атомик Софт.\nЕженедельные статусы, ежемесячные отчёты руководству.\nФиксированный бюджет и сроки.", GREEN),
]

y = Emu(1000000)
for num, title, desc, color in steps:
    # Number circle
    circle = slide11.shapes.add_shape(MSO_SHAPE.OVAL, Emu(350000), y, Emu(500000), Emu(500000))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide11, Emu(350000), Emu(y + 100000), Emu(500000), Emu(350000),
                 num, font_size=24, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Title
    add_text_box(slide11, Emu(1050000), Emu(y + 20000), Emu(7500000), Emu(300000),
                 title, font_size=16, font_color=BLUE, bold=True)
    
    # Description
    add_text_box(slide11, Emu(1050000), Emu(y + 280000), Emu(7500000), Emu(500000),
                 desc, font_size=11, font_color=DARK_SLATE)
    
    y += Emu(1100000)


# ═══════════════════════════════════════════════
# SLIDE 12: КОНТАКТЫ
# ═══════════════════════════════════════════════
slide12 = prs.slides.add_slide(blank_layout)

# Dark background
add_shape_rect(slide12, 0, 0, SLIDE_W, SLIDE_H, fill_color=DARK_BG)

# Decorative element
if deco_blob:
    img_stream = io.BytesIO(deco_blob)
    slide12.shapes.add_picture(img_stream, Emu(-800000), Emu(100000), Emu(4900000), Emu(4900000))

# Closing text
txBox = slide12.shapes.add_textbox(Emu(500000), Emu(800000), Emu(8000000), Emu(1200000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Будущее автоматизации"
p.font.size = Pt(32)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = 'Arial'
p2 = tf.add_paragraph()
p2.text = "в ваших руках."
p2.font.size = Pt(32)
p2.font.color.rgb = CYAN
p2.font.bold = True
p2.font.name = 'Arial'
p2.space_before = Pt(6)

# Contact info
contacts = [
    ("АО «Атомик Софт»", 14, WHITE, True),
    ("634050 Томск, пр. Ленина 60/1", 12, RGBColor(0xBB, 0xCC, 0xDD), False),
    ("", 8, WHITE, False),
    ("+7 (3822) 281 914", 12, RGBColor(0xBB, 0xCC, 0xDD), False),
    ("info@automiq.ru", 12, CYAN, False),
    ("www.automiq.ru", 12, CYAN, False),
]

txBox2 = slide12.shapes.add_textbox(Emu(2000000), Emu(2800000), Emu(3500000), Emu(1800000))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, (text, size, color, bold) in enumerate(contacts):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.space_before = Pt(2)

# Personal contact
personal = slide12.shapes.add_textbox(Emu(6000000), Emu(3200000), Emu(3000000), Emu(1000000))
tf3 = personal.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "Ваш контакт:"
p.font.size = Pt(10)
p.font.color.rgb = MED_GRAY
p.font.name = 'Arial'
p2 = tf3.add_paragraph()
p2.text = "Станислав Павловский"
p2.font.size = Pt(13)
p2.font.color.rgb = WHITE
p2.font.bold = True
p2.font.name = 'Arial'
p2.space_before = Pt(4)
p3 = tf3.add_paragraph()
p3.text = "Telegram: @Integraleus"
p3.font.size = Pt(11)
p3.font.color.rgb = CYAN
p3.font.name = 'Arial'
p3.space_before = Pt(2)

# Logo
if logo_blob:
    img_stream = io.BytesIO(logo_blob)
    slide12.shapes.add_picture(img_stream, Emu(500000), Emu(4300000), Emu(2200000), Emu(340000))

# Accent line bottom
add_shape_rect(slide12, 0, Emu(5093500), SLIDE_W, Emu(50000), fill_color=CYAN)


# ── Save ──
output_path = "NOVATEK_VP_IT_v2_new.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
