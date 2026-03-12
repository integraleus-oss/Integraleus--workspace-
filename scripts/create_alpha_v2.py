#!/usr/bin/env python3
"""Generate Alpha Licensing presentation - clean, reliable PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === COLORS ===
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1A, 0x24, 0x3B)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def dark_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK

def rect(slide, l, t, w, h, color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = align
    return tb

def bullets(slide, l, t, w, h, items, size=14, color=LIGHT, lead="•  "):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{lead}{item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(4)
    return tb

def title_bar(slide, text, subtitle=None):
    rect(slide, Inches(0), Inches(0), prs.slide_width, Pt(4), BLUE)
    txt(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
        text, size=32, color=WHITE, bold=True)
    if subtitle:
        txt(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4),
            subtitle, size=16, color=GRAY)

def card(slide, l, t, w, h, title, items, accent=BLUE):
    rect(slide, l, t, w, h, BG_CARD, accent)
    # accent strip at top
    rect(slide, l, t, w, Pt(3), accent)
    txt(slide, l + Inches(0.25), t + Inches(0.2), w - Inches(0.5), Inches(0.4),
        title, size=16, color=accent, bold=True)
    bullets(slide, l + Inches(0.25), t + Inches(0.65), w - Inches(0.5), h - Inches(0.8),
            items, size=12, color=LIGHT)


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
rect(s, Inches(0), Inches(0), prs.slide_width, Pt(4), BLUE)

txt(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1),
    'АЛЬФА ПЛАТФОРМА', size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.2), Inches(11.3), Inches(0.6),
    'Руководство по лицензированию', size=26, color=BLUE, align=PP_ALIGN.CENTER)

rect(s, Inches(5.5), Inches(4.1), Inches(2.3), Pt(2), TEAL)

txt(s, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5),
    'Выбор семейства  |  Архитектура  |  Расчёт  |  Пресейл', size=16, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.0), Inches(11.3), Inches(0.4),
    'Тариф 2026  |  Внутренний документ', size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Product Family
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_bar(s, 'Линейка продуктов Alpha', 'Три семейства — от простого к масштабируемому')

cw = Inches(3.7)
ch = Inches(4.3)
cy = Inches(1.9)
gap = Inches(0.45)
x0 = Inches(0.7)

card(s, x0, cy, cw, ch,
     'Alpha.One+',
     ['Односерверное решение',
      'До 50 000 тегов',
      '1 пользователь + 1 WEB',
      'Без резервирования',
      'Минимальный порог входа',
      'Идеален для малых объектов'],
     TEAL)

card(s, x0 + cw + gap, cy, cw, ch,
     'Alpha.SCADA',
     ['Типовой промышленный контур',
      'Основной + резервный сервер',
      'Множественные клиенты',
      'Full / WEB / Terminal доступ',
      'Historian и отчёты',
      'Классическая SCADA-архитектура'],
     BLUE)

card(s, x0 + 2*(cw + gap), cy, cw, ch,
     'Alpha.Platform',
     ['Мультисерверная архитектура',
      'Масштабирование без ограничений',
      'Драйверы включены в поставку',
      'Alpha.Link для внутреннего обмена',
      'Архитектурная гибкость',
      'Готовность к росту системы'],
     ORANGE)

# ============================================================
# SLIDE 3: Selection Matrix
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_bar(s, 'Матрица выбора семейства', 'Ключевые переключатели для принятия решения')

cw2 = Inches(3.7)
ch2 = Inches(1.8)
y1 = Inches(1.9)

card(s, Inches(0.5), y1, cw2, ch2,
     'Нужен резерв?  -->  SCADA',
     ['Базовый вариант Alpha.SCADA',
      'Резерв = 2x серверная лицензия одного уровня',
      'Historian резерв — отдельная лицензия'],
     RED)

card(s, Inches(0.5) + cw2 + Inches(0.4), y1, cw2, ch2,
     'До 50К тегов, 1 клиент?  -->  One+',
     ['Без резервирования',
      'Простой однопользовательский сценарий',
      'Допускается 1 WEB-клиент'],
     TEAL)

card(s, Inches(0.5) + 2*(cw2 + Inches(0.4)), y1, cw2, ch2,
     'Мультисервер / рост?  -->  Platform',
     ['Сложные интеграции',
      'Выход за рамки SCADA-подхода',
      'Масштабируемость на будущее'],
     ORANGE)

# Rules box
ry = Inches(4.2)
rect(s, Inches(0.5), ry, Inches(12.3), Inches(2.8), BG_CARD, BLUE)
txt(s, Inches(1), ry + Inches(0.2), Inches(11), Inches(0.4),
    'Правила формирования КП', size=20, color=BLUE, bold=True)

bullets(s, Inches(1), ry + Inches(0.7), Inches(5.5), Inches(2),
    ['По умолчанию: минимум 2 варианта в каждом КП',
     'Обычно: SCADA + Platform',
     'One+ — только для малого/простого сценария',
     'Финал: обязательная сверка policy + тариф'],
    size=13, color=LIGHT)

bullets(s, Inches(7), ry + Inches(0.7), Inches(5.5), Inches(2),
    ['Мультисерверный сценарий за пределами SCADA -> Platform',
     'Резервирование в ТЗ -> начинаем с SCADA',
     'Всегда указывать год тарифа',
     'Открытые вопросы — явно выписывать'],
    size=13, color=LIGHT)

# ============================================================
# SLIDE 4: Integration & Tags
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_bar(s, 'Интеграция и расчёт тегов', 'Правила обмена данными между системами')

card(s, Inches(0.5), Inches(1.9), Inches(5.8), Inches(2.3),
     'Межсистемный обмен',
     ['One+/SCADA -> Platform = внешняя интеграция (OPC UA)',
      'Теги считаются внешними у принимающей стороны',
      'Platform <-> Platform по Alpha.Link = внутренние теги',
      'По OPC/Modbus/IEC и др. — всегда внешние'],
     BLUE)

card(s, Inches(6.8), Inches(1.9), Inches(5.8), Inches(2.3),
     'Протоколы и драйверы',
     ['OPC UA / OPC DA',
      'Modbus TCP / RTU',
      'IEC 60870-5-101/104 / IEC 61850',
      'S7, BACnet, MQTT и другие',
      'Alpha.Platform: все драйверы включены'],
     TEAL)

card(s, Inches(0.5), Inches(4.6), Inches(5.8), Inches(2.5),
     'Информационная мощность (коэффициенты)',
     ['Simple:    DI x2    DO x2    AI x3     AO x1',
      'Medium:  DI x2    DO x5    AI x10   AO x1',
      'Complex: DI x3    DO x10  AI x10   AO x15',
      'Custom:   задаётся явно',
      'Порядок: сначала сигналы, потом уровень мощности'],
     ORANGE)

card(s, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.5),
     'Режимы расчёта тегов',
     ['signals_only — пересчёт из DI/DO/AI/AO',
      'tags_only — вручную заданные внешние теги',
      'combine — сигналы + ручные теги',
      'Выбирать ближайшую верхнюю ступень тарифа',
      'Минимум 2 сценария: базовый + консервативный'],
     GREEN)

# ============================================================
# SLIDE 5: Clients, Historian, Redundancy
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_bar(s, 'Клиенты, Historian и резервирование')

card(s, Inches(0.5), Inches(1.9), Inches(3.8), Inches(2.5),
     'Клиентский доступ',
     ['Full-клиент (CL-F)',
      'WEB-клиент (CL-RO)',
      'Terminal / RDP',
      'Обязательно уточнить распределение',
      'Или: «уточняется при заказе ключей»'],
     BLUE)

card(s, Inches(4.7), Inches(1.9), Inches(3.8), Inches(2.5),
     'WEB-клиенты',
     ['<= 5 одновременных — стандарт',
      '> 5 — выделенный WEB-сервер',
      '+ WEB-PORTAL обязательно',
      'One+ — только 1 WEB-клиент'],
     TEAL)

card(s, Inches(8.9), Inches(1.9), Inches(3.8), Inches(2.5),
     'Резервирование',
     ['SCADA: 2x серверная лицензия',
      'Historian: отдельная лицензия x2',
      'Platform: по правилу комплекта',
      'Без двойного учёта тегов'],
     RED)

card(s, Inches(0.5), Inches(4.8), Inches(5.8), Inches(2.3),
     'Historian',
     ['Рассчитывается отдельно от внешних тегов',
      'При задании в %: ceil(total_tags x percent)',
      'Ближайшая верхняя ступень тарифа',
      'Резерв Historian = отдельная лицензия'],
     ORANGE)

card(s, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.3),
     'Reports и SLA',
     ['Reports: серверный + клиентский профиль RPT',
      'SLA уровни: BASE / STD / OPT / PRM',
      'Не указан SLA -> по умолчанию BASE',
      'НДС 22% применяется к SLA'],
     GREEN)

# ============================================================
# SLIDE 6: Presale Pipeline
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_bar(s, 'Пресейл-пайплайн: 10 этапов ТКП')

steps = [
    ('0', 'Квалификация', 'Заказчик, объект, цель, дедлайн'),
    ('1', 'Преселекция', 'Выбор семейства: One+ / SCADA / Platform'),
    ('2', 'Тех. контур', 'Протоколы, драйверы, характер обмена'),
    ('3', 'Архитектура', 'Топология, резерв, серверные роли'),
    ('4', 'Модель данных', 'DI/DO/AI/AO или внешние теги, пересчёт'),
    ('5', 'Клиенты', 'Full / WEB / Terminal, распределение'),
    ('6', 'Модули', 'Historian, Reports, SLA, хранение'),
    ('7', 'Проверка', 'Что подтверждено / оценочно / требует уточнений'),
    ('8', 'Расчёт', 'Минимум 2 варианта: A (рекоменд.) + B (альтерн.)'),
    ('9', 'Валидация', 'Сверка policy + тариф + интеграционные правила'),
]

for i, (num, title, desc) in enumerate(steps):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    x = Inches(0.7) + col * Inches(6.3)
    y = Inches(1.9) + row * Inches(0.95)

    # number box
    nb = rect(s, x, y, Inches(0.45), Inches(0.45), BLUE if col == 0 else TEAL)
    tf = nb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER

    txt(s, x + Inches(0.6), y, Inches(2.5), Inches(0.35),
        title, size=16, color=WHITE, bold=True)
    txt(s, x + Inches(0.6), y + Inches(0.3), Inches(5.2), Inches(0.3),
        desc, size=12, color=GRAY)

# Final step bar
rect(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), BG_CARD, BLUE)
txt(s, Inches(1), Inches(6.75), Inches(11), Inches(0.4),
    '10 ->  Сборка ТКП: входные данные, варианты A/B, что входит / не входит, риски, год тарифа',
    size=14, color=LIGHT)

# ============================================================
# SLIDE 7: KP Format & Financial
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_bar(s, 'Формат коммерческого предложения')

card(s, Inches(0.5), Inches(1.9), Inches(5.8), Inches(4.8),
     'Структура ТКП',
     ['1. Подтверждённые входные данные',
      '2. Расчёт тегов и выбор ступеней тарифа',
      '3. Вариант A — рекомендованный',
      '4. Вариант B — альтернативный',
      '5. Финансовый блок:',
      '      Лицензии (без НДС)',
      '      SLA (без НДС)',
      '      НДС 22% на SLA',
      '      SLA (с НДС)',
      '      Итого к оплате',
      '6. Риски и открытые вопросы',
      '7. Год тарифа — указывать обязательно'],
     BLUE)

card(s, Inches(6.8), Inches(1.9), Inches(5.8), Inches(2.2),
     'Чек-лист перед отправкой',
     ['Протоколы и драйверы подтверждены',
      'Резерв подтверждён',
      'Профиль клиентов зафиксирован',
      'Внешние теги посчитаны корректно',
      'Даны 2 варианта',
      'Policy + тариф сверены'],
     GREEN)

card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.2),
     'Финансовая безопасность',
     ['Артикулы (SKU) — можно показывать',
      'Построчные цены по SKU — НЕЛЬЗЯ',
      'Допустима только общая цена системы',
      'Тарифные детали — только прямой контур',
      'Для рабочих чатов — share-safe версия'],
     RED)

# ============================================================
# SLIDE 8: Client Questionnaire
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_bar(s, 'Опросник для клиента', 'Стандартная форма сбора данных под ТКП')

cw3 = Inches(2.9)
ch3 = Inches(2.3)
r1 = Inches(1.9)
r2 = Inches(4.5)
g = Inches(0.3)

card(s, Inches(0.4), r1, cw3, ch3,
     'Общая рамка',
     ['Заказчик / объект',
      'Отрасль / тип объекта',
      'Цель проекта',
      'Срок / дедлайн ТКП'],
     BLUE)

card(s, Inches(0.4) + cw3 + g, r1, cw3, ch3,
     'Архитектура',
     ['Семейство: One+/SCADA/Platform',
      'Кол-во серверов',
      'Нужен ли резерв и где',
      'Количество площадок'],
     TEAL)

card(s, Inches(0.4) + 2*(cw3 + g), r1, cw3, ch3,
     'Данные и теги',
     ['Оценка внешних тегов',
      'Или сигналы DI/DO/AI/AO',
      'Historian: да/нет',
      'Объём тегов в историю'],
     ORANGE)

card(s, Inches(0.4) + 3*(cw3 + g), r1, cw3, ch3,
     'Интеграции',
     ['Протоколы: OPC, Modbus...',
      'Источник -> Получатель',
      'Внутренний или внешний',
      'Нестандартные протоколы'],
     GREEN)

card(s, Inches(0.4), r2, Inches(3.9), Inches(2.4),
     'Клиенты и доступ',
     ['Full-клиенты (кол-во)',
      'WEB одновременных',
      'Terminal / RDP ?',
      'Распределение по типам'],
     BLUE)

card(s, Inches(4.7), r2, Inches(3.9), Inches(2.4),
     'Доп. требования',
     ['Отчёты / аналитика',
      'SLA / поддержка',
      'Обучение персонала',
      'ОС, сеть, DMZ'],
     TEAL)

card(s, Inches(9.0), r2, Inches(3.7), Inches(2.4),
     'Открытые вопросы',
     ['Что пока неизвестно',
      'Кто подтверждает тех.часть',
      'Когда подтверждение',
      'Бюджетный коридор'],
     RED)

# ============================================================
# SLIDE 9: Key Principles
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
rect(s, Inches(0), Inches(0), prs.slide_width, Pt(4), BLUE)

txt(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
    'Ключевые принципы', size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
rect(s, Inches(5.5), Inches(2.4), Inches(2.3), Pt(2), TEAL)

principles = [
    ('Минимум 2 варианта в каждом КП', 'Рекомендованный + альтернативный', BLUE),
    ('Обязательная сверка перед отправкой', 'Policy + тариф + интеграционные правила', TEAL),
    ('Финансовая дисциплина', 'Общая цена — да, построчные цены по SKU — нет', ORANGE),
    ('Год тарифа', 'Всегда указывать актуальный тарифный год в ТКП', GREEN),
    ('Драйверы в Platform', 'Входят в поставку, отдельно не докупаются', RED),
]

for i, (title, desc, color) in enumerate(principles):
    y = Inches(2.9) + i * Inches(0.85)
    rect(s, Inches(1.5), y, Inches(10.3), Inches(0.7), BG_CARD, color)
    rect(s, Inches(1.5), y, Pt(5), Inches(0.7), color)
    txt(s, Inches(1.8), y + Inches(0.08), Inches(4.5), Inches(0.35),
        title, size=16, color=WHITE, bold=True)
    txt(s, Inches(1.8), y + Inches(0.38), Inches(9.5), Inches(0.3),
        desc, size=13, color=GRAY)

txt(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
    'Альфа Платформа  |  Тариф 2026  |  Внутренний документ', size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Save
out = '/root/.openclaw/workspace/agents/main/Alpha_Licensing_Guide_2026_v2.pptx'
prs.save(out)
print(f'OK: {out}')
