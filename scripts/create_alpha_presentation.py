#!/usr/bin/env python3
"""Generate a professional Alpha Platform Licensing presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)      # Deep navy
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)   # Bright blue
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xA0)   # Teal
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
CARD_BG = RGBColor(0x24, 0x24, 0x3E)
SUBTLE = RGBColor(0x99, 0x99, 0xBB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, items, accent=ACCENT_BLUE):
    card = add_shape(slide, left, top, width, height, CARD_BG, accent)
    add_accent_line(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), accent)
    add_text(slide, left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), Inches(0.5),
             title, font_size=18, color=accent, bold=True)
    add_bullet_list(slide, left + Inches(0.3), top + Inches(0.85), width - Inches(0.6), height - Inches(1.1),
                    items, font_size=13, color=LIGHT_GRAY, spacing=Pt(6))

# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Gradient-like accent bar at top
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
         'АЛЬФА ПЛАТФОРМА', font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
         'Руководство по лицензированию', font_size=30, color=ACCENT_BLUE, bold=False, align=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(4.0), Inches(3.3), ACCENT_TEAL)

add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
         'Выбор семейства  •  Архитектура  •  Расчёт  •  Пресейл', font_size=18, color=SUBTLE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
         '2026  |  Внутренний документ', font_size=14, color=SUBTLE, align=PP_ALIGN.CENTER)

# ===== SLIDE 2: Product Family Overview =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         'Линейка продуктов Alpha', font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_TEAL)
add_text(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.5),
         'Три семейства — от простого к масштабируемому', font_size=16, color=SUBTLE)

# Three product cards
card_w = Inches(3.6)
card_h = Inches(4.5)
gap = Inches(0.5)
start_x = Inches(0.8)
card_y = Inches(2.1)

# One+
add_card(slide, start_x, card_y, card_w, card_h,
         '⚡ Alpha.One+',
         ['Односерверное решение',
          'До 50 000 тегов',
          '1 пользователь + 1 WEB',
          'Без резервирования',
          'Идеален для малых объектов',
          'Минимальный порог входа'],
         ACCENT_TEAL)

# SCADA
add_card(slide, start_x + card_w + gap, card_y, card_w, card_h,
         '🏭 Alpha.SCADA',
         ['Типовой промышленный контур',
          'Основной + резервный сервер',
          'Множественные клиенты',
          'Full / WEB / Terminal доступ',
          'Historian и отчёты',
          'Классическая SCADA-архитектура'],
         ACCENT_BLUE)

# Platform
add_card(slide, start_x + 2*(card_w + gap), card_y, card_w, card_h,
         '🚀 Alpha.Platform',
         ['Мультисерверная архитектура',
          'Масштабирование без ограничений',
          'Драйверы включены в поставку',
          'Alpha.Link для внутреннего обмена',
          'Архитектурная гибкость',
          'Готовность к росту системы'],
         ORANGE)

# ===== SLIDE 3: Selection Matrix =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         'Матрица выбора семейства', font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_TEAL)

# Decision tree as cards
row1_y = Inches(1.8)
row2_y = Inches(4.2)

# Decision 1
add_card(slide, Inches(0.5), row1_y, Inches(3.7), Inches(2.0),
         '🔒 Нужен резерв?',
         ['ДА → базовый вариант Alpha.SCADA',
          'Резерв = 2× серверная лицензия',
          'одного уровня'],
         RED)

# Decision 2
add_card(slide, Inches(4.6), row1_y, Inches(3.7), Inches(2.0),
         '📊 До 50К тегов, 1 клиент?',
         ['ДА (без резерва) → Alpha.One+',
          'Простой сценарий',
          'Однопользовательский формат'],
         ACCENT_TEAL)

# Decision 3
add_card(slide, Inches(8.7), row1_y, Inches(4.0), Inches(2.0),
         '🌐 Мультисервер / рост?',
         ['ДА → Alpha.Platform',
          'Сложные интеграции',
          'Масштабируемость на будущее'],
         ORANGE)

# Golden rule
add_shape(slide, Inches(0.5), row2_y, Inches(12.2), Inches(2.5), CARD_BG, ACCENT_BLUE)
add_text(slide, Inches(1.0), row2_y + Inches(0.3), Inches(11.2), Inches(0.5),
         '📋 Жёсткие переключатели', font_size=22, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(1.0), row2_y + Inches(0.9), Inches(5.5), Inches(1.4),
    ['Мультисерверный сценарий за пределами SCADA → Platform',
     'Резервирование в требованиях → начинаем с SCADA',
     'Без резерва, ≤50К тегов, 1 клиент → One+ допустим'],
    font_size=14, color=LIGHT_GRAY)
add_bullet_list(slide, Inches(6.8), row2_y + Inches(0.9), Inches(5.5), Inches(1.4),
    ['По умолчанию: минимум 2 варианта в КП',
     'Обычно: SCADA + Platform',
     'Финал: обязательная сверка policy + тариф'],
    font_size=14, color=LIGHT_GRAY)

# ===== SLIDE 4: Integration Rules =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         'Правила интеграции и тегов', font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_TEAL)

# Integration rules
add_card(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.5),
         '🔗 Межсистемный обмен',
         ['One+/SCADA → Platform = внешняя интеграция (OPC UA)',
          'Теги считаются внешними у принимающей стороны',
          'Platform ↔ Platform по Alpha.Link = внутренние теги',
          'По OPC/Modbus/IEC — всегда внешние'],
         ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.5),
         '📡 Протоколы и драйверы',
         ['OPC UA / OPC DA',
          'Modbus TCP / RTU',
          'IEC 60870-5-101/104',
          'IEC 61850, S7, BACnet, MQTT',
          'Alpha.Platform: драйверы включены!'],
         ACCENT_TEAL)

# Tag calculation
add_card(slide, Inches(0.5), Inches(4.6), Inches(5.8), Inches(2.5),
         '🧮 Информационная мощность (коэффициенты)',
         ['Simple:   DI×2  DO×2  AI×3   AO×1',
          'Medium:  DI×2  DO×5  AI×10  AO×1',
          'Complex: DI×3  DO×10 AI×10  AO×15',
          'Custom:  задаётся явно'],
         ORANGE)

add_card(slide, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.5),
         '🔢 Режимы расчёта',
         ['signals_only — пересчёт DI/DO/AI/AO',
          'tags_only — вручную заданные внешние теги',
          'combine — сигналы + ручные теги',
          'Порядок: сначала сигналы → потом уровень мощности'],
         GREEN)

# ===== SLIDE 5: Clients & Historian =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         'Клиенты, Historian и резервирование', font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_TEAL)

add_card(slide, Inches(0.5), Inches(1.7), Inches(3.8), Inches(2.8),
         '👤 Клиентский доступ',
         ['Full-клиент (CL-F)',
          'WEB-клиент (CL-RO)',
          'Terminal / RDP',
          'Обязательно: уточнить',
          'terminal/RDP распределение',
          'Или пометка «при заказе»'],
         ACCENT_BLUE)

add_card(slide, Inches(4.7), Inches(1.7), Inches(3.8), Inches(2.8),
         '🌐 WEB-клиенты',
         ['≤5 одновременных — стандарт',
          '>5 — выделенный WEB-сервер',
          '+ WEB-PORTAL обязательно',
          'One+ — только 1 WEB-клиент',
          ''],
         ACCENT_TEAL)

add_card(slide, Inches(8.9), Inches(1.7), Inches(3.8), Inches(2.8),
         '🔒 Резервирование',
         ['SCADA: 2× серверная лицензия',
          'Historian: отдельная лицензия ×2',
          'Platform: по правилу комплекта',
          'Без двойного учёта тегов',
          'при одном внешнем источнике'],
         RED)

# Historian card
add_card(slide, Inches(0.5), Inches(4.9), Inches(5.8), Inches(2.2),
         '📦 Historian',
         ['Отдельно от внешних тегов системы',
          'При задании в % : ceil(total_tags × percent)',
          'Ближайшая верхняя ступень тарифа',
          'Резерв Historian = отдельная лицензия'],
         ORANGE)

add_card(slide, Inches(6.8), Inches(4.9), Inches(5.8), Inches(2.2),
         '📊 Reports & SLA',
         ['Reports: серверный + клиентский профиль RPT',
          'SLA уровни: BASE / STD / OPT / PRM',
          'По умолчанию (не указан) → BASE',
          'НДС 22% применяется к SLA'],
         GREEN)

# ===== SLIDE 6: Presale Pipeline =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         'Пресейл-пайплайн: 10 этапов', font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_TEAL)

steps_left = [
    ('0', 'Квалификация', 'Заказчик, объект, цель, дедлайн'),
    ('1', 'Преселекция', 'One+ / SCADA / Platform'),
    ('2', 'Тех. контур', 'Протоколы, драйверы, обмен'),
    ('3', 'Архитектура', 'Топология, резерв, серверы'),
    ('4', 'Модель данных', 'DI/DO/AI/AO, внешние теги'),
]
steps_right = [
    ('5', 'Клиенты', 'Full / WEB / Terminal'),
    ('6', 'Модули', 'Historian, Reports, SLA'),
    ('7', 'Проверка', 'Допущения и неизвестные'),
    ('8', 'Расчёт', 'Мин. 2 варианта: A + B'),
    ('9', 'Валидация', 'Policy + тариф + правила'),
]

for i, (num, title, desc) in enumerate(steps_left):
    y = Inches(1.7) + i * Inches(1.05)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_text(slide, Inches(1.5), y + Inches(0.02), Inches(2), Inches(0.35),
             title, font_size=17, color=WHITE, bold=True)
    add_text(slide, Inches(1.5), y + Inches(0.33), Inches(4.5), Inches(0.3),
             desc, font_size=13, color=SUBTLE)

for i, (num, title, desc) in enumerate(steps_right):
    y = Inches(1.7) + i * Inches(1.05)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_TEAL
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_text(slide, Inches(7.8), y + Inches(0.02), Inches(2), Inches(0.35),
             title, font_size=17, color=WHITE, bold=True)
    add_text(slide, Inches(7.8), y + Inches(0.33), Inches(4.5), Inches(0.3),
             desc, font_size=13, color=SUBTLE)

# Step 10
add_shape(slide, Inches(0.5), Inches(6.95), Inches(12.2), Inches(0.45), CARD_BG, ACCENT_BLUE)
add_text(slide, Inches(1.0), Inches(6.97), Inches(11), Inches(0.4),
         '10  ➤  Сборка ТКП: входные данные, архитектура, варианты A/B, что входит/не входит, риски, год тарифа',
         font_size=15, color=LIGHT_GRAY, bold=False)

# ===== SLIDE 7: Financial Output Format =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         'Формат коммерческого предложения', font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_TEAL)

# Structure of KP
add_card(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(5.0),
         '📄 Структура ТКП',
         ['1. Подтверждённые входные данные',
          '2. Расчёт тегов и выбор ступеней',
          '3. Вариант A — рекомендованный',
          '4. Вариант B — альтернативный',
          '5. Финансовый блок:',
          '     • Лицензии (без НДС)',
          '     • SLA (без НДС)',
          '     • НДС 22% на SLA',
          '     • SLA (с НДС)',
          '     • Итого к оплате',
          '6. Риски и открытые вопросы',
          '7. Год тарифа — обязательно'],
         ACCENT_BLUE)

# Rules
add_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.3),
         '✅ Чек-лист перед отправкой',
         ['Протоколы/драйверы подтверждены',
          'Резерв подтверждён',
          'Профиль клиентов подтверждён',
          'Внешние теги посчитаны корректно',
          'Даны 2 варианта (SCADA + Platform)',
          'Сверка policy + тариф выполнена'],
         GREEN)

add_card(slide, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.4),
         '⚠️ Финансовая безопасность',
         ['Артикулы (SKU) — можно показывать',
          'Построчные цены по SKU — НЕЛЬЗЯ',
          'Только общая цена системы',
          'Цены/тарифы — только прямой контур',
          'В публичных чатах — share-safe версия'],
         RED)

# ===== SLIDE 8: Questionnaire =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         'Опросник для клиента (под ТКП)', font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3), ACCENT_TEAL)

# 4 category cards
cw = Inches(2.85)
ch = Inches(2.6)
r1y = Inches(1.7)
r2y = Inches(4.6)

add_card(slide, Inches(0.4), r1y, cw, ch,
         '🏢 Общая рамка',
         ['Заказчик / объект',
          'Отрасль / тип объекта',
          'Цель проекта',
          'Срок / дедлайн',
          '1 или 2 варианта КП'],
         ACCENT_BLUE)

add_card(slide, Inches(0.4 + cw + 0.25*72*1), r1y, cw, ch,
         '🏗 Архитектура',
         ['Семейство: One+/SCADA/Platform',
          'Кол-во серверов',
          'Нужен резерв?',
          'Где именно резерв',
          'Кол-во площадок'],
         ACCENT_TEAL)

add_card(slide, Inches(0.4) + 2*(cw + Inches(0.25)), r1y, cw, ch,
         '📊 Данные',
         ['Оценка внешних тегов',
          'Или сигналы DI/DO/AI/AO',
          'Есть Historian?',
          'Сколько тегов в историю',
          'Режим расчёта'],
         ORANGE)

add_card(slide, Inches(0.4) + 3*(cw + Inches(0.25)), r1y, cw, ch,
         '🔗 Интеграции',
         ['OPC UA/DA, Modbus',
          'IEC 101/104, 61850',
          'S7, BACnet, MQTT',
          'Источник → Получатель',
          'Внутренний или внешний'],
         GREEN)

add_card(slide, Inches(0.4), r2y, Inches(4), Inches(2.4),
         '👤 Клиенты и доступ',
         ['Full-клиенты (кол-во)',
          'WEB одновременных',
          'Terminal/RDP?',
          'Распределение по клиентам'],
         ACCENT_BLUE)

add_card(slide, Inches(4.7), r2y, Inches(4), Inches(2.4),
         '📋 Доп. требования',
         ['Отчёты / аналитика',
          'SLA / поддержка',
          'Обучение',
          'ОС, виртуализация, DMZ'],
         ACCENT_TEAL)

add_card(slide, Inches(9.0), r2y, Inches(3.7), Inches(2.4),
         '❓ Неизвестные',
         ['Что пока не уточнено',
          'Кто подтверждает тех.часть',
          'Дата подтверждения',
          'Бюджетный коридор'],
         RED)

# ===== SLIDE 9: Summary =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.0),
         'Ключевые принципы', font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5), Inches(2.5), Inches(3.3), ACCENT_TEAL)

principles = [
    ('🎯', 'Минимум 2 варианта', 'В каждом КП: рекомендованный + альтернативный'),
    ('🔍', 'Сверка перед отправкой', 'Policy + тариф + интеграционные правила'),
    ('🔒', 'Финансовая дисциплина', 'Общая цена — да, построчные SKU-цены — нет'),
    ('📅', 'Год тарифа', 'Всегда указывать актуальный тарифный год'),
]

for i, (emoji, title, desc) in enumerate(principles):
    y = Inches(3.0) + i * Inches(1.0)
    add_shape(slide, Inches(2), y, Inches(9.3), Inches(0.8), CARD_BG)
    add_text(slide, Inches(2.3), y + Inches(0.1), Inches(0.6), Inches(0.6),
             emoji, font_size=24, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.0), y + Inches(0.1), Inches(3), Inches(0.4),
             title, font_size=18, color=WHITE, bold=True)
    add_text(slide, Inches(3.0), y + Inches(0.42), Inches(8), Inches(0.35),
             desc, font_size=14, color=SUBTLE)

out_path = '/root/.openclaw/workspace/agents/main/Alpha_Licensing_Guide_2026.pptx'
prs.save(out_path)
print(f'Saved to {out_path}')
