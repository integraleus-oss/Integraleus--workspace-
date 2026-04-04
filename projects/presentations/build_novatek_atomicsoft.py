#!/usr/bin/env python3
"""
Rebuild NOVATEK PI→Alpha migration presentation in Атомик Софт corporate style.
Uses atomicsoft-about.pptx as template base for theme colors and master slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# === Атомик Софт Color Palette ===
DEEP_BLUE = RGBColor(0x00, 0x54, 0x97)      # accent1 — primary brand
MEDIUM_BLUE = RGBColor(0x33, 0x75, 0xAD)     # accent2
TEAL = RGBColor(0x09, 0x97, 0xC8)            # accent3 — highlights
LIGHT_CYAN = RGBColor(0x52, 0xCC, 0xF7)      # accent4
GREEN = RGBColor(0x82, 0xC4, 0x44)           # accent5
LIGHT_GREEN = RGBColor(0xCD, 0xE7, 0xB4)     # accent6
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_SLATE = RGBColor(0x44, 0x54, 0x69)      # body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
AMBER = RGBColor(0xE8, 0x9C, 0x00)           # for warnings ⚠️
RED = RGBColor(0xCC, 0x33, 0x33)             # for gaps ❌

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.625)

# Margins
LEFT_M = Inches(0.55)
RIGHT_M = Inches(9.45)
TOP_CONTENT = Inches(1.15)
CONTENT_W = Inches(8.9)

TEMPLATE = '/root/.openclaw/workspace/agents/main/skills/presentation-designer/templates/atomicsoft-about.pptx'
OUTPUT = '/root/.openclaw/workspace/agents/main/projects/presentations/NOVATEK_PI_to_Alpha_v2.pptx'


def remove_all_slides(prs):
    """Remove all slides from a presentation."""
    sldIdLst = prs.slides._sldIdLst
    while len(sldIdLst):
        sldId = sldIdLst[0]
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def add_header_bar(slide, title_text, y_offset=Inches(0), height=Inches(0.95)):
    """Add the blue header bar with title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), y_offset, SLIDE_W, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DEEP_BLUE
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(LEFT_M, y_offset + Inches(0.15), Inches(8.5), height - Inches(0.15))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Arial'
    return bar


def add_accent_line(slide, y):
    """Add a teal accent line below header."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, LEFT_M, y, Inches(1.2), Inches(0.035)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_bullets(slide, bullets, x=None, y=None, width=None, height=None, font_size=12, color=None, spacing=6, line_spacing=None):
    """Add bullet points to slide. Returns the textbox shape."""
    x = x or Inches(0.75)
    y = y or Inches(1.3)
    width = width or Inches(8.5)
    height = height or Inches(4.0)
    color = color or DARK_SLATE

    txBox = slide.shapes.add_textbox(x, y, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle bold prefix (text before first —)
        if '—' in bullet:
            parts = bullet.split('—', 1)
            run1 = p.add_run()
            run1.text = parts[0].strip()
            run1.font.bold = True
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.name = 'Arial'
            run2 = p.add_run()
            run2.text = ' — ' + parts[1].strip()
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = 'Arial'
        else:
            p.text = bullet
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = 'Arial'

        p.space_after = Pt(spacing)
        if line_spacing:
            p.line_spacing = Pt(line_spacing)

    return txBox


def add_text(slide, text, x, y, w, h, size=12, bold=False, color=None, align=PP_ALIGN.LEFT, name='Arial'):
    """Add a simple text box."""
    color = color or DARK_SLATE
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return txBox


def add_metric_block(slide, value, label, x, y, w=Inches(2.0), val_color=None):
    """Add a large metric number with label below it."""
    val_color = val_color or TEAL
    # Value
    add_text(slide, value, x, y, w, Inches(0.6), size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    # Label
    add_text(slide, label, x, y + Inches(0.55), w, Inches(0.5), size=10, color=DARK_SLATE, align=PP_ALIGN.CENTER)


def add_two_column_table(slide, left_title, left_items, right_title, right_items, y_start=Inches(1.3)):
    """Add two-column comparison layout."""
    col_w = Inches(4.2)
    gap = Inches(0.3)
    left_x = LEFT_M
    right_x = LEFT_M + col_w + gap

    # Left column background
    lbg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, y_start, col_w, Inches(3.8))
    lbg.fill.solid()
    lbg.fill.fore_color.rgb = LIGHT_GRAY
    lbg.line.fill.background()

    # Right column background
    rbg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, y_start, col_w, Inches(3.8))
    rbg.fill.solid()
    rbg.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)  # light blue tint
    rbg.line.fill.background()

    # Left title
    add_text(slide, left_title, left_x + Inches(0.15), y_start + Inches(0.1), col_w - Inches(0.3), Inches(0.4),
             size=14, bold=True, color=DARK_SLATE)
    # Right title
    add_text(slide, right_title, right_x + Inches(0.15), y_start + Inches(0.1), col_w - Inches(0.3), Inches(0.4),
             size=14, bold=True, color=DEEP_BLUE)

    # Left bullets
    add_bullets(slide, left_items, left_x + Inches(0.15), y_start + Inches(0.5), col_w - Inches(0.3), Inches(3.2),
                font_size=10, spacing=4)
    # Right bullets
    add_bullets(slide, right_items, right_x + Inches(0.15), y_start + Inches(0.5), col_w - Inches(0.3), Inches(3.2),
                font_size=10, color=DEEP_BLUE, spacing=4)


def build():
    prs = Presentation(TEMPLATE)
    remove_all_slides(prs)

    # Use General Slide layout (index 3) for most slides
    layout_general = prs.slide_layouts[3]
    layout_title = prs.slide_layouts[0]  # Helium_Break_2
    layout_contact = prs.slide_layouts[4]  # 1_Contact Us
    layout_nofooter = prs.slide_layouts[2]  # No Footer

    # ================================================================
    # SLIDE 1 — Title
    # ================================================================
    slide = prs.slides.add_slide(layout_title)

    # Dark overlay
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = DEEP_BLUE
    overlay.line.fill.background()

    # Teal accent bar at top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    add_text(slide, 'Миграция с AVEVA PI System\nна Альфа платформу',
             LEFT_M, Inches(1.3), Inches(7.5), Inches(1.8),
             size=32, bold=True, color=WHITE)

    add_text(slide, 'Анализ возможностей, пробелов и план перехода',
             LEFT_M, Inches(3.2), Inches(7.5), Inches(0.5),
             size=16, color=LIGHT_CYAN)

    add_text(slide, '2026',
             LEFT_M, Inches(3.8), Inches(2), Inches(0.4),
             size=14, color=TEAL)

    # Bottom bar
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.0), SLIDE_W, Inches(0.625))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = RGBColor(0x00, 0x3D, 0x6E)
    bottom.line.fill.background()

    add_text(slide, 'АО «Атомик Софт»  |  automiq.ru',
             LEFT_M, Inches(5.1), Inches(5), Inches(0.4),
             size=11, color=RGBColor(0x80, 0xB0, 0xD0))

    # ================================================================
    # SLIDE 2 — Текущая ситуация
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Текущая ситуация: почему миграция актуальна')
    add_accent_line(slide, Inches(0.98))
    add_bullets(slide, [
        '• Санкции — AVEVA прекратила поддержку и продажу в РФ',
        '• AVEVA Flex — подписочная модель: ~$108K/год за 25-50K тегов (по данным публичных контрактов)',
        '• При прекращении оплаты — ПО перестаёт работать (в отличие от бессрочных лицензий)',
        '• Невозможность получить обновления, патчи безопасности, новые версии',
        '• Требования регуляторов: ФЗ-187 о КИИ требует перехода на отечественное ПО',
        '• Риск отключения облачных сервисов AVEVA (PI Cloud, AVEVA Connect)',
        '• Кадровый вопрос: новые специалисты обучаются на российском ПО',
    ], font_size=12, spacing=5)

    # ================================================================
    # SLIDE 3 — PI System архитектура
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'PI System — архитектура и компоненты')
    add_accent_line(slide, Inches(0.98))
    add_bullets(slide, [
        '• PI Server — Data Archive (архив тегов) + Asset Framework (модель активов) + Analysis Service + Notifications + Event Frames',
        '• Сбор данных — 450+ PI Interfaces, PI Connectors (OPC UA, UFL, Relay), PI Buffer Subsystem',
        '• Визуализация — PI Vision (веб), PI ProcessBook (десктоп), PI DataLink (Excel)',
        '• Аналитика — PI Analytics, Asset Analytics, PI Integrator for BA',
        '• Инфраструктура — PI to PI (репликация), PI Web API (REST), PI SDK/.NET',
        '• Облако — AVEVA Data Hub (PI Cloud), AVEVA Insight',
    ], font_size=12, spacing=5)

    # ================================================================
    # SLIDE 4 — Альфа платформа архитектура
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Альфа платформа — архитектура и модули')
    add_accent_line(slide, Inches(0.98))
    add_bullets(slide, [
        '• Ядро — Alpha.Server 6.4 (сигналы, ~30 коммуникационных модулей, резервирование, кластер)',
        '• Архив — Alpha.Historian 4.0 (форматы 3x/4x, журнал транзакций, сжатие LZMA)',
        '• Визуализация — Alpha.HMI 2.0 (дизайнер + визуализатор) + WebViewer 2.0 (веб)',
        '• Тревоги — Alpha.HMI.Alarms 3.3 (кроссплатформенный: Windows, Linux, веб)',
        '• Графики — встроенные графики Alpha.HMI (alpha.hmi.charts)',
        '• Отчёты — Alpha.Report.Generator 2.0 + Diagnostics.AuditReport',
        '• Разработка — DevStudio (объектная модель, Alpha.Om, JavaScript)',
        '• Безопасность — ФСТЭК, IEC 62443, подсистема безопасности',
    ], font_size=11, spacing=4)

    # ================================================================
    # SLIDE 5 — Покрытие: что заменяется полностью
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Покрытие функционала: что заменяется полностью')
    add_accent_line(slide, Inches(0.98))

    items = [
        ('✅ PI Data Archive', 'Alpha.Historian 4.0 — полноценный архив с журналом транзакций и сжатием'),
        ('✅ PI Notifications', 'Alpha.HMI.Alarms 3.3 + Diagnostics.AuditReport — тревоги + рассылка SMTP/Syslog/OPC'),
        ('✅ PI ProcessBook', 'Alpha.HMI 2.0 + WebViewer — полная замена десктопной и веб-визуализации'),
        ('✅ PI System Management', 'DevStudio + Alpha.Diag — конфигурирование, диагностика, мониторинг'),
        ('✅ PI Buffer Subsystem', 'Встроенная буферизация Alpha.Server — автоматическая, без настройки'),
        ('✅ PI SDK (программный доступ)', 'Alpha.Link + OPC UA SDK + SQL Connector — несколько вариантов доступа'),
    ]

    y = Inches(1.25)
    for label, desc in items:
        # Green marker
        marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, LEFT_M, y, Inches(0.08), Inches(0.45))
        marker.fill.solid()
        marker.fill.fore_color.rgb = GREEN
        marker.line.fill.background()

        add_text(slide, label, Inches(0.75), y, Inches(3.0), Inches(0.35), size=11, bold=True, color=DEEP_BLUE)
        add_text(slide, desc, Inches(3.8), y, Inches(5.5), Inches(0.45), size=10, color=DARK_SLATE)
        y += Inches(0.55)

    # ================================================================
    # SLIDE 6 — Покрытие: что заменяется частично
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Покрытие функционала: что заменяется частично')
    add_accent_line(slide, Inches(0.98))

    items = [
        ('⚠️ PI Interfaces (450+)', 'Alpha.Server (~30 модулей): OPC UA/DA/HDA/AE, Modbus TCP/RTU, IEC 104/101/61850, S7, EtherNet/IP, FINS, BACnet, MQTT, SNMP, SQL Connector и др.'),
        ('⚠️ PI Vision (self-service)', 'Alpha.HMI.WebViewer — веб-доступ к мнемосхемам и трендам, но без self-service построителя дашбордов'),
        ('⚠️ PI AF (Asset Framework)', 'DevStudio — объектная модель, но design-time (не runtime, как AF)'),
        ('⚠️ PI DataLink (Excel)', 'Alpha.Report.Generator 2.0 — отчёты, но не интерактивный Excel add-in'),
    ]

    y = Inches(1.25)
    for label, desc in items:
        marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, LEFT_M, y, Inches(0.08), Inches(0.55))
        marker.fill.solid()
        marker.fill.fore_color.rgb = AMBER
        marker.line.fill.background()

        add_text(slide, label, Inches(0.75), y, Inches(3.2), Inches(0.35), size=11, bold=True, color=DEEP_BLUE)
        add_text(slide, desc, Inches(0.75), y + Inches(0.3), Inches(8.3), Inches(0.35), size=10, color=DARK_SLATE)
        y += Inches(0.7)

    # ================================================================
    # SLIDE 7 — Функциональные пробелы
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Функциональные пробелы: что отсутствует')
    add_accent_line(slide, Inches(0.98))

    items = [
        ('❌ Event Frames', 'Контекстные события с привязкой к активам, вложенностью, захватом данных за период'),
        ('❌ PI Web API (REST)', 'Нет открытого REST API; доступ через Alpha.Link (TCP), OPC UA, SQL Connector'),
        ('❌ PI to PI (репликация)', 'Нет встроенной репликации данных между архивами'),
        ('❌ Edge Data Store', 'Нет лёгкого edge-решения для удалённых площадок'),
        ('❌ Manual Logger', 'Нет встроенного интерфейса ручного ввода данных'),
        ('❌ PI Integrator for BA', 'Нет прямой интеграции с BI (Power BI, Tableau)'),
    ]

    y = Inches(1.25)
    for label, desc in items:
        marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, LEFT_M, y, Inches(0.08), Inches(0.45))
        marker.fill.solid()
        marker.fill.fore_color.rgb = RED
        marker.line.fill.background()

        add_text(slide, label, Inches(0.75), y, Inches(3.0), Inches(0.35), size=11, bold=True, color=DEEP_BLUE)
        add_text(slide, desc, Inches(3.8), y, Inches(5.5), Inches(0.35), size=10, color=DARK_SLATE)
        y += Inches(0.55)

    # ================================================================
    # SLIDE 8 — Стратегия закрытия пробелов
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Стратегия закрытия пробелов')
    add_accent_line(slide, Inches(0.98))

    # Table-like layout
    headers = ['Пробел', 'Решение']
    rows = [
        ('Event Frames', 'Кастомная доработка на Alpha.Om + Alpha.HMI.Alarms + PostgreSQL'),
        ('REST API', 'Разработка REST-обёртки над Alpha.RMap (SQL → REST)'),
        ('Репликация данных', 'Периодический экспорт из Historian → PostgreSQL → импорт'),
        ('Edge Data Store', 'Alpha.Server в минимальной конфигурации + буферизация'),
        ('Ручной ввод данных', 'Веб-форма через Alpha.HMI.WebViewer + Alpha.Om скрипт'),
        ('BI-интеграция', 'SQL Connector → PostgreSQL → Power BI / Grafana'),
        ('450+ интерфейсов', 'Покрытие ~80% промышленных протоколов; остальное — OPC UA шлюзы'),
    ]

    # Header row
    hdr_y = Inches(1.2)
    add_text(slide, headers[0], LEFT_M, hdr_y, Inches(3.0), Inches(0.35), size=12, bold=True, color=WHITE)
    add_text(slide, headers[1], Inches(3.8), hdr_y, Inches(5.5), Inches(0.35), size=12, bold=True, color=WHITE)
    hdr_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_M, hdr_y, Inches(8.9), Inches(0.35))
    hdr_bar.fill.solid()
    hdr_bar.fill.fore_color.rgb = MEDIUM_BLUE
    hdr_bar.line.fill.background()
    # Re-add text on top
    add_text(slide, headers[0], LEFT_M + Inches(0.1), hdr_y + Inches(0.03), Inches(3.0), Inches(0.3), size=11, bold=True, color=WHITE)
    add_text(slide, headers[1], Inches(3.8), hdr_y + Inches(0.03), Inches(5.5), Inches(0.3), size=11, bold=True, color=WHITE)

    y = hdr_y + Inches(0.4)
    for i, (gap, solution) in enumerate(rows):
        bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_M, y, Inches(8.9), Inches(0.42))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = bg_color
        row_bg.line.fill.background()

        add_text(slide, gap, LEFT_M + Inches(0.1), y + Inches(0.05), Inches(2.8), Inches(0.3), size=10, bold=True, color=DEEP_BLUE)
        add_text(slide, solution, Inches(3.8), y + Inches(0.05), Inches(5.3), Inches(0.3), size=10, color=DARK_SLATE)
        y += Inches(0.42)

    # ================================================================
    # SLIDE 9 — Ключевое архитектурное различие
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Ключевое архитектурное различие')
    add_accent_line(slide, Inches(0.98))

    add_two_column_table(
        slide,
        'AVEVA PI System',
        [
            '• Data Infrastructure — инфраструктура данных',
            '• Фокус: сбор → хранение → доставка данных',
            '• Визуализация — вторична (PI Vision)',
            '• Asset Framework — runtime-модель активов',
            '• REST API, SDK — открытая экосистема',
            '• Облако (AVEVA Data Hub, AVEVA Connect)',
        ],
        'Альфа платформа',
        [
            '• SCADA-платформа — управление техпроцессом',
            '• Фокус: визуализация + управление + данные',
            '• HMI/SCADA — ядро платформы',
            '• DevStudio — design-time модель объектов',
            '• Проприетарные протоколы + OPC UA + SQL',
            '• Только on-premises',
        ],
    )

    # ================================================================
    # SLIDE 10 — Преимущества перехода
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Преимущества перехода на Альфа платформу')
    add_accent_line(slide, Inches(0.98))
    add_bullets(slide, [
        '• Реестр Минцифры — соответствие требованиям ФЗ-187 о КИИ',
        '• Бессрочная лицензия — нет зависимости от ежегодных подписок',
        '• Linux: Astra Linux, РЕД ОС, Альт, Ubuntu — полная поддержка',
        '• Единая платформа: SCADA + HMI + Historian + Alarms + Reports',
        '• Техподдержка на русском языке 24/7, вендор в Томске',
        '• 6500+ инсталляций, 80+ партнёров-интеграторов',
        '• Лицензия ФСТЭК на разработку СЗКИ',
        '• Соответствие п.29.3 Приказа №239 ФСТЭК',
    ], font_size=12, spacing=5)

    # ================================================================
    # SLIDE 11 — Сравнение протоколов
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Сравнение протоколов и коннекторов')
    add_accent_line(slide, Inches(0.98))
    add_bullets(slide, [
        '• OPC UA/DA/HDA/AE — клиент и сервер (UA кроссплатформенный, DA/AE только Windows)',
        '• Modbus TCP/RTU — Master и Slave',
        '• IEC 60870-5-104/101 — Master и Slave',
        '• IEC 61850 — Client',
        '• Siemens S7 — Client (прямое подключение к S7-300/400/1200/1500)',
        '• EtherNet/IP (Allen-Bradley) — Client',
        '• FINS (Omron) — Client',
        '• BACnet IP — Client (здания и инфраструктура)',
        '• MQTT — Publisher/Subscriber (IIoT, Edge)',
        '• SNMP — Manager (мониторинг сетевого оборудования)',
        '• SQL Connector — PostgreSQL, MS SQL, Oracle, MySQL',
    ], font_size=11, spacing=3)

    # ================================================================
    # SLIDE 12 — План миграции: 5 фаз
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'План миграции: 5 фаз')
    add_accent_line(slide, Inches(0.98))

    phases = [
        ('1', 'АУДИТ', '1-2 мес.', 'Инвентаризация PI-инфраструктуры:\nтеги, AF-модели, интерфейсы,\nотчёты, Event Frames, интеграции'),
        ('2', 'АРХИТЕКТУРА', '1 мес.', 'Проектирование целевой архитектуры,\nвыбор редакции, расчёт лицензий,\nопределение доработок'),
        ('3', 'ПИЛОТ', '2-3 мес.', 'Развёртывание на пилотном объекте,\nмиграция тегов и мнемосхем,\nтестирование'),
        ('4', 'МИГРАЦИЯ', '3-6 мес.', 'Последовательный перевод объектов,\nпараллельная работа PI + Альфа,\nобучение персонала'),
        ('5', 'ЗАКРЫТИЕ', '1-2 мес.', 'Отключение PI System,\nперенос исторических данных,\nфинальная документация'),
    ]

    # Phase blocks
    block_w = Inches(1.65)
    gap = Inches(0.12)
    start_x = LEFT_M
    y_top = Inches(1.2)

    for i, (num, name, duration, desc) in enumerate(phases):
        x = start_x + i * (block_w + gap)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), y_top, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        add_text(slide, num, x + Inches(0.55), y_top + Inches(0.05), Inches(0.5), Inches(0.4),
                size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Phase name
        add_text(slide, name, x, y_top + Inches(0.6), block_w, Inches(0.35),
                size=12, bold=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER)

        # Duration
        add_text(slide, duration, x, y_top + Inches(0.9), block_w, Inches(0.25),
                size=10, color=TEAL, align=PP_ALIGN.CENTER)

        # Description box
        desc_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_top + Inches(1.2), block_w, Inches(2.0))
        desc_bg.fill.solid()
        desc_bg.fill.fore_color.rgb = LIGHT_GRAY
        desc_bg.line.fill.background()

        add_text(slide, desc, x + Inches(0.1), y_top + Inches(1.3), block_w - Inches(0.2), Inches(1.8),
                size=9, color=DARK_SLATE)

        # Arrow between phases (except last)
        if i < len(phases) - 1:
            arrow_x = x + block_w + Inches(0.01)
            add_text(slide, '›', arrow_x, y_top + Inches(1.8), Inches(0.12), Inches(0.4),
                    size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 13 — Оценка рисков
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Оценка рисков и митигация')
    add_accent_line(slide, Inches(0.98))

    risks = [
        ('⚠️ Потеря функционала Event Frames', 'Аудит использования EF, при необходимости — кастомная доработка'),
        ('⚠️ Несовместимость формата архива', 'Параллельная запись, экспорт/импорт исторических данных'),
        ('⚠️ Переобучение персонала', 'Программа обучения, 5 лабораторий в ВУЗах, онлайн-курсы'),
        ('⚠️ Простой при переходе', 'Параллельная работа PI + Альфа, поэтапная миграция'),
        ('⚠️ Недостающие протоколы', 'OPC UA шлюзы для экзотических устройств'),
        ('⚠️ Потеря облачных сервисов', 'On-prem аналитика, Grafana/Power BI через SQL Connector'),
    ]

    y = Inches(1.2)
    for risk, mitigation in risks:
        add_text(slide, risk, LEFT_M, y, Inches(4.0), Inches(0.35), size=11, bold=True, color=DEEP_BLUE)
        add_text(slide, '→ ' + mitigation, Inches(4.8), y, Inches(4.6), Inches(0.4), size=10, color=DARK_SLATE)
        y += Inches(0.5)

    # ================================================================
    # SLIDE 14 — Экономика TCO
    # ================================================================
    slide = prs.slides.add_slide(layout_nofooter)
    add_header_bar(slide, 'Экономика: TCO сравнение (ориентировочно)')
    add_accent_line(slide, Inches(0.98))

    # Metric blocks
    add_metric_block(slide, '~$108K', 'AVEVA Flex / год\n(25-50K тегов)', Inches(0.5), Inches(1.5), Inches(2.5), RED)
    add_metric_block(slide, '~$325K', 'AVEVA Flex / 3 года\n(публичный контракт)', Inches(3.2), Inches(1.5), Inches(2.5), RED)

    # vs divider
    add_text(slide, 'VS', Inches(5.9), Inches(1.8), Inches(0.7), Inches(0.5),
            size=18, bold=True, color=MEDIUM_BLUE, align=PP_ALIGN.CENTER)

    add_metric_block(slide, '₽', 'Альфа платформа\nбессрочная лицензия', Inches(6.8), Inches(1.5), Inches(2.5), GREEN)

    # Additional info box
    info_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, LEFT_M, Inches(3.2), Inches(8.9), Inches(1.8))
    info_bg.fill.solid()
    info_bg.fill.fore_color.rgb = LIGHT_GRAY
    info_bg.line.fill.background()

    add_bullets(slide, [
        '• Разовая закупка vs ежегодная подписка — принципиально разная модель владения',
        '• Нет валютных рисков — расчёты в рублях',
        '• Нет risk of deactivation — лицензия бессрочная, работает без интернета',
        '• Стоимость миграции окупается за 1.5-3 года в зависимости от масштаба',
    ], x=Inches(0.75), y=Inches(3.3), width=Inches(8.3), height=Inches(1.6), font_size=11, spacing=4)

    # ================================================================
    # SLIDE 15 — Следующие шаги
    # ================================================================
    slide = prs.slides.add_slide(layout_general)
    add_header_bar(slide, 'Следующие шаги')
    add_accent_line(slide, Inches(0.98))

    steps = [
        'Согласование объёма аудита PI-инфраструктуры НОВАТЭК',
        'Определение пилотного объекта / участка для миграции',
        'Составление перечня критичных функций PI System (Event Frames? PI DataLink? API?)',
        'Расчёт лицензий Альфа платформы под целевую архитектуру',
        'Формирование рабочей группы: НОВАТЭК + Атомик Софт + интегратор',
        'Разработка детального плана-графика и бюджета',
    ]

    y = Inches(1.3)
    for i, step in enumerate(steps):
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, LEFT_M, y, Inches(0.4), Inches(0.4))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = TEAL if i < 3 else MEDIUM_BLUE
        num_circle.line.fill.background()
        add_text(slide, str(i + 1), LEFT_M + Inches(0.02), y + Inches(0.04), Inches(0.36), Inches(0.32),
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, step, Inches(1.1), y + Inches(0.03), Inches(8.0), Inches(0.35),
                size=12, color=DARK_SLATE)
        y += Inches(0.55)

    # ================================================================
    # SLIDE 16 — Контакты
    # ================================================================
    slide = prs.slides.add_slide(layout_contact)

    # Dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DEEP_BLUE
    bg.line.fill.background()

    # Teal accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    add_text(slide, 'Контакты', LEFT_M, Inches(1.0), Inches(5), Inches(0.7),
            size=32, bold=True, color=WHITE)

    add_text(slide, 'Станислав Павловский', LEFT_M, Inches(2.2), Inches(5), Inches(0.5),
            size=18, bold=True, color=WHITE)

    add_text(slide, 'Telegram: @Integraleus', LEFT_M, Inches(2.8), Inches(5), Inches(0.4),
            size=14, color=LIGHT_CYAN)

    # Company info
    add_text(slide, 'АО «Атомик Софт»\n634050 Томск, пр. Ленина 60/1\n+7 (3822) 281 914\ninfo@automiq.ru | automiq.ru',
            LEFT_M, Inches(3.6), Inches(5), Inches(1.5),
            size=12, color=RGBColor(0x80, 0xB0, 0xD0))

    add_text(slide, 'Будущее автоматизации\nв ваших руках.',
            Inches(6.0), Inches(2.0), Inches(3.5), Inches(1.0),
            size=20, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)

    prs.save(OUTPUT)
    print(f'✅ Saved: {OUTPUT}')
    print(f'   Slides: {len(prs.slides)}')


if __name__ == '__main__':
    build()
