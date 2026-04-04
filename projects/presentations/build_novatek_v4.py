#!/usr/bin/env python3
"""
NOVATEK PI→Alpha v4 — True Атомик Софт brand style.
Uses actual brand elements: header with flanking lines, white bg, metric tiles,
real logo, architecture diagram, proper fonts and colors.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# === True Атомик Софт Brand Colors ===
BRAND_BLUE = RGBColor(0x00, 0x54, 0x97)      # accent1 — headers, titles
BLUE_TEXT = RGBColor(0x00, 0x55, 0x9A)        # main heading color
CYAN = RGBColor(0x11, 0xA0, 0xD7)            # accent highlights
MEDIUM_BLUE = RGBColor(0x33, 0x75, 0xAD)     # accent2
TEAL = RGBColor(0x09, 0x97, 0xC8)            # accent3 — links
GREEN = RGBColor(0x82, 0xC4, 0x44)           # accent5
BODY_TEXT = RGBColor(0x44, 0x54, 0x69)        # body text #445469
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
LINE_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
AMBER = RGBColor(0xE8, 0x9C, 0x00)
RED_ACCENT = RGBColor(0xCC, 0x33, 0x33)

ASSETS = '/tmp/atomicsoft_assets'
TEMPLATE = '/root/.openclaw/workspace/agents/main/skills/presentation-designer/templates/atomicsoft-about.pptx'
OUTPUT = '/root/.openclaw/workspace/agents/main/projects/presentations/NOVATEK_PI_to_Alpha_v4.pptx'

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.625)


def remove_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    while len(sldIdLst):
        rId = sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldIdLst[0])


# ─────────────────────────────────────────────────────────────────────
# Brand elements
# ─────────────────────────────────────────────────────────────────────

def add_brand_header(slide, title_text):
    """
    Атомик Софт header: blue bold text with horizontal lines on both sides.
    Matches the pattern from their corporate presentations.
    """
    y = Inches(0.22)
    h = Inches(0.50)
    
    # Left line
    left_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), y + Inches(0.25), Inches(0.5), Inches(0.015))
    left_line.fill.solid()
    left_line.fill.fore_color.rgb = BRAND_BLUE
    left_line.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.55), y, Inches(8.5), h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.font.name = 'Arial'

    # Right line
    right_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(8.5), y + Inches(0.25), Inches(1.5), Inches(0.015))
    right_line.fill.solid()
    right_line.fill.fore_color.rgb = BRAND_BLUE
    right_line.line.fill.background()


def add_logo(slide, x=Inches(0.30), y=Inches(0.27)):
    """Add the Атомик Софт logo."""
    logo_path = os.path.join(ASSETS, 'logo.png')
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, x, y, Inches(2.48), Inches(0.38))


def add_dotted_separator(slide, x, y, width):
    """Row of small circles as separator (like in the original)."""
    dot_d = Inches(0.025)
    gap = Inches(0.04)
    n = int(width / (dot_d + gap * 914400) * 914400)
    n = min(n, 40)
    for i in range(n):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + i * (dot_d + Inches(0.035)), y, dot_d, dot_d)
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL
        dot.line.fill.background()


def txt(slide, text, x, y, w, h, size=12, bold=False, color=None, align=PP_ALIGN.LEFT, name='Arial'):
    color = color or BODY_TEXT
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return txBox


def rich_txt(slide, parts, x, y, w, h, align=PP_ALIGN.LEFT):
    """Parts: list of (text, size, bold, color)"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (text, size, bold, color) in enumerate(parts):
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = 'Arial'
    return txBox


def metric_tile(slide, number, description, x, y, w=Inches(2.7), h=Inches(0.85), title_size=26):
    """
    Атомик Софт metric tile: large blue number on top, description below.
    Exact pattern from slide 2 of their corporate presentation.
    """
    txt(slide, number, x, y, w, Inches(0.4), size=title_size, bold=True, color=BLUE_TEXT)
    txt(slide, description, x, y + Inches(0.38), w, Inches(0.45), size=11, color=BODY_TEXT)


def capability_block(slide, title, description, icon_path, x, y, w=Inches(4.1), h=Inches(1.1)):
    """
    Capability block with icon + bold title + description.
    Pattern from slide 6 (capabilities).
    """
    # Icon
    if icon_path and os.path.exists(icon_path):
        slide.shapes.add_picture(icon_path, x, y + Inches(0.05), Inches(0.4), Inches(0.4))
    
    # Title (bold, blue)
    txt(slide, title, x + Inches(0.5), y, w - Inches(0.5), Inches(0.3),
        size=13, bold=True, color=BRAND_BLUE)
    # Description
    txt(slide, description, x + Inches(0.5), y + Inches(0.3), w - Inches(0.5), h - Inches(0.3),
        size=10, color=BODY_TEXT)


def box(slide, x, y, w, h, fill_color, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s


def thin_line(slide, x, y, w, color=None):
    color = color or LINE_GRAY
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.01))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def status_marker(slide, status, x, y):
    """✅ ⚠️ ❌ marker as colored circle with symbol."""
    colors = {'ok': GREEN, 'warn': AMBER, 'gap': RED_ACCENT}
    symbols = {'ok': '✓', 'warn': '!', 'gap': '✕'}
    clr = colors.get(status, GREEN)
    sym = symbols.get(status, '?')
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.22), Inches(0.22))
    c.fill.solid()
    c.fill.fore_color.rgb = clr
    c.line.fill.background()
    txt(slide, sym, x + Inches(0.02), y + Inches(0.01), Inches(0.18), Inches(0.18),
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation(TEMPLATE)
    remove_all_slides(prs)

    layout_general = prs.slide_layouts[3]  # General Slide (white bg, small logo bottom-right)
    layout_title = prs.slide_layouts[0]     # Helium_Break_2
    layout_nf = prs.slide_layouts[2]        # No Footer
    layout_contact = prs.slide_layouts[4]   # 1_Contact Us

    # ════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title (photo background like real АС title)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_title)

    # Background photo (clipped to slide height)
    bg_path = os.path.join(ASSETS, 'title_bg.jpg')
    if os.path.exists(bg_path):
        s.shapes.add_picture(bg_path, 0, 0, SLIDE_W, SLIDE_H)

    # Logo
    add_logo(s)

    # Title
    txt(s, 'Миграция с AVEVA PI System\nна Альфа платформу',
        Inches(0.30), Inches(1.80), Inches(5.9), Inches(1.4),
        size=28, bold=True, color=WHITE)

    # Dotted separator
    add_dotted_separator(s, Inches(0.54), Inches(3.45), Inches(2.0))

    # Subtitle
    txt(s, 'анализ возможностей, пробелов и план перехода',
        Inches(0.30), Inches(3.55), Inches(5.0), Inches(0.5),
        size=16, color=WHITE)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 2 — Почему миграция (metric tiles like slide 2 of АС)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Почему миграция с PI System актуальна')

    # Row 1 — metric tiles
    metric_tile(s, '~$108K/год', 'стоимость AVEVA Flex\nподписки за 25-50K тегов', Inches(0.6), Inches(1.0))
    metric_tile(s, '~$325K', 'AVEVA Flex за 3 года\n(публичный контракт)', Inches(3.5), Inches(1.0))
    metric_tile(s, '0', 'обновлений и патчей\nот AVEVA для РФ', Inches(6.8), Inches(1.0))

    # Divider
    thin_line(s, Inches(0.6), Inches(2.1), Inches(8.8), TEAL)

    # Row 2 — text blocks (like slide 4 capability blocks)
    blocks = [
        ('САНКЦИИ', 'AVEVA прекратила поддержку\nи продажу в России.\nНет патчей, нет новых версий.'),
        ('ФЗ-187 О КИИ', 'Требование перехода на\nотечественное ПО для объектов\nкритической информационной\nинфраструктуры.'),
        ('RISK OF DEACTIVATION', 'Подписочная модель Flex:\nпри неоплате ПО\nперестаёт работать.\nОблачные сервисы — отключены.'),
        ('КАДРЫ', 'Новые специалисты\nобучаются на российском ПО.\nЭкспертиза по PI System\nсокращается.'),
    ]

    bw = Inches(2.1)
    gap = Inches(0.15)
    for i, (title, desc) in enumerate(blocks):
        x = Inches(0.6) + i * (bw + gap)
        y = Inches(2.3)
        txt(s, title, x, y, bw, Inches(0.3), size=13, bold=True, color=BRAND_BLUE)
        txt(s, desc, x, y + Inches(0.35), bw, Inches(1.5), size=10, color=BODY_TEXT)

    # Bottom labels row (like аналитика / разработка / тестирование)
    labels = ['подписка', 'обновления', 'поддержка', 'облако', 'кадры']
    for i, label in enumerate(labels):
        lx = Inches(0.6) + i * Inches(1.8)
        add_dotted_separator(s, lx, Inches(4.9), Inches(0.3))
        txt(s, label, lx + Inches(0.4), Inches(4.85), Inches(1.2), Inches(0.25),
            size=9, color=BODY_TEXT)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 3 — PI System архитектура (layered, clean style)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'PI System — архитектура и компоненты')

    layers = [
        ('Облако', 'AVEVA Data Hub  •  AVEVA Insight  •  AVEVA Connect', RGBColor(0x7B, 0x68, 0xAC)),
        ('Аналитика', 'PI Analytics  •  Asset Analytics  •  PI Integrator for BA', MEDIUM_BLUE),
        ('Визуализация', 'PI Vision (веб)  •  PI ProcessBook (десктоп)  •  PI DataLink (Excel)', TEAL),
        ('Сервер', 'Data Archive  •  Asset Framework  •  Notifications  •  Event Frames  •  Analysis', BRAND_BLUE),
        ('Сбор данных', '450+ PI Interfaces  •  PI Connectors (OPC UA, UFL, Relay)  •  Buffer Subsystem', RGBColor(0x3D, 0x7A, 0x44)),
    ]

    ly = Inches(1.0)
    lh = Inches(0.65)
    lgap = Inches(0.08)
    lx = Inches(0.6)
    lw = Inches(8.8)

    for label, content, clr in layers:
        # Light background
        bg = box(s, lx, ly, lw, lh, LIGHT_GRAY, radius=False)
        # Colored left accent
        box(s, lx, ly, Inches(0.05), lh, clr, radius=False)
        # Label
        txt(s, label, lx + Inches(0.15), ly + Inches(0.08), Inches(2.0), Inches(0.3),
            size=12, bold=True, color=clr)
        # Content
        txt(s, content, lx + Inches(2.2), ly + Inches(0.08), Inches(6.4), Inches(0.5),
            size=10, color=BODY_TEXT)
        ly += lh + lgap

    # Side annotations
    txt(s, 'PI Web API (REST)\nPI SDK (.NET)\nPI to PI', Inches(0.6), ly + Inches(0.15), Inches(3.0), Inches(0.7),
        size=9, color=TEAL)
    txt(s, 'Открытая экосистема:\nSDK, REST, 450+ коннекторов', Inches(5.0), ly + Inches(0.15), Inches(4.0), Inches(0.5),
        size=9, color=BODY_TEXT)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 4 — Альфа платформа архитектура
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Альфа платформа — архитектура и модули')

    # Use the real architecture diagram from Атомик Софт
    arch_path = os.path.join(ASSETS, 'architecture.png')
    if os.path.exists(arch_path):
        s.shapes.add_picture(arch_path, Inches(0.5), Inches(0.65), Inches(8.0), Inches(4.9))

    # ════════════════════════════════════════════════════════════════
    # SLIDE 5 — Что такое Альфа платформа (like slide 5 of АС)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_nf)

    # SCADA background
    scada_path = os.path.join(ASSETS, 'scada_bg.jpg')
    if os.path.exists(scada_path):
        s.shapes.add_picture(scada_path, 0, 0, SLIDE_W, SLIDE_H)

    # Semi-transparent overlay on left
    overlay = box(s, 0, 0, Inches(6.0), SLIDE_H, BRAND_BLUE, radius=False)
    # Make semi-transparent
    solidFill = overlay._element.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_elem = srgb.makeelement(qn('a:alpha'), {'val': '85000'})
            srgb.append(alpha_elem)

    # Header
    txt(s, 'Альфа платформа — замена PI System', Inches(0.5), Inches(0.30), Inches(5.0), Inches(0.5),
        size=20, bold=True, color=WHITE)
    thin_line(s, Inches(0.5), Inches(0.75), Inches(1.5), WHITE)

    txt(s, 'supervisory control and data acquisition', Inches(0.5), Inches(0.85), Inches(5.0), Inches(0.35),
        size=11, color=CYAN)

    # Key points with ">" markers (like original slide 5)
    points = [
        ('Сбор данных', '~30 модулей: OPC UA, Modbus, IEC 104, S7, MQTT...'),
        ('Хранение', 'Alpha.Historian 4.0: журнал транзакций, сжатие LZMA'),
        ('Управление', 'Alpha.Server 6.4: сигналы, резервирование, кластер'),
        ('Визуализация', 'Alpha.HMI 2.0 + Alpha.HMI.WebViewer 2.0 + alpha.hmi.charts'),
        ('Тревоги', 'Alpha.HMI.Alarms 3.3: кроссплатформенный'),
        ('Отчёты', 'Alpha.Reports 1.1 + Alpha.Diagnostics 2.2'),
    ]

    py = Inches(1.3)
    for label, desc in points:
        txt(s, '›', Inches(0.5), py, Inches(0.3), Inches(0.3), size=18, bold=True, color=CYAN)
        txt(s, label, Inches(0.8), py, Inches(1.5), Inches(0.3), size=12, bold=True, color=WHITE)
        txt(s, desc, Inches(0.8), py + Inches(0.25), Inches(4.5), Inches(0.25), size=9, color=RGBColor(0xB0, 0xD0, 0xE8))
        py += Inches(0.55)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 6 — Покрытие: полная замена (capability blocks like slide 6)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Покрытие функционала: полная замена')

    icons = [
        'icon_graphic_3.png', 'icon_graphic_5.png', 'icon_graphic_21.png',
        'icon_graphic_23.png', 'icon_graphic_25.png', 'icon_graphic_27.png',
    ]

    items = [
        ('Архив данных → Alpha.Historian 4.0',
         'Все данные с датчиков хранятся надёжно.\nСжатие LZMA — архив занимает меньше места.'),
        ('Тревоги → Alpha.HMI.Alarms 3.3',
         'Оператор видит аварии и подтверждает их.\nРаботает на Windows, Linux и в браузере.'),
        ('Экраны оператора → Alpha.HMI 2.0',
         'Мнемосхемы даже мощнее, чем в PI.\nРисуешь в дизайнере — видишь в визуализаторе.'),
        ('Администрирование → Alpha.DevStudio',
         'Единая среда: настройка, диагностика,\nмониторинг всей системы.'),
        ('Буферизация → встроена',
         'Если связь с сервером пропала — данные\nне теряются. Работает автоматически.'),
        ('Программный доступ → Alpha.Link',
         'Внешние программы подключаются через\nTCP, OPC UA или SQL — на выбор.'),
    ]

    # 2 columns × 3 rows (like slide 6 capability layout)
    col_w = Inches(4.3)
    row_h = Inches(1.1)
    gap_x = Inches(0.3)
    gap_y = Inches(0.15)

    for i, (title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.3) + col * (col_w + gap_x)
        y = Inches(1.0) + row * (row_h + gap_y)
        icon_path = os.path.join(ASSETS, icons[i]) if i < len(icons) else None
        capability_block(s, title, desc, icon_path, x, y, col_w, row_h)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 7 — Частичная замена + пробелы (like задачи slide 4)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Частичная замена и функциональные пробелы')

    # Two column layout matching slide 4 pattern
    col_w = Inches(4.2)

    # Left — Частичная замена
    txt(s, 'ЧАСТИЧНАЯ ЗАМЕНА', Inches(0.6), Inches(1.0), col_w, Inches(0.35),
        size=14, bold=True, color=BRAND_BLUE)

    partial = [
        ('450+ коннекторов → ~30', 'Основные протоколы есть. Для экзотики —\nOPC UA шлюзы. ~80% покрыто.'),
        ('Оператор сам рисует дашборд', 'В Альфе дашборды рисует инженер.\nОператор смотрит готовое в браузере.'),
        ('Модель активов «на лету»', 'В Альфе структуру проектируют заранее\nв Alpha.DevStudio (не меняют на ходу).'),
        ('Excel-отчёты (PI DataLink)', 'Есть Alpha.Reports 1.1, но это не\nинтерактивный Excel add-in.'),
    ]

    py = Inches(1.5)
    for title, desc in partial:
        status_marker(s, 'warn', Inches(0.6), py + Inches(0.02))
        txt(s, title, Inches(0.9), py, Inches(3.5), Inches(0.25), size=11, bold=True, color=BRAND_BLUE)
        txt(s, desc, Inches(0.9), py + Inches(0.25), Inches(3.5), Inches(0.5), size=9, color=BODY_TEXT)
        py += Inches(0.75)

    # Right — Пробелы
    rx = Inches(5.2)
    txt(s, 'ФУНКЦИОНАЛЬНЫЕ ПРОБЕЛЫ', rx, Inches(1.0), col_w, Inches(0.35),
        size=14, bold=True, color=BRAND_BLUE)

    gaps = [
        ('«Запись событий»', 'Нельзя сказать «запиши, что\nбыло с 10:15 до 11:30 на насосе»'),
        ('HTTP-доступ к данным', 'Нет REST API — внешние\nпрограммы идут через TCP/SQL'),
        ('Синхронизация серверов', 'Нет репликации данных\nмежду архивами на разных серверах'),
        ('Лёгкая версия для edge', 'Нет «мини-сервера» для\nудалённых площадок'),
        ('Ручной ввод данных', 'Нет встроенной формы для\nлабораторных анализов и обходов'),
        ('Прямая связь с BI', 'Нет коннектора в Power BI.\nОбход: SQL → PostgreSQL → BI'),
    ]

    py = Inches(1.5)
    for title, desc in gaps:
        status_marker(s, 'gap', rx, py + Inches(0.02))
        txt(s, title, rx + Inches(0.3), py, Inches(1.8), Inches(0.25), size=11, bold=True, color=BRAND_BLUE)
        txt(s, desc, rx + Inches(2.2), py, Inches(2.0), Inches(0.45), size=9, color=BODY_TEXT)
        py += Inches(0.55)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 8 — Стратегия закрытия пробелов
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Стратегия закрытия пробелов')

    rows = [
        ('Запись событий', 'Программируем на Alpha.Om + PostgreSQL', 'Доработка'),
        ('HTTP-доступ', 'Делаем REST-обёртку над SQL-данными', 'Обёртка'),
        ('Синхронизация', 'Экспорт из Historian в PostgreSQL', 'Периодически'),
        ('Edge-сервер', 'Alpha.Server в мини-конфигурации', 'Лёгкая версия'),
        ('Ручной ввод', 'Веб-форма через Alpha.HMI.WebViewer', 'Веб-форма'),
        ('BI-аналитика', 'SQL Connector → Grafana / Power BI', 'Через СУБД'),
        ('Экзотичные протоколы', '~30 модулей покрывают 80%, остальное — OPC UA шлюзы', 'Шлюзы'),
    ]

    # Table header
    hy = Inches(0.95)
    thin_line(s, Inches(0.6), hy, Inches(8.8), BRAND_BLUE)
    txt(s, 'Чего не хватает', Inches(0.6), hy + Inches(0.05), Inches(2.5), Inches(0.3),
        size=11, bold=True, color=BRAND_BLUE)
    txt(s, 'Как решаем', Inches(3.3), hy + Inches(0.05), Inches(3.5), Inches(0.3),
        size=11, bold=True, color=BRAND_BLUE)
    txt(s, 'Подход', Inches(7.0), hy + Inches(0.05), Inches(2.5), Inches(0.3),
        size=11, bold=True, color=BRAND_BLUE)
    thin_line(s, Inches(0.6), hy + Inches(0.35), Inches(8.8), BRAND_BLUE)

    ry = hy + Inches(0.45)
    for i, (gap_name, solution, approach) in enumerate(rows):
        txt(s, gap_name, Inches(0.6), ry, Inches(2.5), Inches(0.35),
            size=10, bold=True, color=BODY_TEXT)
        txt(s, solution, Inches(3.3), ry, Inches(3.5), Inches(0.35),
            size=10, color=BODY_TEXT)
        txt(s, approach, Inches(7.0), ry, Inches(2.5), Inches(0.35),
            size=10, color=TEAL)
        thin_line(s, Inches(0.6), ry + Inches(0.38), Inches(8.8), LINE_GRAY)
        ry += Inches(0.45)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 9 — Архитектурное различие: «Склад» vs «Диспетчерская»
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Ключевое различие: «Склад» vs «Диспетчерская»')

    col_w = Inches(4.2)

    # Left — PI System
    # Blue-ish background box
    box(s, Inches(0.4), Inches(0.9), Inches(4.4), Inches(4.2), LIGHT_GRAY)
    box(s, Inches(0.4), Inches(0.9), Inches(4.4), Inches(0.04), AMBER, radius=False)

    txt(s, '📦  «СКЛАД ДАННЫХ»', Inches(0.6), Inches(1.0), col_w, Inches(0.35),
        size=14, bold=True, color=BRAND_BLUE)
    txt(s, 'AVEVA PI System', Inches(0.6), Inches(1.35), col_w, Inches(0.25),
        size=11, bold=True, color=TEAL)
    txt(s, 'Огромный склад с идеальной системой\n'
        'учёта. Любой товар можно найти,\n'
        'любому курьеру выдать, в любую\n'
        'аналитику подключить.\n'
        'Но сам склад ничем не управляет.',
        Inches(0.6), Inches(1.65), col_w, Inches(1.3), size=10, color=BODY_TEXT)

    txt(s, '• 450+ коннекторов — почти всё на свете\n'
        '• Оператор сам рисует себе дашборд\n'
        '• REST API — данные по HTTP как сайт\n'
        '• Облако — данные из любой точки мира\n'
        '• Event Frames — «что произошло\n'
        '  с 10:15 до 11:30 на насосе №3»',
        Inches(0.6), Inches(3.05), col_w, Inches(1.8), size=10, color=BODY_TEXT)

    # Right — Альфа
    box(s, Inches(5.0), Inches(0.9), Inches(4.6), Inches(4.2), RGBColor(0xE8, 0xF4, 0xFD))
    box(s, Inches(5.0), Inches(0.9), Inches(4.6), Inches(0.04), TEAL, radius=False)

    txt(s, '🖥️  «ДИСПЕТЧЕРСКАЯ»', Inches(5.2), Inches(1.0), col_w, Inches(0.35),
        size=14, bold=True, color=BRAND_BLUE)
    txt(s, 'Альфа платформа', Inches(5.2), Inches(1.35), col_w, Inches(0.25),
        size=11, bold=True, color=TEAL)
    txt(s, 'Диспетчерская с пультом. Оператор\n'
        'видит всё, управляет всем, все тревоги\n'
        'приходят сюда. Склад тоже есть\n'
        '(Historian), но главное — пульт.',
        Inches(5.2), Inches(1.65), col_w, Inches(1.1), size=10, color=BODY_TEXT)

    txt(s, '• Мощный дизайнер мнемосхем\n'
        '• ~30 протоколов — все основные\n'
        '• Напрямую к оборудованию,\n'
        '  без OPC-прослоек\n'
        '• Linux (Astra, РЕД ОС) + Windows\n'
        '• Только on-premises = безопасно\n'
        '• Бессрочная лицензия в рублях',
        Inches(5.2), Inches(2.85), col_w, Inches(2.0), size=10, color=BODY_TEXT)

    # Bottom analogy
    txt(s, 'Миграция = переезд со «склада данных» на «диспетчерскую с собственным складом»',
        Inches(0.6), Inches(5.2), Inches(8.8), Inches(0.3),
        size=9, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 9b — Сравнительная таблица
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Сравнение: PI System vs Альфа платформа')

    # Table-style comparison
    comparison = [
        ('Что это', 'Инфраструктура данных', 'SCADA-платформа'),
        ('Главный фокус', 'Собрать и раздать данные', 'Показать и управлять'),
        ('Визуализация', 'Вторична (PI Vision — простой)', 'Ядро продукта (Alpha.HMI — мощный)'),
        ('Коннекторы', '450+ (почти всё на свете)', '~30 (все основные промышленные)'),
        ('Модель объектов', 'Runtime (меняешь на ходу)', 'Design-time (проектируешь заранее)'),
        ('API', 'REST (как веб-сайт)', 'Alpha.Link (TCP) + OPC UA + SQL'),
        ('Облако', 'Да (AVEVA Data Hub)', 'Нет (только on-prem — безопасно)'),
        ('Event Frames', 'Да (уникальная фича)', 'Нет (нужна доработка)'),
        ('ОС', 'Только Windows', 'Windows + Linux (Astra, РЕД ОС)'),
        ('Лицензия', 'Подписка ~$108K/год', 'Бессрочная, в рублях'),
        ('Поддержка в РФ', 'Нет (санкции)', '24/7 на русском языке'),
    ]

    # Headers
    hy = Inches(0.85)
    box(s, Inches(0.5), hy, Inches(1.8), Inches(0.3), BRAND_BLUE, radius=False)
    box(s, Inches(2.4), hy, Inches(3.4), Inches(0.3), RGBColor(0x80, 0x40, 0x40), radius=False)
    box(s, Inches(5.9), hy, Inches(3.6), Inches(0.3), RGBColor(0x00, 0x6B, 0x3F), radius=False)
    txt(s, '', Inches(0.6), hy + Inches(0.02), Inches(1.6), Inches(0.25), size=10, bold=True, color=WHITE)
    txt(s, 'AVEVA PI System', Inches(2.5), hy + Inches(0.02), Inches(3.2), Inches(0.25), size=10, bold=True, color=WHITE)
    txt(s, 'Альфа платформа', Inches(6.0), hy + Inches(0.02), Inches(3.4), Inches(0.25), size=10, bold=True, color=WHITE)

    ry = hy + Inches(0.35)
    for i, (param, pi_val, alpha_val) in enumerate(comparison):
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        box(s, Inches(0.5), ry, Inches(9.0), Inches(0.35), bg, radius=False)
        txt(s, param, Inches(0.6), ry + Inches(0.03), Inches(1.6), Inches(0.3), size=9, bold=True, color=BRAND_BLUE)
        txt(s, pi_val, Inches(2.5), ry + Inches(0.03), Inches(3.2), Inches(0.3), size=9, color=BODY_TEXT)
        txt(s, alpha_val, Inches(6.0), ry + Inches(0.03), Inches(3.4), Inches(0.3), size=9, color=BODY_TEXT)
        ry += Inches(0.37)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 10 — Преимущества (clean list, no overlap)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Преимущества Альфа платформы')

    benefits = [
        ('Реестр Минцифры', 'соответствие требованиям ФЗ-187 о КИИ'),
        ('Бессрочная лицензия', 'нет зависимости от ежегодных подписок'),
        ('Поддержка Linux', 'Astra Linux, РЕД ОС, Альт, Ubuntu'),
        ('Единая платформа', 'SCADA + HMI + Historian + Alarms + Reports'),
        ('Техподдержка 24/7', 'на русском языке, вендор в Томске'),
        ('6500+ инсталляций', 'в разных отраслях промышленности'),
        ('80+ партнёров', 'развитая сеть интеграторов с компетенциями'),
        ('Лицензия ФСТЭК', 'разработка СЗКИ, п.29.3 Приказа №239'),
    ]

    # Two columns, 4 items each — clean, no overlap
    col_w = Inches(4.0)
    py = Inches(1.0)

    for i, (title, desc) in enumerate(benefits):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * Inches(4.5)
        y = Inches(1.0) + row * Inches(0.9)

        # Teal number/bullet
        txt(s, '✓', x, y + Inches(0.02), Inches(0.3), Inches(0.3),
            size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        # Title
        txt(s, title, x + Inches(0.35), y, col_w - Inches(0.4), Inches(0.3),
            size=13, bold=True, color=BRAND_BLUE)
        # Description
        txt(s, desc, x + Inches(0.35), y + Inches(0.3), col_w - Inches(0.4), Inches(0.4),
            size=10, color=BODY_TEXT)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 11 — Протоколы (clean table)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Протоколы и коннекторы Alpha.Server')

    protocols = [
        ('OPC UA/DA/HDA/AE', 'Клиент и Сервер', 'UA кроссплатформенный, DA/AE — Windows'),
        ('Modbus TCP/RTU', 'Master и Slave', ''),
        ('IEC 60870-5-104/101', 'Master и Slave', ''),
        ('IEC 61850', 'Client', ''),
        ('Siemens S7', 'Client', 'Прямое подключение к S7-300/400/1200/1500'),
        ('EtherNet/IP', 'Client', 'Allen-Bradley'),
        ('FINS', 'Client', 'Omron'),
        ('BACnet IP', 'Client', 'Здания и инфраструктура'),
        ('MQTT', 'Pub/Sub', 'IIoT, Edge'),
        ('SNMP', 'Manager', 'Мониторинг сетевого оборудования'),
        ('SQL Connector', 'PG/MS/Ora/MySQL', 'Двусторонний обмен с СУБД'),
    ]

    # Table (left side, leaving room for metrics on right)
    hy = Inches(0.95)
    thin_line(s, Inches(0.6), hy, Inches(5.5), BRAND_BLUE)
    txt(s, 'Протокол', Inches(0.6), hy + Inches(0.05), Inches(2.2), Inches(0.25), size=11, bold=True, color=BRAND_BLUE)
    txt(s, 'Режим', Inches(2.9), hy + Inches(0.05), Inches(1.5), Inches(0.25), size=11, bold=True, color=BRAND_BLUE)
    txt(s, 'Примечание', Inches(4.5), hy + Inches(0.05), Inches(1.6), Inches(0.25), size=11, bold=True, color=BRAND_BLUE)
    thin_line(s, Inches(0.6), hy + Inches(0.32), Inches(5.5), BRAND_BLUE)

    ry = hy + Inches(0.38)
    for proto, mode, note in protocols:
        txt(s, proto, Inches(0.6), ry, Inches(2.2), Inches(0.3), size=10, bold=False, color=BODY_TEXT)
        txt(s, mode, Inches(2.9), ry, Inches(1.5), Inches(0.3), size=10, color=TEAL)
        if note:
            txt(s, note, Inches(4.5), ry, Inches(1.6), Inches(0.3), size=8, color=BODY_TEXT)
        thin_line(s, Inches(0.6), ry + Inches(0.3), Inches(5.5), LINE_GRAY)
        ry += Inches(0.33)

    # Summary metrics (right side, not overlapping table)
    metric_tile(s, '~30', 'коммуникационных модулей\nв Alpha.Server', Inches(6.5), Inches(1.0), Inches(3.0), title_size=22)
    metric_tile(s, '~80%', 'промышленных протоколов\nпокрыто «из коробки»', Inches(6.5), Inches(2.2), Inches(3.0), title_size=22)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 12 — Лицензии и архитектуры (real diagram)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_nf)
    add_brand_header(s, 'Архитектуры и комплекты лицензий')

    txt(s, 'Альфа платформа позволяет разрабатывать системы различных архитектур\n'
        'от локальной с Alpha.One+ до распределенной или многоуровневой системы с Alpha.Platform.',
        Inches(0.6), Inches(0.85), Inches(8.0), Inches(0.7), size=11, color=BODY_TEXT)

    # Real license architecture diagram
    lic_path = os.path.join(ASSETS, 'license_arch.png')
    if os.path.exists(lic_path):
        s.shapes.add_picture(lic_path, Inches(0.8), Inches(1.7), Inches(8.0), Inches(3.5))

    # ════════════════════════════════════════════════════════════════
    # SLIDE 13 — План миграции (timeline)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'План миграции: 5 фаз')

    phases = [
        ('1', 'АУДИТ', '1-2 мес.', 'Считаем всё: сколько\nтегов, какие протоколы,\nкакие скрипты,\nчто критично'),
        ('2', 'АРХИТЕКТУРА', '1 мес.', 'Рисуем будущую\nсистему: что куда\nставим, какие\nлицензии нужны'),
        ('3', 'ПИЛОТ', '2-3 мес.', 'Ставим Альфу на\nодин объект, гоняем\nпараллельно с PI,\nпроверяем'),
        ('4', 'МИГРАЦИЯ', '3-6 мес.', 'Переводим объекты\nпо одному. PI и Альфа\nработают вместе.\nОбучаем людей'),
        ('5', 'ЗАКРЫТИЕ', '1-2 мес.', 'Выключаем PI.\nПереносим историю.\nПишем документацию.\nСдаём заказчику'),
    ]

    # Timeline (horizontal line with dots)
    tl_y = Inches(2.4)
    thin_line(s, Inches(0.6), tl_y, Inches(8.8), TEAL)

    bw = Inches(1.55)
    bgap = Inches(0.2)

    for i, (num, name, duration, desc) in enumerate(phases):
        x = Inches(0.6) + i * (bw + bgap)

        # Dot on timeline
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + bw/2 - Inches(0.12), tl_y - Inches(0.12), Inches(0.24), Inches(0.24))
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL
        dot.line.fill.background()
        txt(s, num, x + bw/2 - Inches(0.12), tl_y - Inches(0.1), Inches(0.24), Inches(0.2),
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Phase name + duration above
        txt(s, name, x, Inches(1.0), bw, Inches(0.3), size=11, bold=True, color=BRAND_BLUE, align=PP_ALIGN.CENTER)
        txt(s, duration, x, Inches(1.3), bw, Inches(0.25), size=10, color=TEAL, align=PP_ALIGN.CENTER)

        # Description below
        txt(s, desc, x, Inches(2.75), bw, Inches(1.5), size=9, color=BODY_TEXT)

    # Total
    txt(s, 'Общая длительность: 8–14 месяцев', Inches(2.5), Inches(4.8), Inches(5.0), Inches(0.35),
        size=14, bold=True, color=BRAND_BLUE, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 14 — Риски
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Оценка рисков и митигация')

    risks = [
        ('warn', 'Потеря «записи событий»', 'Сначала проверяем: реально ли используются Event Frames. Если да — допрограммируем'),
        ('warn', 'Старые данные', 'Оба архива работают одновременно. Историю переносим отдельным этапом'),
        ('ok', 'Люди не знают Альфу', 'Обучение от Атомик Софт, 5 лабораторий в ВУЗах, онлайн-курсы'),
        ('warn', 'Простой при переходе', 'PI и Альфа работают параллельно — переключаемся без остановки'),
        ('ok', 'Не все протоколы есть', '~80% покрыто из коробки, остальное — через OPC UA шлюзы'),
        ('warn', 'Нет облака', 'Вместо AVEVA Cloud ставим Grafana / Power BI на своих серверах'),
    ]

    # Table
    hy = Inches(0.95)
    thin_line(s, Inches(0.6), hy, Inches(8.8), BRAND_BLUE)
    txt(s, '', Inches(0.6), hy + Inches(0.05), Inches(0.8), Inches(0.25), size=10, bold=True, color=BRAND_BLUE)
    txt(s, 'Что может пойти не так', Inches(1.5), hy + Inches(0.05), Inches(2.8), Inches(0.25), size=10, bold=True, color=BRAND_BLUE)
    txt(s, 'Что делаем', Inches(4.5), hy + Inches(0.05), Inches(5.0), Inches(0.25), size=10, bold=True, color=BRAND_BLUE)
    thin_line(s, Inches(0.6), hy + Inches(0.32), Inches(8.8), BRAND_BLUE)

    ry = hy + Inches(0.4)
    for level, risk, mitigation in risks:
        status_marker(s, level, Inches(0.75), ry + Inches(0.03))
        txt(s, risk, Inches(1.5), ry, Inches(2.8), Inches(0.4), size=10, bold=True, color=BODY_TEXT)
        txt(s, mitigation, Inches(4.5), ry, Inches(5.0), Inches(0.5), size=10, color=BODY_TEXT)
        thin_line(s, Inches(0.6), ry + Inches(0.5), Inches(8.8), LINE_GRAY)
        ry += Inches(0.58)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 15 — TCO (metric tiles)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Экономика: TCO сравнение')

    # AVEVA metrics
    metric_tile(s, '~$108K', 'AVEVA Flex / год\n(25-50K тегов)', Inches(0.6), Inches(1.0), Inches(2.5))
    metric_tile(s, '~$325K', 'AVEVA Flex / 3 года\n(публичный контракт)', Inches(3.3), Inches(1.0), Inches(2.5))

    # Альфа metric
    txt(s, '₽', Inches(7.0), Inches(1.0), Inches(2.5), Inches(0.6),
        size=40, bold=True, color=GREEN)
    txt(s, 'Альфа платформа\nбессрочная лицензия', Inches(7.0), Inches(1.55), Inches(2.5), Inches(0.5),
        size=11, color=BODY_TEXT)

    thin_line(s, Inches(0.6), Inches(2.3), Inches(8.8), TEAL)

    # Comparison points
    points = [
        ('Купил — и твоё', 'не платишь каждый год. Лицензия бессрочная'),
        ('Рубли, не доллары', 'курс не влияет. Бюджет предсказуемый'),
        ('Не отключат', 'работает без интернета, нет «деактивации» при неоплате'),
        ('Окупаемость 1.5–3 года', 'стоимость миграции отбивается за счёт отмены подписки AVEVA'),
    ]

    py = Inches(2.5)
    for title, desc in points:
        status_marker(s, 'ok', Inches(0.6), py + Inches(0.02))
        txt(s, title, Inches(0.9), py, Inches(2.3), Inches(0.3), size=11, bold=True, color=BRAND_BLUE)
        txt(s, desc, Inches(3.3), py, Inches(6.0), Inches(0.3), size=10, color=BODY_TEXT)
        py += Inches(0.45)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 16 — Следующие шаги
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_general)
    add_brand_header(s, 'Следующие шаги')

    steps = [
        'Согласование объёма аудита PI-инфраструктуры НОВАТЭК',
        'Определение пилотного объекта / участка для миграции',
        'Составление перечня критичных функций PI System\n(Event Frames? PI DataLink? API?)',
        'Расчёт лицензий Альфа платформы под целевую архитектуру',
        'Формирование рабочей группы: НОВАТЭК + Атомик Софт + интегратор',
        'Разработка детального плана-графика и бюджета',
    ]

    py = Inches(1.0)
    for i, step in enumerate(steps):
        # Number
        txt(s, str(i + 1), Inches(0.6), py, Inches(0.4), Inches(0.4),
            size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        # Step text
        txt(s, step, Inches(1.1), py + Inches(0.02), Inches(8.0), Inches(0.5),
            size=12, color=BODY_TEXT)
        if i < len(steps) - 1:
            thin_line(s, Inches(0.6), py + Inches(0.55), Inches(8.8), LINE_GRAY)
        py += Inches(0.63)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 17 — Контакты (like АС contact slide)
    # ════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(layout_contact)

    # Decorative circle (clipped to slide bounds)
    circle_path = os.path.join(ASSETS, 'contact_circle.png')
    if os.path.exists(circle_path):
        s.shapes.add_picture(circle_path, Inches(0), Inches(0.2), Inches(4.0), Inches(4.0))

    # Tagline
    txt(s, 'Будущее автоматизации', Inches(0.5), Inches(0.90), Inches(9.0), Inches(0.7),
        size=28, bold=True, color=BRAND_BLUE, align=PP_ALIGN.CENTER)
    txt(s, 'в ваших руках.', Inches(0.5), Inches(1.50), Inches(9.0), Inches(0.7),
        size=28, bold=True, color=BRAND_BLUE, align=PP_ALIGN.CENTER)

    # QR / decorative
    qr_path = os.path.join(ASSETS, 'qr_or_deco.png')
    if os.path.exists(qr_path):
        s.shapes.add_picture(qr_path, Inches(7.5), Inches(2.4), Inches(1.7), Inches(1.7))

    # Contacts — 3 columns
    txt(s, 'АО «Атомик Софт»\n634050 Томск\nпр. Ленина 60/1',
        Inches(0.6), Inches(4.40), Inches(2.5), Inches(1.1), size=11, color=BODY_TEXT)
    txt(s, '+7 (3822) 281 914\ninfo@automiq.ru\nwww.automiq.ru',
        Inches(3.5), Inches(4.40), Inches(2.5), Inches(1.1), size=11, color=BODY_TEXT)
    txt(s, 'Станислав Павловский\nTelegram: @Integraleus',
        Inches(6.5), Inches(4.40), Inches(3.0), Inches(1.1), size=11, color=BODY_TEXT)

    prs.save(OUTPUT)
    print(f'✅ Saved: {OUTPUT}')
    print(f'   Slides: {len(prs.slides)}')


if __name__ == '__main__':
    build()
