#!/usr/bin/env python3
"""
NOVATEK PI→Alpha v3 — Атомик Софт style with infographics & brainstorm diagrams.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math

# === Palette ===
DEEP_BLUE = RGBColor(0x00, 0x54, 0x97)
MEDIUM_BLUE = RGBColor(0x33, 0x75, 0xAD)
TEAL = RGBColor(0x09, 0x97, 0xC8)
LIGHT_CYAN = RGBColor(0x52, 0xCC, 0xF7)
GREEN = RGBColor(0x82, 0xC4, 0x44)
LIGHT_GREEN = RGBColor(0xCD, 0xE7, 0xB4)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_SLATE = RGBColor(0x44, 0x54, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
LIGHTER_GRAY = RGBColor(0xF8, 0xF8, 0xF8)
AMBER = RGBColor(0xE8, 0x9C, 0x00)
RED = RGBColor(0xCC, 0x33, 0x33)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)
LIGHT_TEAL_BG = RGBColor(0xE0, 0xF7, 0xFA)
LIGHT_GREEN_BG = RGBColor(0xF0, 0xF8, 0xE8)
LIGHT_AMBER_BG = RGBColor(0xFF, 0xF3, 0xE0)
LIGHT_RED_BG = RGBColor(0xFD, 0xE8, 0xE8)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
LIGHT_PURPLE_BG = RGBColor(0xF3, 0xE5, 0xF5)
DARK_BLUE_BG = RGBColor(0x00, 0x3D, 0x6E)

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.625)
LEFT_M = Inches(0.55)

TEMPLATE = '/root/.openclaw/workspace/agents/main/skills/presentation-designer/templates/atomicsoft-about.pptx'
OUTPUT = '/root/.openclaw/workspace/agents/main/projects/presentations/NOVATEK_PI_to_Alpha_v3.pptx'


def remove_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    while len(sldIdLst):
        sldId = sldIdLst[0]
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def add_header_bar(slide, title_text, height=Inches(0.85)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = DEEP_BLUE
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(LEFT_M, Inches(0.12), Inches(8.5), height - Inches(0.12))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Arial'
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_M, height, Inches(1.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def txt(slide, text, x, y, w, h, size=12, bold=False, color=None, align=PP_ALIGN.LEFT, name='Arial', anchor=None):
    color = color or DARK_SLATE
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.auto_size = None
        # Set vertical centering via XML
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', anchor)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return txBox


def multitext(slide, lines, x, y, w, h, size=11, color=None, spacing=4, bold_first=False):
    """Multiple paragraphs in one textbox."""
    color = color or DARK_SLATE
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(spacing)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def box(slide, x, y, w, h, fill_color, radius=True):
    """Rounded rectangle background."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s


def circle(slide, x, y, d, fill_color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s


def arrow_right(slide, x, y, w=Inches(0.4), h=Inches(0.25), color=None):
    """Draw a right-pointing arrow."""
    color = color or TEAL
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def arrow_down(slide, x, y, w=Inches(0.25), h=Inches(0.3), color=None):
    color = color or TEAL
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def connector_line(slide, x, y, w, h, color=None, dash=False):
    """Simple line using a thin rectangle."""
    color = color or TEAL
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, max(w, Inches(0.015)), max(h, Inches(0.015)))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def hexagon(slide, x, y, w, h, fill_color):
    s = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s


def build():
    prs = Presentation(TEMPLATE)
    remove_all_slides(prs)

    layout = prs.slide_layouts[3]  # General Slide
    layout_title = prs.slide_layouts[0]
    layout_nf = prs.slide_layouts[2]
    layout_contact = prs.slide_layouts[4]

    # ================================================================
    # SLIDE 1 — Title (dramatic)
    # ================================================================
    s = prs.slides.add_slide(layout_title)
    box(s, 0, 0, SLIDE_W, SLIDE_H, DEEP_BLUE, radius=False)
    # Decorative diagonal stripe
    box(s, Inches(6.5), 0, Inches(3.5), SLIDE_H, DARK_BLUE_BG, radius=False)
    # Teal top line
    box(s, 0, 0, SLIDE_W, Inches(0.05), TEAL, radius=False)

    txt(s, 'Миграция с AVEVA PI System\nна Альфа платформу', LEFT_M, Inches(1.0), Inches(6.0), Inches(2.0),
        size=32, bold=True, color=WHITE)
    txt(s, 'Анализ возможностей, пробелов и план перехода', LEFT_M, Inches(3.1), Inches(5.5), Inches(0.5),
        size=14, color=LIGHT_CYAN)

    # Decorative circles (brainstorm feel)
    for cx, cy, d, clr in [
        (Inches(7.2), Inches(0.8), Inches(1.8), TEAL),
        (Inches(8.3), Inches(2.3), Inches(1.0), MEDIUM_BLUE),
        (Inches(6.8), Inches(3.0), Inches(0.7), GREEN),
        (Inches(8.5), Inches(3.5), Inches(0.5), LIGHT_CYAN),
    ]:
        c = circle(s, cx, cy, d, clr)
        c.fill.fore_color.rgb = clr
        # Make semi-transparent via XML
        solidFill = c._element.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = qn('a:alpha')
                alpha_elem = srgb.makeelement(alpha, {'val': '25000'})
                srgb.append(alpha_elem)

    # Labels in circles
    txt(s, 'SCADA', Inches(7.5), Inches(1.3), Inches(1.2), Inches(0.4), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, 'HMI', Inches(8.4), Inches(2.55), Inches(0.8), Inches(0.3), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, 'IIoT', Inches(6.85), Inches(3.15), Inches(0.55), Inches(0.25), size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, '2026', LEFT_M, Inches(3.7), Inches(1.5), Inches(0.35), size=12, color=TEAL)
    # Bottom
    box(s, 0, Inches(4.95), SLIDE_W, Inches(0.675), RGBColor(0x00, 0x2D, 0x52), radius=False)
    txt(s, 'АО «Атомик Софт»  •  automiq.ru  •  Томск', LEFT_M, Inches(5.05), Inches(6), Inches(0.4),
        size=10, color=RGBColor(0x80, 0xB0, 0xD0))

    # ================================================================
    # SLIDE 2 — Почему миграция (infographic cards)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Почему миграция актуальна')

    cards = [
        ('🚫', 'Санкции', 'AVEVA прекратила\nподдержку в РФ', RED, LIGHT_RED_BG),
        ('💰', 'Подписка', '~$108K/год\nза 25-50K тегов', AMBER, LIGHT_AMBER_BG),
        ('⚖️', 'ФЗ-187 КИИ', 'Требование перехода\nна отечественное ПО', DEEP_BLUE, LIGHT_BLUE_BG),
        ('🔒', 'Деактивация', 'При неоплате ПО\nперестаёт работать', RED, LIGHT_RED_BG),
        ('🔄', 'Обновления', 'Нет патчей,\nнет новых версий', AMBER, LIGHT_AMBER_BG),
        ('☁️', 'Облако', 'Риск отключения\nAVEVA Cloud', PURPLE, LIGHT_PURPLE_BG),
    ]

    card_w = Inches(2.8)
    card_h = Inches(1.7)
    gap = Inches(0.2)
    start_x = LEFT_M + Inches(0.05)
    row1_y = Inches(1.1)
    row2_y = Inches(3.0)

    for i, (icon, title, desc, title_clr, bg_clr) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap)
        y = row1_y if row == 0 else row2_y

        box(s, x, y, card_w, card_h, bg_clr)
        # Left accent bar
        box(s, x, y, Inches(0.06), card_h, title_clr, radius=False)
        # Icon
        txt(s, icon, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5), size=22, color=title_clr)
        # Title
        txt(s, title, x + Inches(0.7), y + Inches(0.15), Inches(1.9), Inches(0.35), size=13, bold=True, color=title_clr)
        # Desc
        txt(s, desc, x + Inches(0.7), y + Inches(0.55), Inches(1.9), Inches(1.0), size=10, color=DARK_SLATE)

    # ================================================================
    # SLIDE 3 — PI System архитектура (visual diagram)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'PI System — архитектура')

    # Layered architecture diagram
    layers = [
        ('Облако', 'AVEVA Data Hub  •  AVEVA Insight  •  AVEVA Connect', RGBColor(0xA0, 0x60, 0xC0), Inches(1.0)),
        ('Аналитика', 'PI Analytics  •  Asset Analytics  •  PI Integrator for BA', MEDIUM_BLUE, Inches(1.75)),
        ('Визуализация', 'PI Vision (веб)  •  PI ProcessBook  •  PI DataLink (Excel)', TEAL, Inches(2.5)),
        ('Сервер', 'Data Archive  •  Asset Framework  •  Notifications  •  Event Frames', DEEP_BLUE, Inches(3.25)),
        ('Сбор данных', '450+ PI Interfaces  •  PI Connectors  •  Buffer Subsystem', RGBColor(0x3D, 0x7A, 0x44), Inches(4.0)),
    ]

    lx = Inches(0.8)
    lw = Inches(8.4)
    lh = Inches(0.6)

    for label, content, clr, ly in layers:
        # Main bar
        box(s, lx, ly, lw, lh, clr, radius=False)
        txt(s, label, lx + Inches(0.15), ly + Inches(0.05), Inches(1.8), Inches(0.3), size=12, bold=True, color=WHITE)
        txt(s, content, lx + Inches(2.0), ly + Inches(0.05), Inches(6.0), Inches(0.5), size=10, color=WHITE)

    # Side labels
    txt(s, 'PI Web API\n(REST)', Inches(0.05), Inches(1.5), Inches(0.7), Inches(1.0), size=8, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s, 'PI SDK\n(.NET)', Inches(0.05), Inches(2.8), Inches(0.7), Inches(0.8), size=8, color=TEAL, align=PP_ALIGN.CENTER)
    # Vertical connector
    connector_line(s, Inches(0.35), Inches(1.5), Inches(0.02), Inches(2.7), TEAL)

    # Right side
    txt(s, 'PI to PI\nРепликация', Inches(9.3), Inches(2.0), Inches(0.7), Inches(0.8), size=8, color=MEDIUM_BLUE, align=PP_ALIGN.CENTER)
    connector_line(s, Inches(9.55), Inches(1.8), Inches(0.02), Inches(2.0), MEDIUM_BLUE)

    # ================================================================
    # SLIDE 4 — Альфа платформа архитектура (visual diagram)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Альфа платформа — архитектура')

    layers2 = [
        ('Отчёты', 'Alpha.Report.Generator 2.0  •  Diagnostics.AuditReport', MEDIUM_BLUE, Inches(1.0)),
        ('Визуализация', 'Alpha.HMI 2.0  •  WebViewer 2.0  •  alpha.hmi.charts  •  Alpha.HMI.Alarms 3.3', TEAL, Inches(1.75)),
        ('Ядро', 'Alpha.Server 6.4  •  Резервирование  •  Кластер  •  Alpha.Om / JS', DEEP_BLUE, Inches(2.5)),
        ('Архив', 'Alpha.Historian 4.0  •  Журнал транзакций  •  Сжатие LZMA', RGBColor(0x06, 0x71, 0x96), Inches(3.25)),
        ('Сбор данных', '~30 модулей: OPC UA/DA, Modbus, IEC 104/101/61850, S7, MQTT, BACnet, SNMP', RGBColor(0x3D, 0x7A, 0x44), Inches(4.0)),
    ]

    for label, content, clr, ly in layers2:
        box(s, lx, ly, lw, lh, clr, radius=False)
        txt(s, label, lx + Inches(0.15), ly + Inches(0.05), Inches(1.8), Inches(0.3), size=12, bold=True, color=WHITE)
        txt(s, content, lx + Inches(2.0), ly + Inches(0.05), Inches(6.0), Inches(0.5), size=10, color=WHITE)

    # Side labels
    txt(s, 'Alpha.Link\n(TCP)', Inches(0.05), Inches(1.8), Inches(0.7), Inches(0.8), size=8, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s, 'OPC UA\nSQL Conn.', Inches(0.05), Inches(2.8), Inches(0.7), Inches(0.8), size=8, color=TEAL, align=PP_ALIGN.CENTER)
    connector_line(s, Inches(0.35), Inches(1.8), Inches(0.02), Inches(2.0), TEAL)

    txt(s, 'DevStudio\nИнструменты\nразработки', Inches(9.3), Inches(1.5), Inches(0.7), Inches(1.2), size=8, color=MEDIUM_BLUE, align=PP_ALIGN.CENTER)
    connector_line(s, Inches(9.55), Inches(1.5), Inches(0.02), Inches(2.5), MEDIUM_BLUE)

    txt(s, 'ФСТЭК  •  IEC 62443', Inches(9.25), Inches(4.2), Inches(0.8), Inches(0.4), size=7, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 5 — Mind map: что заменяется (brainstorm diagram)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Карта замены компонентов PI System')

    # Central node
    cx, cy = Inches(4.5), Inches(2.8)
    cw, ch = Inches(1.6), Inches(0.7)
    box(s, cx, cy, cw, ch, DEEP_BLUE)
    txt(s, 'PI System\n→ Альфа', cx + Inches(0.05), cy + Inches(0.05), cw - Inches(0.1), ch - Inches(0.1),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Surrounding nodes — (x, y, w, h, label, sublabel, bg_color, accent_color, category)
    nodes = [
        # Top row — полная замена (green)
        (Inches(0.3), Inches(1.0), Inches(2.0), Inches(0.6), 'Data Archive', '→ Alpha.Historian 4.0', LIGHT_GREEN_BG, GREEN),
        (Inches(2.6), Inches(1.0), Inches(2.0), Inches(0.6), 'Notifications', '→ Alpha.HMI.Alarms 3.3', LIGHT_GREEN_BG, GREEN),
        (Inches(5.1), Inches(1.0), Inches(2.2), Inches(0.6), 'ProcessBook', '→ Alpha.HMI 2.0 + WebViewer', LIGHT_GREEN_BG, GREEN),
        (Inches(7.6), Inches(1.0), Inches(2.1), Inches(0.6), 'Buffer Subsystem', '→ Встроен в Alpha.Server', LIGHT_GREEN_BG, GREEN),

        # Left — частичная замена (amber)
        (Inches(0.3), Inches(2.6), Inches(2.0), Inches(0.7), '450+ Interfaces', '→ ~30 модулей\n+ OPC UA шлюзы', LIGHT_AMBER_BG, AMBER),
        (Inches(0.3), Inches(3.6), Inches(2.0), Inches(0.7), 'PI Vision', '→ WebViewer\n(без self-service)', LIGHT_AMBER_BG, AMBER),

        # Right — пробелы (red)
        (Inches(7.6), Inches(2.6), Inches(2.1), Inches(0.7), 'Event Frames', '❌ Нет аналога\nНужна доработка', LIGHT_RED_BG, RED),
        (Inches(7.6), Inches(3.6), Inches(2.1), Inches(0.7), 'REST API', '❌ Нет\nAlpha.Link + OPC UA', LIGHT_RED_BG, RED),

        # Bottom — пробелы (red)
        (Inches(1.5), Inches(4.5), Inches(2.0), Inches(0.6), 'PI to PI', '❌ Нет репликации', LIGHT_RED_BG, RED),
        (Inches(4.0), Inches(4.5), Inches(2.0), Inches(0.6), 'Edge Data Store', '❌ Нет edge', LIGHT_RED_BG, RED),
        (Inches(6.5), Inches(4.5), Inches(2.0), Inches(0.6), 'Integrator for BA', '❌ Нет прямой BI', LIGHT_RED_BG, RED),
    ]

    for nx, ny, nw, nh, label, sublabel, bg, accent in nodes:
        box(s, nx, ny, nw, nh, bg)
        box(s, nx, ny, Inches(0.05), nh, accent, radius=False)
        txt(s, label, nx + Inches(0.12), ny + Inches(0.03), nw - Inches(0.15), Inches(0.25),
            size=9, bold=True, color=accent)
        txt(s, sublabel, nx + Inches(0.12), ny + Inches(0.25), nw - Inches(0.15), nh - Inches(0.3),
            size=8, color=DARK_SLATE)

    # Connector lines from center to nodes
    center_x = cx + cw / 2
    center_y = cy + ch / 2

    # Legend
    for i, (clr, lbl) in enumerate([(GREEN, '✅ Полная замена'), (AMBER, '⚠️ Частичная'), (RED, '❌ Пробел')]):
        lx = Inches(0.3) + i * Inches(1.6)
        circle(s, lx, Inches(0.92), Inches(0.12), clr)
        txt(s, lbl, lx + Inches(0.15), Inches(0.9), Inches(1.3), Inches(0.2), size=8, bold=True, color=clr)

    # ================================================================
    # SLIDE 6 — Покрытие (visual tiles)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Полное покрытие — что заменяется 1:1')

    tiles = [
        ('PI Data Archive', 'Alpha.Historian 4.0', 'Архив с журналом транзакций\nи сжатием LZMA', '📊'),
        ('PI Notifications', 'Alpha.HMI.Alarms 3.3', 'Тревоги + рассылка\nSMTP/Syslog/OPC', '🔔'),
        ('PI ProcessBook', 'Alpha.HMI 2.0', 'Десктоп + веб\nвизуализация', '🖥️'),
        ('PI System Mgmt', 'DevStudio + Alpha.Diag', 'Конфигурация,\nдиагностика', '⚙️'),
        ('PI Buffer', 'Alpha.Server', 'Встроенная буферизация\nбез настройки', '💾'),
        ('PI SDK', 'Alpha.Link + OPC UA', 'Несколько вариантов\nпрограммного доступа', '🔌'),
    ]

    tw = Inches(2.8)
    th = Inches(1.7)
    gap = Inches(0.18)
    sx = LEFT_M + Inches(0.05)
    r1y = Inches(1.05)
    r2y = Inches(2.95)

    for i, (old, new, desc, icon) in enumerate(tiles):
        col = i % 3
        row = i // 3
        x = sx + col * (tw + gap)
        y = r1y if row == 0 else r2y

        box(s, x, y, tw, th, LIGHT_GREEN_BG)
        # Top accent bar
        box(s, x, y, tw, Inches(0.04), GREEN, radius=False)
        # Icon
        txt(s, icon, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4), size=20)
        # Old → New
        txt(s, old, x + Inches(0.6), y + Inches(0.12), Inches(1.0), Inches(0.3), size=9, color=DARK_SLATE)
        arrow_right(s, x + Inches(1.55), y + Inches(0.18), Inches(0.25), Inches(0.15), GREEN)
        txt(s, new, x + Inches(1.85), y + Inches(0.12), Inches(0.9), Inches(0.3), size=9, bold=True, color=DEEP_BLUE)
        # Desc
        txt(s, desc, x + Inches(0.15), y + Inches(0.55), tw - Inches(0.3), Inches(1.0), size=9, color=DARK_SLATE)

    # ================================================================
    # SLIDE 7 — Частичная замена + пробелы (split view)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Частичная замена и функциональные пробелы')

    # Left column — partial
    txt(s, '⚠️  Частичная замена', LEFT_M, Inches(0.95), Inches(4.2), Inches(0.3), size=13, bold=True, color=AMBER)
    partial = [
        ('PI Interfaces (450+)', '~30 модулей Alpha.Server\n+ OPC UA шлюзы для остальных'),
        ('PI Vision (self-service)', 'WebViewer — есть, но без\nконструктора дашбордов'),
        ('PI AF (Asset Framework)', 'DevStudio — design-time\n(не runtime как AF)'),
        ('PI DataLink (Excel)', 'Report.Generator 2.0\n(не интерактивный Excel add-in)'),
    ]

    py = Inches(1.35)
    for title, desc in partial:
        box(s, LEFT_M, py, Inches(4.2), Inches(0.85), LIGHT_AMBER_BG)
        box(s, LEFT_M, py, Inches(0.05), Inches(0.85), AMBER, radius=False)
        txt(s, title, LEFT_M + Inches(0.15), py + Inches(0.05), Inches(3.9), Inches(0.25), size=10, bold=True, color=DEEP_BLUE)
        txt(s, desc, LEFT_M + Inches(0.15), py + Inches(0.3), Inches(3.9), Inches(0.5), size=9, color=DARK_SLATE)
        py += Inches(0.93)

    # Right column — gaps
    rx = Inches(5.2)
    txt(s, '❌  Функциональные пробелы', rx, Inches(0.95), Inches(4.3), Inches(0.3), size=13, bold=True, color=RED)
    gaps = [
        ('Event Frames', 'Контекстные события\nс привязкой к активам'),
        ('PI Web API (REST)', 'Нет открытого REST API'),
        ('PI to PI (репликация)', 'Нет встроенной репликации\nмежду архивами'),
        ('Edge Data Store', 'Нет лёгкого edge-решения'),
        ('Manual Logger', 'Нет ручного ввода данных'),
        ('Integrator for BA', 'Нет прямой BI-интеграции'),
    ]

    py = Inches(1.35)
    for title, desc in gaps:
        box(s, rx, py, Inches(4.3), Inches(0.58), LIGHT_RED_BG)
        box(s, rx, py, Inches(0.05), Inches(0.58), RED, radius=False)
        txt(s, title, rx + Inches(0.15), py + Inches(0.03), Inches(1.8), Inches(0.2), size=9, bold=True, color=RED)
        txt(s, desc, rx + Inches(2.0), py + Inches(0.03), Inches(2.1), Inches(0.5), size=8, color=DARK_SLATE)
        py += Inches(0.62)

    # ================================================================
    # SLIDE 8 — Стратегия закрытия пробелов (flow diagram)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Стратегия закрытия пробелов')

    rows = [
        ('Event Frames', 'Alpha.Om + Alarms + PostgreSQL', 'Кастомная доработка', TEAL),
        ('REST API', 'Alpha.RMap → SQL → REST', 'REST-обёртка', TEAL),
        ('Репликация', 'Historian → PostgreSQL → Import', 'Периодический экспорт', MEDIUM_BLUE),
        ('Edge Data Store', 'Alpha.Server (мин. конфиг)', 'Лёгкая инсталляция', MEDIUM_BLUE),
        ('Ручной ввод', 'WebViewer + Alpha.Om скрипт', 'Веб-форма', GREEN),
        ('BI-интеграция', 'SQL Connector → Grafana/PowerBI', 'Через PostgreSQL', GREEN),
        ('450+ интерфейсов', '~30 модулей + OPC UA шлюзы', '~80% покрытие', AMBER),
    ]

    sy = Inches(1.0)
    rh = Inches(0.5)
    gap = Inches(0.08)

    # Header
    box(s, LEFT_M, sy, Inches(2.5), Inches(0.35), DEEP_BLUE, radius=False)
    box(s, Inches(3.15), sy, Inches(3.3), Inches(0.35), DEEP_BLUE, radius=False)
    box(s, Inches(6.55), sy, Inches(2.9), Inches(0.35), DEEP_BLUE, radius=False)
    txt(s, 'Пробел', LEFT_M + Inches(0.1), sy + Inches(0.03), Inches(2.3), Inches(0.3), size=10, bold=True, color=WHITE)
    txt(s, 'Решение', Inches(3.25), sy + Inches(0.03), Inches(3.1), Inches(0.3), size=10, bold=True, color=WHITE)
    txt(s, 'Подход', Inches(6.65), sy + Inches(0.03), Inches(2.7), Inches(0.3), size=10, bold=True, color=WHITE)

    sy += Inches(0.42)
    for i, (gap_name, solution, approach, clr) in enumerate(rows):
        bg = LIGHTER_GRAY if i % 2 == 0 else WHITE
        box(s, LEFT_M, sy, Inches(8.9), rh, bg, radius=False)

        # Accent dot
        circle(s, LEFT_M + Inches(0.08), sy + Inches(0.17), Inches(0.15), clr)
        txt(s, gap_name, LEFT_M + Inches(0.3), sy + Inches(0.08), Inches(2.1), Inches(0.35), size=10, bold=True, color=DEEP_BLUE)

        arrow_right(s, Inches(2.8), sy + Inches(0.15), Inches(0.25), Inches(0.18), clr)
        txt(s, solution, Inches(3.15), sy + Inches(0.08), Inches(3.2), Inches(0.35), size=9, color=DARK_SLATE)

        # Approach tag
        tag_w = Inches(1.8)
        tag = box(s, Inches(6.65), sy + Inches(0.08), tag_w, Inches(0.3), clr)
        txt(s, approach, Inches(6.7), sy + Inches(0.1), tag_w - Inches(0.1), Inches(0.25), size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        sy += rh + gap

    # ================================================================
    # SLIDE 9 — Архитектурное различие (visual comparison)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Ключевое архитектурное различие')

    # VS badge in center
    vs_circle = circle(s, Inches(4.55), Inches(2.55), Inches(0.6), TEAL)
    txt(s, 'VS', Inches(4.55), Inches(2.6), Inches(0.6), Inches(0.45), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Left — PI System
    box(s, Inches(0.3), Inches(1.0), Inches(4.1), Inches(3.9), LIGHT_GRAY)
    box(s, Inches(0.3), Inches(1.0), Inches(4.1), Inches(0.5), RGBColor(0x6A, 0x6A, 0x6A), radius=False)
    txt(s, 'AVEVA PI System', Inches(0.45), Inches(1.05), Inches(3.8), Inches(0.4), size=14, bold=True, color=WHITE)

    pi_items = [
        ('📦', 'Data Infrastructure'),
        ('🎯', 'Фокус: сбор → хранение → доставка'),
        ('📊', 'Визуализация вторична'),
        ('🏗️', 'Asset Framework (runtime)'),
        ('🔌', 'REST API, SDK — открытая экосистема'),
        ('☁️', 'Облачные сервисы (AVEVA Data Hub)'),
    ]
    py = Inches(1.65)
    for icon, item in pi_items:
        txt(s, f'{icon}  {item}', Inches(0.5), py, Inches(3.7), Inches(0.3), size=10, color=DARK_SLATE)
        py += Inches(0.38)

    # Right — Альфа
    box(s, Inches(5.3), Inches(1.0), Inches(4.4), Inches(3.9), LIGHT_BLUE_BG)
    box(s, Inches(5.3), Inches(1.0), Inches(4.4), Inches(0.5), DEEP_BLUE, radius=False)
    txt(s, 'Альфа платформа', Inches(5.45), Inches(1.05), Inches(4.1), Inches(0.4), size=14, bold=True, color=WHITE)

    alpha_items = [
        ('🏭', 'SCADA-платформа'),
        ('🎯', 'Фокус: визуализация + управление'),
        ('🖥️', 'HMI/SCADA — ядро'),
        ('🛠️', 'DevStudio (design-time)'),
        ('🔗', 'Alpha.Link + OPC UA + SQL'),
        ('🔒', 'Только on-premises (безопасно)'),
    ]
    py = Inches(1.65)
    for icon, item in alpha_items:
        txt(s, f'{icon}  {item}', Inches(5.5), py, Inches(4.0), Inches(0.3), size=10, color=DEEP_BLUE)
        py += Inches(0.38)

    # ================================================================
    # SLIDE 10 — Преимущества (hexagon infographic)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Преимущества перехода на Альфа платформу')

    benefits = [
        ('Реестр\nМинцифры', 'ФЗ-187 о КИИ', DEEP_BLUE),
        ('Бессрочная\nлицензия', 'Нет подписок', GREEN),
        ('Linux\nподдержка', 'Astra, РЕД, Альт', TEAL),
        ('Единая\nплатформа', 'SCADA+HMI+Hist', MEDIUM_BLUE),
        ('24/7\nподдержка', 'На русском языке', RGBColor(0x06, 0x71, 0x96)),
        ('6500+\nинсталляций', '80+ партнёров', GREEN),
        ('ФСТЭК\nлицензия', 'Разработка СЗКИ', DEEP_BLUE),
        ('п.29.3\nПриказ №239', 'Соответствие', TEAL),
    ]

    hw = Inches(1.05)
    hh = Inches(1.05)
    cols = 4
    start_x = Inches(0.6)
    gap_x = Inches(0.18)
    r1y = Inches(1.05)
    r2y = Inches(3.0)

    for i, (title, sub, clr) in enumerate(benefits):
        col = i % cols
        row = i // cols
        x = start_x + col * (Inches(2.2) + gap_x)
        y = r1y if row == 0 else r2y

        # Hexagon
        hexagon(s, x, y, hw, hh, clr)
        txt(s, title, x + Inches(0.08), y + Inches(0.2), hw - Inches(0.16), Inches(0.6),
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Label to the right
        txt(s, sub, x + hw + Inches(0.1), y + Inches(0.3), Inches(1.0), Inches(0.4),
            size=9, color=DARK_SLATE)

    # ================================================================
    # SLIDE 11 — Протоколы (visual matrix)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Протоколы и коннекторы')

    protocols = [
        ('OPC UA', 'C+S', '✅', GREEN),
        ('OPC DA/HDA/AE', 'C+S', '✅ Win', MEDIUM_BLUE),
        ('Modbus TCP/RTU', 'M+S', '✅', GREEN),
        ('IEC 60870-5-104', 'M+S', '✅', GREEN),
        ('IEC 60870-5-101', 'M+S', '✅', GREEN),
        ('IEC 61850', 'Client', '✅', GREEN),
        ('Siemens S7', 'Client', '✅', GREEN),
        ('EtherNet/IP', 'Client', '✅', GREEN),
        ('FINS (Omron)', 'Client', '✅', GREEN),
        ('BACnet IP', 'Client', '✅', GREEN),
        ('MQTT', 'Pub/Sub', '✅', GREEN),
        ('SNMP', 'Manager', '✅', GREEN),
        ('SQL Connector', 'PG/MS/Ora', '✅', GREEN),
    ]

    # Headers
    hy = Inches(0.95)
    box(s, LEFT_M, hy, Inches(3.0), Inches(0.3), DEEP_BLUE, radius=False)
    box(s, Inches(3.6), hy, Inches(1.5), Inches(0.3), DEEP_BLUE, radius=False)
    box(s, Inches(5.15), hy, Inches(1.0), Inches(0.3), DEEP_BLUE, radius=False)
    txt(s, 'Протокол', LEFT_M + Inches(0.1), hy + Inches(0.02), Inches(2.8), Inches(0.25), size=10, bold=True, color=WHITE)
    txt(s, 'Режим', Inches(3.7), hy + Inches(0.02), Inches(1.3), Inches(0.25), size=10, bold=True, color=WHITE)
    txt(s, 'Статус', Inches(5.25), hy + Inches(0.02), Inches(0.8), Inches(0.25), size=10, bold=True, color=WHITE)

    ry = hy + Inches(0.35)
    for i, (proto, mode, status, clr) in enumerate(protocols):
        bg = LIGHTER_GRAY if i % 2 == 0 else WHITE
        box(s, LEFT_M, ry, Inches(5.6), Inches(0.3), bg, radius=False)
        txt(s, proto, LEFT_M + Inches(0.1), ry + Inches(0.02), Inches(2.8), Inches(0.25), size=9, color=DEEP_BLUE)
        txt(s, mode, Inches(3.7), ry + Inches(0.02), Inches(1.3), Inches(0.25), size=9, color=DARK_SLATE)
        txt(s, status, Inches(5.25), ry + Inches(0.02), Inches(0.8), Inches(0.25), size=9, bold=True, color=clr)
        ry += Inches(0.3)

    # Right side — summary infographic
    box(s, Inches(6.5), Inches(1.1), Inches(3.2), Inches(4.0), LIGHT_BLUE_BG)
    txt(s, '~30', Inches(6.7), Inches(1.3), Inches(2.8), Inches(0.6), size=36, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s, 'коммуникационных\nмодулей', Inches(6.7), Inches(1.9), Inches(2.8), Inches(0.5), size=12, color=DEEP_BLUE, align=PP_ALIGN.CENTER)

    txt(s, '~80%', Inches(6.7), Inches(2.7), Inches(2.8), Inches(0.6), size=36, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, 'промышленных\nпротоколов', Inches(6.7), Inches(3.3), Inches(2.8), Inches(0.5), size=12, color=DEEP_BLUE, align=PP_ALIGN.CENTER)

    txt(s, 'Остальное через\nOPC UA шлюзы', Inches(6.7), Inches(4.1), Inches(2.8), Inches(0.5), size=10, color=DARK_SLATE, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 12 — План миграции: Timeline (horizontal)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'План миграции: 5 фаз')

    phases = [
        ('1', 'АУДИТ', '1-2 мес', 'Инвентаризация PI:\nтеги, AF-модели,\nинтерфейсы, отчёты', TEAL),
        ('2', 'АРХИТЕКТУРА', '1 мес', 'Целевая архитектура,\nрасчёт лицензий,\nперечень доработок', MEDIUM_BLUE),
        ('3', 'ПИЛОТ', '2-3 мес', 'Пилотный объект,\nмиграция тегов,\nтестирование', DEEP_BLUE),
        ('4', 'МИГРАЦИЯ', '3-6 мес', 'Последовательный\nперевод объектов,\nобучение персонала', RGBColor(0x06, 0x71, 0x96)),
        ('5', 'ЗАКРЫТИЕ', '1-2 мес', 'Отключение PI,\nперенос истории,\nдокументация', GREEN),
    ]

    # Timeline bar
    bar_y = Inches(2.35)
    box(s, Inches(0.5), bar_y, Inches(9.0), Inches(0.06), MEDIUM_BLUE, radius=False)

    bw = Inches(1.55)
    gap = Inches(0.25)
    sx = Inches(0.5)

    for i, (num, name, duration, desc, clr) in enumerate(phases):
        x = sx + i * (bw + gap)

        # Circle on timeline
        cd = Inches(0.4)
        circle(s, x + bw/2 - cd/2, bar_y - cd/2 + Inches(0.03), cd, clr)
        txt(s, num, x + bw/2 - cd/2, bar_y - cd/2 + Inches(0.07), cd, Inches(0.3),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Phase card above
        card_h = Inches(1.1)
        card_y = bar_y - card_h - Inches(0.3)
        box(s, x, card_y, bw, card_h, LIGHT_BLUE_BG)
        box(s, x, card_y, bw, Inches(0.04), clr, radius=False)
        txt(s, name, x + Inches(0.08), card_y + Inches(0.08), bw - Inches(0.16), Inches(0.25),
            size=10, bold=True, color=clr, align=PP_ALIGN.CENTER)
        txt(s, duration, x + Inches(0.08), card_y + Inches(0.3), bw - Inches(0.16), Inches(0.2),
            size=9, color=TEAL, align=PP_ALIGN.CENTER)

        # Connector line down from circle
        connector_line(s, x + bw/2, bar_y + Inches(0.25), Inches(0.015), Inches(0.25), clr)

        # Description below
        desc_y = bar_y + Inches(0.55)
        box(s, x, desc_y, bw, Inches(1.6), LIGHTER_GRAY)
        txt(s, desc, x + Inches(0.1), desc_y + Inches(0.1), bw - Inches(0.2), Inches(1.4),
            size=9, color=DARK_SLATE)

        # Arrows between phases
        if i < len(phases) - 1:
            ax = x + bw + Inches(0.03)
            arrow_right(s, ax, bar_y - Inches(0.08), Inches(0.18), Inches(0.2), MEDIUM_BLUE)

    # Total duration
    txt(s, 'Общая длительность: 8-14 месяцев', Inches(2.5), Inches(5.0), Inches(5.0), Inches(0.35),
        size=12, bold=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 13 — Риски (visual risk matrix)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Оценка рисков и митигация')

    risks = [
        ('Event Frames', 'Аудит использования EF → кастомная доработка', 'Средний', AMBER),
        ('Формат архива', 'Параллельная запись, экспорт/импорт', 'Средний', AMBER),
        ('Переобучение', '5 лабораторий в ВУЗах, онлайн-курсы', 'Низкий', GREEN),
        ('Простой', 'Параллельная работа PI + Альфа', 'Средний', AMBER),
        ('Недостающие протоколы', 'OPC UA шлюзы для экзотики', 'Низкий', GREEN),
        ('Облачные сервисы', 'On-prem аналитика, Grafana + SQL', 'Высокий', RED),
    ]

    ry = Inches(1.0)
    for i, (risk, mitigation, level, clr) in enumerate(risks):
        bg = LIGHTER_GRAY if i % 2 == 0 else WHITE
        box(s, LEFT_M, ry, Inches(8.9), Inches(0.6), bg, radius=False)

        # Risk level indicator
        level_w = Inches(0.8)
        box(s, LEFT_M + Inches(0.05), ry + Inches(0.12), level_w, Inches(0.35), clr)
        txt(s, level, LEFT_M + Inches(0.05), ry + Inches(0.14), level_w, Inches(0.3),
            size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Risk name
        txt(s, risk, LEFT_M + Inches(1.0), ry + Inches(0.05), Inches(2.5), Inches(0.5),
            size=11, bold=True, color=DEEP_BLUE)

        # Arrow
        arrow_right(s, Inches(3.8), ry + Inches(0.2), Inches(0.25), Inches(0.18), TEAL)

        # Mitigation
        txt(s, mitigation, Inches(4.2), ry + Inches(0.08), Inches(5.0), Inches(0.45),
            size=10, color=DARK_SLATE)

        ry += Inches(0.65)

    # Legend
    for i, (lbl, clr) in enumerate([('Низкий', GREEN), ('Средний', AMBER), ('Высокий', RED)]):
        lx = Inches(0.6) + i * Inches(1.5)
        circle(s, lx, Inches(5.0), Inches(0.15), clr)
        txt(s, lbl, lx + Inches(0.2), Inches(4.97), Inches(1.0), Inches(0.2), size=9, color=clr)

    # ================================================================
    # SLIDE 14 — TCO (infographic comparison)
    # ================================================================
    s = prs.slides.add_slide(layout_nf)
    add_header_bar(s, 'Экономика: TCO сравнение')

    # AVEVA side — bar chart visual
    box(s, Inches(0.5), Inches(1.1), Inches(4.0), Inches(3.8), LIGHT_RED_BG)
    txt(s, 'AVEVA PI System', Inches(0.7), Inches(1.15), Inches(3.6), Inches(0.35), size=14, bold=True, color=RED)

    # Bars representing costs
    bar_items = [
        ('~$108K / год', Inches(2.8), '25-50K тегов'),
        ('~$325K / 3 года', Inches(3.6), 'публичный контракт'),
    ]
    by = Inches(1.7)
    for label, bar_w, sub in bar_items:
        box(s, Inches(0.7), by, bar_w, Inches(0.5), RED, radius=False)
        txt(s, label, Inches(0.8), by + Inches(0.07), bar_w - Inches(0.2), Inches(0.35), size=12, bold=True, color=WHITE)
        txt(s, sub, Inches(0.7) + bar_w + Inches(0.1), by + Inches(0.1), Inches(1.0), Inches(0.3), size=9, color=DARK_SLATE)
        by += Inches(0.7)

    # Recurring payment warning
    txt(s, '⚠️ Ежегодная подписка\n⚠️ Валютные риски ($)\n⚠️ Risk of deactivation', Inches(0.8), Inches(3.3), Inches(3.5), Inches(0.8),
        size=10, color=RED)

    # VS
    circle(s, Inches(4.55), Inches(2.5), Inches(0.6), TEAL)
    txt(s, 'VS', Inches(4.55), Inches(2.55), Inches(0.6), Inches(0.45), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Alpha side
    box(s, Inches(5.5), Inches(1.1), Inches(4.2), Inches(3.8), LIGHT_GREEN_BG)
    txt(s, 'Альфа платформа', Inches(5.7), Inches(1.15), Inches(3.8), Inches(0.35), size=14, bold=True, color=GREEN)

    # Big price
    txt(s, '₽', Inches(5.7), Inches(1.7), Inches(3.8), Inches(1.0), size=60, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, 'Бессрочная лицензия', Inches(5.7), Inches(2.6), Inches(3.8), Inches(0.3), size=14, bold=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER)

    txt(s, '✅ Разовая закупка\n✅ Расчёты в рублях\n✅ Работает без интернета\n✅ Окупаемость: 1.5-3 года',
        Inches(5.8), Inches(3.1), Inches(3.5), Inches(1.2), size=10, color=GREEN)

    # ================================================================
    # SLIDE 15 — Следующие шаги (numbered flow)
    # ================================================================
    s = prs.slides.add_slide(layout)
    add_header_bar(s, 'Следующие шаги')

    steps = [
        ('Согласование объёма аудита\nPI-инфраструктуры НОВАТЭК', TEAL),
        ('Определение пилотного\nобъекта для миграции', TEAL),
        ('Перечень критичных функций PI\n(Event Frames? DataLink? API?)', MEDIUM_BLUE),
        ('Расчёт лицензий Альфа\nпод целевую архитектуру', MEDIUM_BLUE),
        ('Рабочая группа:\nНОВАТЭК + АС + интегратор', DEEP_BLUE),
        ('Детальный план-график\nи бюджет', DEEP_BLUE),
    ]

    # Two columns of 3
    for i, (step, clr) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(4.8)
        y = Inches(1.1) + row * Inches(1.35)

        # Number circle
        cd = Inches(0.5)
        circle(s, x, y + Inches(0.15), cd, clr)
        txt(s, str(i + 1), x + Inches(0.02), y + Inches(0.2), cd, Inches(0.4),
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Step card
        cw = Inches(3.8)
        box(s, x + Inches(0.65), y, cw, Inches(1.05), LIGHT_BLUE_BG)
        box(s, x + Inches(0.65), y, Inches(0.04), Inches(1.05), clr, radius=False)
        txt(s, step, x + Inches(0.8), y + Inches(0.15), cw - Inches(0.25), Inches(0.8), size=11, color=DARK_SLATE)

        # Arrow down between rows (except last row)
        if row < 2 and col == 0:
            arrow_down(s, x + Inches(0.15), y + Inches(1.05), Inches(0.2), Inches(0.2), RGBColor(0xCC, 0xCC, 0xCC))
        if row < 2 and col == 1:
            arrow_down(s, x + Inches(0.15), y + Inches(1.05), Inches(0.2), Inches(0.2), RGBColor(0xCC, 0xCC, 0xCC))

    # ================================================================
    # SLIDE 16 — Контакты
    # ================================================================
    s = prs.slides.add_slide(layout_contact)
    box(s, 0, 0, SLIDE_W, SLIDE_H, DEEP_BLUE, radius=False)
    box(s, 0, 0, SLIDE_W, Inches(0.05), TEAL, radius=False)

    # Decorative circles
    for cx, cy, d in [(Inches(7.5), Inches(1.0), Inches(2.5)), (Inches(8.5), Inches(3.0), Inches(1.5))]:
        c = circle(s, cx, cy, d, TEAL)
        solidFill = c._element.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_elem = srgb.makeelement(qn('a:alpha'), {'val': '15000'})
                srgb.append(alpha_elem)

    txt(s, 'Контакты', LEFT_M, Inches(0.8), Inches(5), Inches(0.6), size=28, bold=True, color=WHITE)
    txt(s, 'Станислав Павловский', LEFT_M, Inches(1.8), Inches(5), Inches(0.4), size=18, bold=True, color=WHITE)
    txt(s, 'Telegram: @Integraleus', LEFT_M, Inches(2.3), Inches(5), Inches(0.35), size=14, color=LIGHT_CYAN)

    # Divider line
    box(s, LEFT_M, Inches(2.9), Inches(2.5), Inches(0.02), TEAL, radius=False)

    txt(s, 'АО «Атомик Софт»\n634050 Томск, пр. Ленина 60/1\n+7 (3822) 281 914\ninfo@automiq.ru  •  automiq.ru',
        LEFT_M, Inches(3.1), Inches(4), Inches(1.5), size=11, color=RGBColor(0x80, 0xB0, 0xD0))

    txt(s, 'Будущее\nавтоматизации\nв ваших руках.', Inches(6.0), Inches(1.5), Inches(3.5), Inches(1.5),
        size=20, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)

    prs.save(OUTPUT)
    print(f'✅ Saved: {OUTPUT}')
    print(f'   Slides: {len(prs.slides)}')


if __name__ == '__main__':
    build()
