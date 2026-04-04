#!/usr/bin/env python3
"""
NOVATEK PI→Alpha — Executive version for VP of IT.
Focus: strategy, risks, money, timeline. Minimal tech details.
Атомик Софт brand style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# === Brand ===
BRAND_BLUE = RGBColor(0x00, 0x54, 0x97)
BLUE_TEXT = RGBColor(0x00, 0x55, 0x9A)
TEAL = RGBColor(0x09, 0x97, 0xC8)
MEDIUM_BLUE = RGBColor(0x33, 0x75, 0xAD)
GREEN = RGBColor(0x82, 0xC4, 0x44)
BODY_TEXT = RGBColor(0x44, 0x54, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
LINE_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
AMBER = RGBColor(0xE8, 0x9C, 0x00)
RED_ACCENT = RGBColor(0xCC, 0x33, 0x33)
DARK_BLUE_BG = RGBColor(0x00, 0x3D, 0x6E)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.625)
ASSETS = '/tmp/atomicsoft_assets'
TEMPLATE = '/root/.openclaw/workspace/agents/main/skills/presentation-designer/templates/atomicsoft-about.pptx'
OUTPUT = '/root/.openclaw/workspace/agents/main/projects/presentations/NOVATEK_VP_IT.pptx'


def remove_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    while len(sldIdLst):
        rId = sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldIdLst[0])


def add_brand_header(slide, title_text):
    y = Inches(0.22)
    h = Inches(0.50)
    left_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y + Inches(0.25), Inches(0.5), Inches(0.015))
    left_line.fill.solid()
    left_line.fill.fore_color.rgb = BRAND_BLUE
    left_line.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.55), y, Inches(8.5), h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.font.name = 'Arial'
    right_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), y + Inches(0.25), Inches(1.5), Inches(0.015))
    right_line.fill.solid()
    right_line.fill.fore_color.rgb = BRAND_BLUE
    right_line.line.fill.background()


def txt(slide, text, x, y, w, h, size=12, bold=False, color=None, align=PP_ALIGN.LEFT, name='Arial'):
    color = color or BODY_TEXT
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return txBox


def box(slide, x, y, w, h, fill_color, radius=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s


def thin_line(slide, x, y, w, color=None):
    color = color or LINE_GRAY
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.01))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def metric_big(slide, value, label, x, y, w=Inches(2.5), val_color=None, val_size=36):
    val_color = val_color or TEAL
    txt(slide, value, x, y, w, Inches(0.6), size=val_size, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(0.55), w, Inches(0.5), size=11, color=BODY_TEXT, align=PP_ALIGN.CENTER)


def add_logo(slide):
    logo_path = os.path.join(ASSETS, 'logo.png')
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.30), Inches(0.27), Inches(2.48), Inches(0.38))


def build():
    prs = Presentation(TEMPLATE)
    remove_all_slides(prs)
    
    layout = prs.slide_layouts[3]   # General
    layout_title = prs.slide_layouts[0]  # Helium
    layout_nf = prs.slide_layouts[2]     # No Footer
    layout_contact = prs.slide_layouts[4]

    # ═══════════════════════════════════════════
    # SLIDE 1 — Title
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout_title)
    bg_path = os.path.join(ASSETS, 'title_bg.jpg')
    if os.path.exists(bg_path):
        s.shapes.add_picture(bg_path, 0, 0, SLIDE_W, SLIDE_H)
    add_logo(s)
    
    txt(s, 'Переход с AVEVA PI System\nна Альфа платформу',
        Inches(0.30), Inches(1.6), Inches(6.0), Inches(1.4),
        size=28, bold=True, color=WHITE)
    txt(s, 'Стратегическое обоснование для руководства',
        Inches(0.30), Inches(3.1), Inches(5.0), Inches(0.4),
        size=14, color=RGBColor(0x80, 0xD0, 0xF0))
    
    box(s, 0, Inches(4.95), SLIDE_W, Inches(0.675), RGBColor(0x00, 0x2D, 0x52), radius=False)
    txt(s, 'АО «Атомик Софт»  •  automiq.ru  •  Конфиденциально',
        Inches(0.55), Inches(5.05), Inches(6), Inches(0.4),
        size=10, color=RGBColor(0x80, 0xB0, 0xD0))

    # ═══════════════════════════════════════════
    # SLIDE 2 — Проблема (executive summary)
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout)
    add_brand_header(s, 'Ситуация: почему нужно действовать сейчас')

    # Three risk cards
    risks = [
        ('🚫', 'Санкционный риск', 
         'AVEVA прекратила поддержку\nи продажу в РФ.\nНет обновлений. Нет патчей\nбезопасности. Нет техподдержки.',
         RED_ACCENT),
        ('💰', 'Финансовый риск',
         'AVEVA Flex: ~$108K/год\nза подписку.\nПри неоплате — ПО\nперестаёт работать.\nВалютный риск.',
         AMBER),
        ('⚖️', 'Регуляторный риск',
         'ФЗ-187 о КИИ требует\nперехода на отечественное ПО.\nСроки ужесточаются.\nФСТЭК-аудиты.',
         BRAND_BLUE),
    ]

    for i, (icon, title, desc, clr) in enumerate(risks):
        x = Inches(0.5) + i * Inches(3.15)
        w = Inches(2.95)
        y = Inches(0.95)
        
        box(s, x, y, w, Inches(3.8), LIGHT_GRAY)
        box(s, x, y, w, Inches(0.05), clr, radius=False)
        txt(s, icon, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4), size=22)
        txt(s, title, x + Inches(0.6), y + Inches(0.2), w - Inches(0.75), Inches(0.3),
            size=14, bold=True, color=clr)
        txt(s, desc, x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), Inches(2.8),
            size=11, color=BODY_TEXT)

    # Bottom urgency line
    txt(s, 'Каждый месяц без действий — это $9K подписки + растущий риск деактивации',
        Inches(0.5), Inches(4.9), Inches(9.0), Inches(0.3),
        size=11, bold=True, color=RED_ACCENT, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════
    # SLIDE 3 — Решение в одном слайде
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout)
    add_brand_header(s, 'Решение: Альфа платформа')

    txt(s, 'Российская SCADA-платформа №1\nот АО «Атомик Софт» (Томск)',
        Inches(0.6), Inches(0.9), Inches(9.0), Inches(0.6),
        size=14, bold=True, color=BRAND_BLUE)

    # Key metrics in a row
    metrics = [
        ('6500+', 'инсталляций\nв промышленности', TEAL),
        ('80+', 'партнёров-\nинтеграторов', TEAL),
        ('140+', 'сотрудников\nразработки', TEAL),
        ('10 лет', 'на рынке\nАСУ ТП', TEAL),
    ]
    
    for i, (val, label, clr) in enumerate(metrics):
        x = Inches(0.5) + i * Inches(2.3)
        metric_big(s, val, label, x, Inches(1.6), Inches(2.1), clr, 28)

    thin_line(s, Inches(0.6), Inches(2.8), Inches(8.8), TEAL)

    # Key selling points
    points = [
        'Реестр Минцифры — полное соответствие ФЗ-187 о КИИ',
        'Бессрочная лицензия в рублях — предсказуемый бюджет',
        'Windows + Linux (Astra, РЕД ОС) — сертифицированные ОС',
        'Лицензия ФСТЭК на разработку СЗКИ',
        'Техподдержка 24/7 на русском языке, вендор в России',
    ]
    
    py = Inches(2.95)
    for point in points:
        txt(s, '→', Inches(0.6), py, Inches(0.3), Inches(0.25), size=12, bold=True, color=TEAL)
        txt(s, point, Inches(0.9), py, Inches(8.5), Inches(0.25), size=11, color=BODY_TEXT)
        py += Inches(0.35)

    # ═══════════════════════════════════════════
    # SLIDE 4 — Суть различия (для руководства)
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout)
    add_brand_header(s, 'Что меняется: суть перехода')

    # Left
    box(s, Inches(0.4), Inches(0.9), Inches(4.4), Inches(2.5), LIGHT_GRAY)
    box(s, Inches(0.4), Inches(0.9), Inches(4.4), Inches(0.04), AMBER, radius=False)
    txt(s, 'AVEVA PI System', Inches(0.6), Inches(1.0), Inches(4.0), Inches(0.3),
        size=14, bold=True, color=BRAND_BLUE)
    txt(s, '«Склад данных»', Inches(0.6), Inches(1.3), Inches(4.0), Inches(0.25),
        size=12, bold=True, color=AMBER)
    txt(s, 'Собирает данные и раздаёт —\n'
        'аналитикам, Excel, облаку, BI.\n'
        'Управление процессом — не его задача.\n'
        'Подписка в долларах. Облако за рубежом.',
        Inches(0.6), Inches(1.7), Inches(4.0), Inches(1.5), size=11, color=BODY_TEXT)

    # Right
    box(s, Inches(5.0), Inches(0.9), Inches(4.6), Inches(2.5), LIGHT_BLUE_BG)
    box(s, Inches(5.0), Inches(0.9), Inches(4.6), Inches(0.04), TEAL, radius=False)
    txt(s, 'Альфа платформа', Inches(5.2), Inches(1.0), Inches(4.2), Inches(0.3),
        size=14, bold=True, color=BRAND_BLUE)
    txt(s, '«Диспетчерская с пультом»', Inches(5.2), Inches(1.3), Inches(4.2), Inches(0.25),
        size=12, bold=True, color=TEAL)
    txt(s, 'Оператор видит и управляет процессом.\n'
        'Хранение данных тоже есть (Historian).\n'
        'Но главное — пульт управления.\n'
        'Бессрочная лицензия в рублях. On-prem.',
        Inches(5.2), Inches(1.7), Inches(4.2), Inches(1.5), size=11, color=BODY_TEXT)

    # Bottom — what this means for business
    thin_line(s, Inches(0.5), Inches(3.6), Inches(9.0), TEAL)
    txt(s, 'Что это значит для бизнеса:', Inches(0.6), Inches(3.7), Inches(9.0), Inches(0.3),
        size=12, bold=True, color=BRAND_BLUE)
    txt(s, '• Переход с «хранилища данных» на «систему управления с хранилищем» — функционал расширяется\n'
        '• Некоторые «курьерские сервисы» (REST API, облако, Event Frames) потребуют доработки\n'
        '• Основные функции (архив, визуализация, тревоги, протоколы) заменяются полностью',
        Inches(0.6), Inches(4.0), Inches(9.0), Inches(1.3), size=10, color=BODY_TEXT)

    # ═══════════════════════════════════════════
    # SLIDE 5 — Что заменяется: светофор
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout)
    add_brand_header(s, 'Покрытие функционала')

    # Three columns: green / amber / red
    cols = [
        ('✅ Полная замена', GREEN, [
            'Архив данных',
            'Экраны оператора',
            'Тревоги и события',
            'Буферизация',
            'Администрирование',
            'Программный доступ',
        ]),
        ('⚠️ Частичная', AMBER, [
            'Коннекторы (450→30,\nно 80% покрыто)',
            'Self-service дашборды\n(есть, но проще)',
            'Модель активов\n(не «на лету»)',
            'Excel-отчёты\n(другой формат)',
        ]),
        ('❌ Требует доработки', RED_ACCENT, [
            'Запись событий\n(Event Frames)',
            'HTTP-доступ (REST API)',
            'Синхронизация серверов',
            'BI-интеграция\n(через SQL обход)',
        ]),
    ]

    for i, (title, clr, items) in enumerate(cols):
        x = Inches(0.4) + i * Inches(3.15)
        w = Inches(2.95)
        
        box(s, x, Inches(0.9), w, Inches(0.4), clr, radius=False)
        txt(s, title, x + Inches(0.1), Inches(0.93), w - Inches(0.2), Inches(0.3),
            size=12, bold=True, color=WHITE)
        
        py = Inches(1.4)
        for item in items:
            txt(s, '•  ' + item, x + Inches(0.1), py, w - Inches(0.2), Inches(0.5),
                size=10, color=BODY_TEXT)
            lines = item.count('\n') + 1
            py += Inches(0.25 + 0.15 * lines)

    # ═══════════════════════════════════════════
    # SLIDE 6 — TCO: деньги
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout)
    add_brand_header(s, 'Экономика: стоимость владения')

    # AVEVA cost
    box(s, Inches(0.4), Inches(0.9), Inches(4.4), Inches(2.0), LIGHT_GRAY)
    box(s, Inches(0.4), Inches(0.9), Inches(4.4), Inches(0.04), RED_ACCENT, radius=False)
    txt(s, 'AVEVA PI System (текущие затраты)', Inches(0.6), Inches(1.0), Inches(4.0), Inches(0.3),
        size=12, bold=True, color=RED_ACCENT)
    txt(s, '~$108K', Inches(0.6), Inches(1.4), Inches(1.8), Inches(0.5),
        size=32, bold=True, color=RED_ACCENT)
    txt(s, 'в год', Inches(2.3), Inches(1.55), Inches(1.0), Inches(0.3),
        size=14, color=BODY_TEXT)
    txt(s, '~$325K за 3 года (публичный контракт)\nПодписка в долларах. Валютный риск.\nПри неоплате — деактивация.',
        Inches(0.6), Inches(2.0), Inches(4.0), Inches(0.8), size=10, color=BODY_TEXT)

    # Alpha cost
    box(s, Inches(5.0), Inches(0.9), Inches(4.6), Inches(2.0), RGBColor(0xF0, 0xF8, 0xE8))
    box(s, Inches(5.0), Inches(0.9), Inches(4.6), Inches(0.04), GREEN, radius=False)
    txt(s, 'Альфа платформа (целевые затраты)', Inches(5.2), Inches(1.0), Inches(4.2), Inches(0.3),
        size=12, bold=True, color=GREEN)
    txt(s, '₽', Inches(5.2), Inches(1.4), Inches(1.0), Inches(0.5),
        size=32, bold=True, color=GREEN)
    txt(s, 'разово', Inches(6.0), Inches(1.55), Inches(1.5), Inches(0.3),
        size=14, color=BODY_TEXT)
    txt(s, 'Бессрочная лицензия. Рубли.\nНе зависит от курса. Не отключат.\nТехподдержка — отдельный договор.',
        Inches(5.2), Inches(2.0), Inches(4.2), Inches(0.8), size=10, color=BODY_TEXT)

    # ROI section
    thin_line(s, Inches(0.5), Inches(3.1), Inches(9.0), TEAL)
    txt(s, 'Окупаемость', Inches(0.6), Inches(3.25), Inches(9.0), Inches(0.3),
        size=14, bold=True, color=BRAND_BLUE)

    roi_points = [
        ('Год 1', 'Затраты на миграцию + лицензии Альфы. Экономия подписки AVEVA начинается.'),
        ('Год 2', 'Миграция завершена. Полная экономия $108K/год = ~₽10M/год.'),
        ('Год 3+', 'Чистая экономия. TCO перехода окупился. Только затраты на поддержку.'),
    ]

    py = Inches(3.6)
    for year, desc in roi_points:
        txt(s, year, Inches(0.6), py, Inches(0.8), Inches(0.25), size=11, bold=True, color=TEAL)
        txt(s, desc, Inches(1.5), py, Inches(8.0), Inches(0.25), size=10, color=BODY_TEXT)
        py += Inches(0.35)

    txt(s, 'Окупаемость: 1.5–3 года в зависимости от масштаба инфраструктуры',
        Inches(0.6), Inches(4.8), Inches(9.0), Inches(0.3),
        size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════
    # SLIDE 7 — Сроки и фазы (executive view)
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout)
    add_brand_header(s, 'План перехода: сроки и ресурсы')

    phases = [
        ('1', 'АУДИТ', '1-2 мес.', 'Считаем: сколько тегов,\nкакие протоколы, какие\nскрипты, что критично', TEAL),
        ('2', 'ПРОЕКТ', '1 мес.', 'Архитектура, лицензии,\nбюджет, рабочая группа', TEAL),
        ('3', 'ПИЛОТ', '2-3 мес.', 'Один объект — проверяем\nвсё на практике', MEDIUM_BLUE),
        ('4', 'МИГРАЦИЯ', '3-6 мес.', 'Переводим объекты.\nPI + Альфа параллельно.\nОбучаем людей', BRAND_BLUE),
        ('5', 'СДАЧА', '1-2 мес.', 'Отключаем PI.\nДокументация. SLA.', GREEN),
    ]

    # Timeline bar
    tl_y = Inches(2.3)
    thin_line(s, Inches(0.5), tl_y, Inches(9.0), TEAL)

    bw = Inches(1.6)
    for i, (num, name, dur, desc, clr) in enumerate(phases):
        x = Inches(0.5) + i * Inches(1.8)
        
        # Dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + bw/2 - Inches(0.15), tl_y - Inches(0.15), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = clr
        dot.line.fill.background()
        txt(s, num, x + bw/2 - Inches(0.13), tl_y - Inches(0.12), Inches(0.26), Inches(0.24),
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        
        # Above: name + duration
        txt(s, name, x, Inches(1.0), bw, Inches(0.25), size=11, bold=True, color=clr, align=PP_ALIGN.CENTER)
        txt(s, dur, x, Inches(1.25), bw, Inches(0.2), size=10, color=TEAL, align=PP_ALIGN.CENTER)
        
        # Below: description
        txt(s, desc, x, Inches(2.65), bw, Inches(1.5), size=9, color=BODY_TEXT)

    # Key message
    box(s, Inches(0.5), Inches(4.5), Inches(9.0), Inches(0.6), LIGHT_BLUE_BG)
    txt(s, 'Общая длительность: 8–14 месяцев.  Нулевой простой — PI и Альфа работают параллельно.',
        Inches(0.7), Inches(4.55), Inches(8.6), Inches(0.45),
        size=12, bold=True, color=BRAND_BLUE, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════
    # SLIDE 8 — Риски (executive view)
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout)
    add_brand_header(s, 'Управление рисками')

    risks_exec = [
        ('Высокий', RED_ACCENT,
         'Потеря Event Frames',
         'Уникальная функция PI для анализа событий. Проводим аудит: если используется — допрограммируем. Если нет — риск нулевой.'),
        ('Средний', AMBER,
         'Переобучение персонала',
         'Обучение от Атомик Софт: курсы, лаборатории, сертификация. Альфа проще PI в ежедневном использовании.'),
        ('Средний', AMBER,
         'Нет облачных сервисов',
         'PI Cloud отключён для РФ. Ставим Grafana + Power BI на своих серверах — данные остаются внутри периметра.'),
        ('Низкий', GREEN,
         'Технический простой',
         'PI и Альфа работают параллельно на каждом этапе. Переключение — после 2-4 недель валидации.'),
        ('Низкий', GREEN,
         'Неполное покрытие протоколов',
         '~80% промышленных протоколов — из коробки. Остальное — через стандартные OPC UA шлюзы.'),
    ]

    ry = Inches(0.9)
    for level, clr, risk, mitigation in risks_exec:
        box(s, Inches(0.5), ry, Inches(0.9), Inches(0.3), clr, radius=False)
        txt(s, level, Inches(0.55), ry + Inches(0.02), Inches(0.8), Inches(0.25),
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, risk, Inches(1.5), ry, Inches(2.5), Inches(0.3),
            size=11, bold=True, color=BRAND_BLUE)
        txt(s, mitigation, Inches(1.5), ry + Inches(0.3), Inches(8.0), Inches(0.4),
            size=10, color=BODY_TEXT)
        thin_line(s, Inches(0.5), ry + Inches(0.75), Inches(9.0))
        ry += Inches(0.82)

    # ═══════════════════════════════════════════
    # SLIDE 9 — Сравнительная таблица (краткая)
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout)
    add_brand_header(s, 'Сравнение: ключевые параметры')

    comparison = [
        ('Лицензия', '~$108K/год, подписка', 'Разовая покупка, рубли'),
        ('ОС', 'Только Windows', 'Windows + Linux (Astra, РЕД ОС)'),
        ('Реестр Минцифры', 'Нет', 'Да — соответствие ФЗ-187'),
        ('Поддержка в РФ', 'Нет (санкции)', '24/7 на русском, вендор в Томске'),
        ('Визуализация', 'PI Vision — простая', 'Alpha.HMI 2.0 — мощная'),
        ('Коннекторы', '450+ протоколов', '~30 + OPC UA шлюзы (~80%)'),
        ('Облако', 'AVEVA Data Hub', 'Нет (on-prem = безопаснее)'),
        ('Event Frames', 'Есть', 'Нет (доработка если нужно)'),
        ('REST API', 'Есть', 'Нет (доступ через TCP/SQL)'),
        ('ФСТЭК', 'Нет', 'Лицензия на разработку СЗКИ'),
    ]

    hy = Inches(0.85)
    box(s, Inches(0.5), hy, Inches(2.3), Inches(0.32), BRAND_BLUE, radius=False)
    box(s, Inches(2.9), hy, Inches(3.1), Inches(0.32), RGBColor(0x80, 0x40, 0x40), radius=False)
    box(s, Inches(6.1), hy, Inches(3.4), Inches(0.32), RGBColor(0x00, 0x6B, 0x3F), radius=False)
    txt(s, 'Параметр', Inches(0.6), hy + Inches(0.03), Inches(2.1), Inches(0.25), size=10, bold=True, color=WHITE)
    txt(s, 'PI System', Inches(3.0), hy + Inches(0.03), Inches(2.9), Inches(0.25), size=10, bold=True, color=WHITE)
    txt(s, 'Альфа платформа', Inches(6.2), hy + Inches(0.03), Inches(3.2), Inches(0.25), size=10, bold=True, color=WHITE)

    ry = hy + Inches(0.37)
    for i, (param, pi_val, alpha_val) in enumerate(comparison):
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        box(s, Inches(0.5), ry, Inches(9.0), Inches(0.35), bg, radius=False)
        txt(s, param, Inches(0.6), ry + Inches(0.04), Inches(2.1), Inches(0.27), size=9, bold=True, color=BRAND_BLUE)
        txt(s, pi_val, Inches(3.0), ry + Inches(0.04), Inches(2.9), Inches(0.27), size=9, color=BODY_TEXT)
        txt(s, alpha_val, Inches(6.2), ry + Inches(0.04), Inches(3.2), Inches(0.27), size=9, color=BODY_TEXT)
        ry += Inches(0.37)

    # ═══════════════════════════════════════════
    # SLIDE 10 — Что просим: решение
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout)
    add_brand_header(s, 'Запрос на решение')

    txt(s, 'Предлагаем утвердить:', Inches(0.6), Inches(0.9), Inches(9.0), Inches(0.35),
        size=14, bold=True, color=BRAND_BLUE)

    decisions = [
        ('1', 'Запуск аудита PI-инфраструктуры', 
         'Инвентаризация: теги, протоколы, скрипты, отчёты, Event Frames.\nСрок: 1-2 месяца. Результат: отчёт с рекомендациями.'),
        ('2', 'Определение пилотного объекта',
         'Один объект для проверки на практике.\nPI и Альфа работают параллельно 2-4 недели.'),
        ('3', 'Формирование рабочей группы',
         'НОВАТЭК (КИПиА + IT) + Атомик Софт + интегратор.\nЕженедельные статус-встречи.'),
        ('4', 'Бюджет на пилот',
         'Лицензии Альфа платформы (пилотный комплект) +\nработы по миграции пилотного объекта.'),
    ]

    py = Inches(1.4)
    for num, title, desc in decisions:
        # Number circle
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), py + Inches(0.05), Inches(0.35), Inches(0.35))
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL
        dot.line.fill.background()
        txt(s, num, Inches(0.62), py + Inches(0.09), Inches(0.31), Inches(0.27),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        
        txt(s, title, Inches(1.1), py, Inches(8.3), Inches(0.3),
            size=13, bold=True, color=BRAND_BLUE)
        txt(s, desc, Inches(1.1), py + Inches(0.3), Inches(8.3), Inches(0.6),
            size=10, color=BODY_TEXT)
        py += Inches(0.95)

    # ═══════════════════════════════════════════
    # SLIDE 11 — Контакты
    # ═══════════════════════════════════════════
    s = prs.slides.add_slide(layout_contact)

    circle_path = os.path.join(ASSETS, 'contact_circle.png')
    if os.path.exists(circle_path):
        s.shapes.add_picture(circle_path, Inches(0), Inches(0.2), Inches(4.0), Inches(4.0))

    txt(s, 'Будущее автоматизации', Inches(0.5), Inches(0.90), Inches(9.0), Inches(0.7),
        size=28, bold=True, color=BRAND_BLUE, align=PP_ALIGN.CENTER)
    txt(s, 'в ваших руках.', Inches(0.5), Inches(1.50), Inches(9.0), Inches(0.7),
        size=28, bold=True, color=BRAND_BLUE, align=PP_ALIGN.CENTER)

    qr_path = os.path.join(ASSETS, 'qr_or_deco.png')
    if os.path.exists(qr_path):
        s.shapes.add_picture(qr_path, Inches(7.5), Inches(2.4), Inches(1.7), Inches(1.7))

    txt(s, 'АО «Атомик Софт»\n634050 Томск, пр. Ленина 60/1\n+7 (3822) 281 914',
        Inches(0.6), Inches(4.40), Inches(2.8), Inches(1.1), size=11, color=BODY_TEXT)
    txt(s, 'info@automiq.ru\nwww.automiq.ru',
        Inches(3.5), Inches(4.40), Inches(2.5), Inches(1.1), size=11, color=BODY_TEXT)
    txt(s, 'Станислав Павловский\nTelegram: @Integraleus',
        Inches(6.5), Inches(4.40), Inches(3.0), Inches(1.1), size=11, color=BODY_TEXT)

    prs.save(OUTPUT)
    print(f'✅ Saved: {OUTPUT}')
    print(f'   Slides: {len(prs.slides)}')


if __name__ == '__main__':
    build()
